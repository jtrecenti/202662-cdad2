"""
Monta os slides da Aula 3: filtrar e resumir.

O deck constroi a estatistica descritiva de baixo para cima: primeiro o que e
uma estatistica, depois um exemplo de uma coluna virando um numero, depois o
mapa entre tipo de variavel e conta possivel, e so entao cada medida de posicao
em seu proprio slide. Fecha juntando tudo na mesma coluna, com o mapa mental, e
lanca o Projeto 1.

Nenhum numero e digitado a mao: todos saem de `dados/tjsp_cjsg_dano_moral.csv`,
a mesma base do notebook da aula, para que o slide e o codigo do aluno digam a
mesma coisa. As figuras vem de `graficos_aula03.py`.

Uso:
    python graficos_aula03.py      # gera as figuras em assets/aula03
    python build_aula03.py         # monta o pptx
"""

from __future__ import annotations

import os

import pandas as pd
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

from insper import (
    AMARELO, BRANCO, CINZA, CINZA_CLARO, CINZA_ESCURO, DISPLAY, FAIXA, LARANJA,
    MARGEM, PRETO, QUASE_BRANCO, ROXO, TURQUESA, VERMELHO, caixa,
    celula, deck_limpo, escrever, gravar, imagem_ajustada, layouts,
    lista_com_barras, slide_capa, slide_com_titulo, slide_secao,
    tabela_sem_estilo, texto_livre,
)

AQUI = os.path.dirname(os.path.abspath(__file__))
SAIDA = os.path.join(AQUI, "aula03.pptx")
FIGURAS = os.path.join(AQUI, "assets", "aula03")
BASE = os.path.join(os.path.dirname(AQUI), "dados", "tjsp_cjsg_dano_moral.csv")

EYEBROW = "Ciência de Dados Aplicada ao Direito II  ·  Aula 3  ·  Filtrar e resumir"
TITULO_DECK = "Aula 3: filtrar e resumir"

# link direto para o arquivo, e nao para a pasta: o aluno clica e baixa, sem
# passar pelo repositorio e sem escolher entre CSV e Excel no meio da aula
BASE_EM_EXCEL = ("https://raw.githubusercontent.com/jtrecenti/202662-cdad2/"
                 "main/dados/tjsp_cjsg_dano_moral.xlsx")

COR_NUMERICA = TURQUESA
COR_CATEGORICA = ROXO


# ------------------------------------------------------------------ numeros

def brl(v: float, casas: int = 0) -> str:
    texto = f"{v:,.{casas}f}".replace(",", "@").replace(".", ",").replace("@", ".")
    return f"R$ {texto}"


def num(v: float, casas: int = 0) -> str:
    """Numero no padrao brasileiro: ponto no milhar, virgula no decimal."""
    return f"{v:,.{casas}f}".replace(",", "@").replace(".", ",").replace("@", ".")


def pctbr(v: float, casas: int = 0) -> str:
    return f"{v:.{casas}f}".replace(".", ",") + "%"


def estatisticas() -> dict:
    """Tudo o que o deck afirma sobre a base, calculado na hora.

    Se a base mudar, os slides mudam junto. Numero de slide que nao vem do
    arquivo e numero que envelhece sem avisar.
    """
    d = pd.read_csv(BASE)
    v = d["valor_indenizacao"].dropna()
    sem_maximo = v[v < v.max()]
    return {
        "n_acordaos": len(d),
        "n": len(v),
        "media": v.mean(),
        "mediana": v.median(),
        "moda": float(v.mode().iloc[0]),
        "n_moda": int((v == v.mode().iloc[0]).sum()),
        "minimo": v.min(),
        "maximo": v.max(),
        "q1": v.quantile(.25),
        "q3": v.quantile(.75),
        "q90": v.quantile(.9),
        "iqr": v.quantile(.75) - v.quantile(.25),
        "amplitude": v.max() - v.min(),
        "variancia": v.var(),
        "desvio": v.std(),
        "queda_sem_maximo": v.mean() - sem_maximo.mean(),
        "prop_dano_moral": float(d["tem_dano_moral"].mean()),
        "prop_acima_10k": float((v > 10000).mean()),
        "n_acima_10k": int((v > 10000).sum()),
        # o que aconteceria se alguem preenchesse o faltante com zero
        "n_faltantes": int(d["valor_indenizacao"].isna().sum()),
        "media_com_zero": d["valor_indenizacao"].fillna(0).mean(),
        # comparacao entre grupos, que e groupby e so aparece na aula 4
        "mediana_por_assunto": (
            d[d["assunto"].isin(d["assunto"].value_counts().head(3).index)]
            .groupby("assunto")["valor_indenizacao"].median()
            .sort_values(ascending=False)
        ),
        "exemplos": d[["processo", "comarca", "valor_indenizacao"]].head(5),
        # o slide de quartis usa o tamanho da ementa, e nao o valor: em
        # valor_indenizacao o Q1 cai em cima da mediana (ver comentario la)
        "q1_palavras": d["n_palavras_ementa"].quantile(.25),
        "mediana_palavras": d["n_palavras_ementa"].median(),
        "q3_palavras": d["n_palavras_ementa"].quantile(.75),
        # proporcao: uma coluna binaria e uma categorica
        "n_dano_moral": int(d["tem_dano_moral"].sum()),
        "n_sp": int((d["comarca"] == "São Paulo").sum()),
        "prop_sp": float((d["comarca"] == "São Paulo").mean()),
        "n_comarca_top": int(d["comarca"].value_counts().iloc[0]),
        # os sete primeiros valores distintos, para a mediana feita na mao
        "sete": [float(x) for x in pd.unique(v)[:7]],
        "periodo": (
            pd.to_datetime(d["data_julgamento"], format="%d/%m/%Y",
                           errors="coerce").min(),
            pd.to_datetime(d["data_julgamento"], format="%d/%m/%Y",
                           errors="coerce").max(),
        ),
    }


E = estatisticas()


# ------------------------------------------------------------------ pecas

def nota(slide, y, titulo, texto, cor=VERMELHO, *, altura=Inches(0.86),
         largura=FAIXA, x=MARGEM, tamanho_texto=12.5):
    """Faixa de destaque com barra colorida a esquerda.

    `texto` aceita uma string ou uma lista de strings, uma por paragrafo.
    """
    caixa(slide, x, y, largura, altura, preenchimento=QUASE_BRANCO)
    caixa(slide, x, y, Inches(0.09), altura, preenchimento=cor)
    tb = texto_livre(slide, x + Inches(0.31), y + Inches(0.13),
                     largura - Inches(0.62), altura - Inches(0.22))
    paragrafos = [texto] if isinstance(texto, str) else list(texto)
    escrever(tb.text_frame, [
        {"texto": titulo, "tamanho": 14.5, "bold": True, "depois": 3},
        *[{"texto": t, "tamanho": tamanho_texto, "cor": CINZA_ESCURO,
           "entrelinhas": 1.12, "depois": 7 if i < len(paragrafos) - 1 else 0}
          for i, t in enumerate(paragrafos)],
    ])


def faixa_do_valor(slide, y, *, rotulo, valor, cor, calculo, codigo,
                   tamanho_valor=25):
    """O numero grande da medida, com a conta e o metodo do pandas ao lado."""
    altura = Inches(0.86)
    largura_cartao = Inches(3.30)
    caixa(slide, MARGEM, y, largura_cartao, altura, preenchimento=QUASE_BRANCO)
    caixa(slide, MARGEM, y, largura_cartao, Inches(0.09), preenchimento=cor)
    tb = texto_livre(slide, MARGEM + Inches(0.24), y + Inches(0.19),
                     largura_cartao - Inches(0.48), Inches(0.60))
    escrever(tb.text_frame, [
        {"texto": rotulo, "tamanho": 11.5, "cor": CINZA_ESCURO, "depois": 1},
        {"texto": valor, "fonte": DISPLAY, "tamanho": tamanho_valor,
         "bold": True, "depois": 0},
    ])

    tb = texto_livre(slide, MARGEM + largura_cartao + Inches(0.34),
                     y + Inches(0.14), Inches(7.50), Inches(0.66))
    escrever(tb.text_frame, [
        {"texto": calculo, "tamanho": 13, "entrelinhas": 1.1, "depois": 5},
        {"texto": codigo, "fonte": "Consolas", "tamanho": 11.5,
         "cor": CINZA_ESCURO, "depois": 0},
    ])


def slide_medida(prs, lays, *, titulo, definicao, rotulo, valor, cor, calculo,
                 codigo, figura, nota_titulo, nota_texto, largo=True,
                 tamanho_valor=25, altura_nota=Inches(2.72)):
    """O molde comum das medidas de posicao: uma por slide, sempre igual.

    Repetir a mesma estrutura sete vezes e de proposito: o aluno para de gastar
    atencao na diagramacao e olha so o que mudou, que e a medida.
    """
    slide = slide_com_titulo(prs, lays, titulo, EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.92), FAIXA, Inches(0.52))
    escrever(tb.text_frame, [
        {"texto": definicao, "tamanho": 14, "cor": CINZA_ESCURO,
         "entrelinhas": 1.12, "depois": 0},
    ])

    faixa_do_valor(slide, Inches(2.56), rotulo=rotulo, valor=valor, cor=cor,
                   calculo=calculo, codigo=codigo, tamanho_valor=tamanho_valor)

    caminho = os.path.join(FIGURAS, figura)
    if largo:
        imagem_ajustada(slide, caminho, MARGEM, Inches(3.46), FAIXA,
                        Inches(2.04), moldura=False)
        nota(slide, Inches(5.56), nota_titulo, nota_texto, cor,
             altura=Inches(1.14))
    else:
        imagem_ajustada(slide, caminho, MARGEM, Inches(3.48), Inches(4.90),
                        Inches(2.72), moldura=False)
        nota(slide, Inches(3.48), nota_titulo, nota_texto, cor,
             altura=altura_nota, largura=FAIXA - Inches(5.20),
             x=MARGEM + Inches(5.20), tamanho_texto=11.5)
    return slide


def s01_capa(prs, lays):
    return slide_capa(
        prs, lays,
        subtitulo="Aula 3: filtrar e resumir",
        subtema="Medidas de posição e de variabilidade",
    )


def s02_secao_estatistica(prs, lays):
    return slide_secao(prs, lays, "O que é uma estatística")


def s03_o_que_e(prs, lays):
    """Sem numero nenhum da base: ela so e apresentada no slide seguinte.

    Na primeira versao este slide ja falava em "303 indenizacoes" e em "302
    outros casos", antes de alguem saber que base era essa. O numero chegava
    sem contexto e o slide seguinte tinha de reapresentar tudo.
    """
    slide = slide_com_titulo(prs, lays, "Estatística é um resumo da amostra",
                             EYEBROW)

    caixa(slide, MARGEM, Inches(2.02), FAIXA, Inches(1.00),
          preenchimento=QUASE_BRANCO)
    caixa(slide, MARGEM, Inches(2.02), Inches(0.09), Inches(1.00),
          preenchimento=VERMELHO)
    tb = texto_livre(slide, MARGEM + Inches(0.34), Inches(2.20), Inches(10.8),
                     Inches(0.70))
    escrever(tb.text_frame, [
        {"texto": "Uma estatística é um número calculado a partir dos dados da "
                  "amostra, que resume uma característica dela.",
         "tamanho": 17, "bold": True, "entrelinhas": 1.12, "depois": 0},
    ])

    passos = [
        (TURQUESA, "muitos valores", "uma coluna inteira da tabela, com um "
                                     "valor em cada linha"),
        (LARANJA, "uma conta", "somar e dividir, ordenar e pegar o do meio, "
                               "contar quantos são"),
        (VERMELHO, "um número", "o que você reporta, e o que o leitor vai "
                                "guardar"),
    ]
    largura = Inches(3.62)
    vao = Inches(0.32)
    y = Inches(3.30)
    for i, (cor, rotulo, detalhe) in enumerate(passos):
        x = MARGEM + i * (largura + vao)
        caixa(slide, x, y, largura, Inches(1.20), preenchimento=QUASE_BRANCO)
        caixa(slide, x, y, largura, Inches(0.09), preenchimento=cor)
        tb = texto_livre(slide, x + Inches(0.26), y + Inches(0.26),
                         largura - Inches(0.52), Inches(0.86))
        escrever(tb.text_frame, [
            {"texto": rotulo, "fonte": DISPLAY, "tamanho": 20, "bold": True,
             "depois": 4},
            {"texto": detalhe, "tamanho": 12.5, "cor": CINZA_ESCURO,
             "entrelinhas": 1.12, "depois": 0},
        ])
        if i < 2:
            seta = texto_livre(slide, x + largura, y + Inches(0.34),
                               vao, Inches(0.40))
            escrever(seta.text_frame, [
                {"texto": "→", "tamanho": 18, "cor": CINZA, "depois": 0},
            ])

    nota(slide, Inches(4.90),
         "Resumir é perder informação de propósito",
         "O número que sobra não guarda nenhum caso individual: depois de "
         "calcular a média, os casos que a geraram somem. A pergunta da aula "
         "não é como calcular, é qual informação vale a pena perder, e isso "
         "depende do tipo da variável e da pergunta de pesquisa.",
         TURQUESA, altura=Inches(0.96))

    tb = texto_livre(slide, MARGEM, Inches(5.95), FAIXA, Inches(0.70))
    escrever(tb.text_frame, [
        {"texto": "Amostra é o que está na sua tabela. População é o conjunto "
                  "sobre o qual você quer falar: no caso de hoje, todos os "
                  "acórdãos do TJSP. Esta aula descreve a amostra. Usar a "
                  "amostra para afirmar algo sobre a população é inferência, e "
                  "fica para a aula 18.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "entrelinhas": 1.15, "depois": 0},
    ])
    return slide


def s04_coluna_vira_numero(prs, lays):
    """Apresenta a base uma unica vez, com link, e so entao calcula algo.

    Todos os numeros do resto do deck se referem a estes 460 acordaos, ou aos
    303 deles que trazem valor. Se o aluno perder este slide, perde o
    denominador de tudo o que vem depois.
    """
    slide = slide_com_titulo(prs, lays, "De uma coluna a um número", EYEBROW)

    inicio, fim = E["periodo"]
    meses = ("janeiro fevereiro março abril maio junho julho agosto setembro "
             "outubro novembro dezembro").split()
    periodo = (f"{meses[inicio.month - 1]} de {inicio.year} a "
               f"{meses[fim.month - 1]} de {fim.year}")

    tb = texto_livre(slide, MARGEM, Inches(1.90), FAIXA, Inches(0.56))
    escrever(tb.text_frame, [
        {"texto": f"A base de hoje: {E['n_acordaos']} acórdãos de apelação "
                  f"cível do TJSP em que se discute indenização por dano moral, "
                  f"julgados de {periodo}. Em {E['n']} deles a ementa traz o "
                  f"valor arbitrado; nos outros {E['n_acordaos'] - E['n']}, não.",
         "tamanho": 14, "cor": CINZA_ESCURO, "entrelinhas": 1.12, "depois": 0},
    ])

    caixa(slide, MARGEM, Inches(2.50), FAIXA, Inches(0.40),
          preenchimento=QUASE_BRANCO)
    caixa(slide, MARGEM, Inches(2.50), Inches(0.09), Inches(0.40),
          preenchimento=TURQUESA)
    tb = texto_livre(slide, MARGEM + Inches(0.30), Inches(2.56), Inches(11.0),
                     Inches(0.30))
    escrever(tb.text_frame, [
        {"texto": "Baixar a tabela inteira em Excel  ·  460 linhas, 16 colunas",
         "tamanho": 13, "bold": True, "cor": VERMELHO, "sublinhado": True,
         "link": BASE_EM_EXCEL, "depois": 0},
    ])

    tb = texto_livre(slide, MARGEM, Inches(3.02), Inches(6.10), Inches(0.26))
    escrever(tb.text_frame, [
        {"texto": "a variável é a coluna", "tamanho": 11.5, "cor": CINZA_ESCURO,
         "depois": 0},
    ])

    linhas = list(E["exemplos"].itertuples(index=False))
    tabela_shape = slide.shapes.add_table(
        len(linhas) + 2, 3, MARGEM, Inches(3.30), Inches(6.10), Inches(2.10)
    )
    tabela = tabela_shape.table
    tabela_sem_estilo(tabela)
    for largura, i in zip((Inches(2.60), Inches(1.70), Inches(1.80)), range(3)):
        tabela.columns[i].width = largura
    for r in range(len(linhas) + 2):
        tabela.rows[r].height = Inches(0.30)

    for c, titulo in enumerate(("processo", "comarca", "valor_indenizacao")):
        celula(tabela, 0, c, titulo, tamanho=10.5, bold=True, cor=BRANCO,
               fundo=PRETO)
    for r, linha in enumerate(linhas, start=1):
        zebra = QUASE_BRANCO if r % 2 == 0 else None
        celula(tabela, r, 0, linha.processo, tamanho=10.5, fundo=zebra)
        celula(tabela, r, 1, linha.comarca, tamanho=10.5, fundo=zebra)
        tem_valor = linha.valor_indenizacao == linha.valor_indenizacao
        celula(tabela, r, 2,
               brl(linha.valor_indenizacao) if tem_valor else "sem valor",
               tamanho=10.5, bold=tem_valor, fundo=zebra,
               cor=VERMELHO if tem_valor else CINZA)
    ultima = len(linhas) + 1
    for c, texto in enumerate((f"e mais {E['n_acordaos'] - len(linhas)} linhas",
                               "", "")):
        celula(tabela, ultima, c, texto, tamanho=10.5, cor=CINZA_ESCURO)

    x_cartao = Inches(7.55)
    largura_cartao = Inches(5.12)
    tb = texto_livre(slide, x_cartao, Inches(3.02), largura_cartao, Inches(0.26))
    escrever(tb.text_frame, [
        {"texto": "as estatísticas são números que saem dela", "tamanho": 11.5,
         "cor": CINZA_ESCURO, "depois": 0},
    ])

    cartoes = [
        (ROXO, "contagem", f"{E['n_comarca_top']}",
         "acórdãos da comarca de São Paulo",
         'df["comarca"].value_counts()'),
        (VERMELHO, "média", brl(E["media"]), "valor médio arbitrado",
         'df["valor_indenizacao"].mean()'),
    ]
    y = Inches(3.30)
    for cor, rotulo, valor, detalhe, codigo in cartoes:
        caixa(slide, x_cartao, y, largura_cartao, Inches(1.00),
              preenchimento=QUASE_BRANCO)
        caixa(slide, x_cartao, y, largura_cartao, Inches(0.09),
              preenchimento=cor)
        tb = texto_livre(slide, x_cartao + Inches(0.26), y + Inches(0.18),
                         largura_cartao - Inches(0.52), Inches(0.76))
        escrever(tb.text_frame, [
            {"texto": f"{rotulo}  ·  {detalhe}", "tamanho": 11.5,
             "cor": CINZA_ESCURO, "depois": 1},
            {"texto": valor, "fonte": DISPLAY, "tamanho": 22, "bold": True,
             "depois": 2},
            {"texto": codigo, "fonte": "Consolas", "tamanho": 11,
             "cor": CINZA_ESCURO, "depois": 0},
        ])
        y = y + Inches(1.10)

    com_valor = next(l for l in linhas
                     if l.valor_indenizacao == l.valor_indenizacao)
    nota(slide, Inches(5.58),
         "Três coisas diferentes, e vale separar os nomes",
         f"A coluna inteira é a variável. Cada célula é um valor observado: "
         f"{brl(com_valor.valor_indenizacao)} naquele acórdão de "
         f"{com_valor.comarca}, e nenhum valor na primeira linha, porque aquela "
         f"ementa não traz. O número que sai da conta sobre a coluna toda é a "
         f"estatística, e ele não é o valor de nenhum acórdão em particular.",
         altura=Inches(1.12))
    return slide


def s05_duas_perguntas(prs, lays):
    slide = slide_com_titulo(prs, lays, "Duas perguntas sobre a mesma coluna",
                             EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.94), FAIXA, Inches(0.42))
    escrever(tb.text_frame, [
        {"texto": "A pergunta de hoje é: quanto costuma valer uma indenização "
                  "por dano moral no TJSP, e o quanto esse valor varia? Toda "
                  "estatística descritiva responde a uma das duas metades.",
         "tamanho": 14, "cor": CINZA_ESCURO, "entrelinhas": 1.12, "depois": 0},
    ])

    colunas = [
        (TURQUESA, "POSIÇÃO", "Onde fica o centro?",
         "Qual é o valor típico da coluna. É o número que responde a pergunta "
         "“quanto costuma ser”.",
         "média · mediana · moda · mínimo · máximo · quartis · proporção"),
        (LARANJA, "VARIABILIDADE", "Como os valores variam?",
         "O quanto os casos se espalham em volta desse centro. É o número que "
         "diz se o típico representa alguém.",
         "amplitude · intervalo interquartil · variância · desvio padrão"),
    ]
    largura = Inches(5.59)
    vao = Inches(0.32)
    y = Inches(2.58)
    altura = Inches(2.30)
    for i, (cor, etiqueta, pergunta, texto, quais) in enumerate(colunas):
        x = MARGEM + i * (largura + vao)
        caixa(slide, x, y, largura, altura, preenchimento=QUASE_BRANCO)
        caixa(slide, x, y, largura, Inches(0.09), preenchimento=cor)
        tb = texto_livre(slide, x + Inches(0.28), y + Inches(0.30),
                         largura - Inches(0.56), altura - Inches(0.48))
        escrever(tb.text_frame, [
            {"texto": etiqueta, "tamanho": 11.5, "cor": CINZA_ESCURO,
             "depois": 3},
            {"texto": pergunta, "fonte": DISPLAY, "tamanho": 22, "bold": True,
             "depois": 8},
            {"texto": texto, "tamanho": 13.5, "entrelinhas": 1.15, "depois": 10},
            {"texto": quais, "tamanho": 12, "cor": CINZA_ESCURO,
             "entrelinhas": 1.15, "depois": 0},
        ])

    nota(slide, Inches(5.14),
         "Uma sem a outra conta meia história",
         f"“A indenização típica é {brl(E['mediana'])}” parece uma resposta "
         f"fechada, mas nesta mesma base o menor valor é {brl(E['minimo'], 2)} "
         f"e o maior é {brl(E['maximo'])}. Sem dizer o quanto varia, o valor "
         f"típico sugere uma regularidade que não existe.",
         altura=Inches(1.04))
    return slide


def s06_mapa_tipos(prs, lays):
    slide = slide_com_titulo(prs, lays, "Que conta cabe em cada tipo de variável",
                             EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.96), FAIXA, Inches(0.40))
    escrever(tb.text_frame, [
        {"texto": "O tipo que você declarou na aula 2 decide o que dá para "
                  "calcular. Média de categoria não é conta difícil, é conta "
                  "sem sentido.",
         "tamanho": 14, "cor": CINZA_ESCURO, "depois": 0},
    ])

    dados = [
        ("numérica contínua", COR_NUMERICA,
         "média, mediana, mínimo, máximo, quartis",
         "variância, desvio padrão, amplitude, IIQ"),
        ("numérica discreta", COR_NUMERICA,
         "as mesmas, e a moda passa a fazer sentido",
         "variância, desvio padrão, amplitude, IIQ"),
        ("categórica ordinal", COR_CATEGORICA,
         "mediana, moda, mínimo, máximo, quartis",
         "amplitude e IIQ, que aqui saem como categorias"),
        ("categórica nominal", COR_CATEGORICA,
         "moda, proporção de cada categoria",
         "nenhuma das quatro"),
        ("categórica binária", COR_CATEGORICA,
         "proporção, que aqui é a própria média",
         "desvio padrão, que sai da proporção"),
        ("identificador", CINZA, "nenhuma", "nenhuma"),
    ]

    tabela_shape = slide.shapes.add_table(
        len(dados) + 1, 4, MARGEM, Inches(2.48), FAIXA, Inches(2.56)
    )
    tabela = tabela_shape.table
    tabela_sem_estilo(tabela)
    for largura, i in zip(
        (Inches(0.13), Inches(2.95), Inches(4.42), Inches(4.00)), range(4)
    ):
        tabela.columns[i].width = largura
    tabela.rows[0].height = Inches(0.30)
    for r in range(1, len(dados) + 1):
        tabela.rows[r].height = Inches(0.37)

    for c, titulo in enumerate(("", "tipo da variável", "posição",
                                "variabilidade")):
        celula(tabela, 0, c, titulo, tamanho=11, bold=True, cor=BRANCO,
               fundo=PRETO)
    for r, (tipo, cor, posicao, variabilidade) in enumerate(dados, start=1):
        zebra = QUASE_BRANCO if r % 2 == 0 else None
        celula(tabela, r, 0, "", fundo=cor)
        celula(tabela, r, 1, tipo, tamanho=12, bold=True, fundo=zebra)
        celula(tabela, r, 2, posicao, tamanho=12, fundo=zebra)
        celula(tabela, r, 3, variabilidade, tamanho=12, fundo=zebra,
               cor=CINZA_ESCURO)

    return slide


def s07_secao_posicao(prs, lays):
    return slide_secao(prs, lays, "Medidas de posição")


def s08_media(prs, lays):
    return slide_medida(
        prs, lays,
        titulo="Média",
        definicao="Soma de todos os valores dividida pelo número de casos. É o "
                  "ponto de equilíbrio da distribuição: cada caso puxa a média "
                  "na proporção do quanto ele vale.",
        rotulo="média das indenizações", valor=brl(E["media"]), cor=VERMELHO,
        calculo=f"Somar os {E['n']} valores e dividir por {E['n']}.",
        codigo='df["valor_indenizacao"].mean()',
        figura="media.png",
        nota_titulo="Quem puxa a média são os casos extremos",
        nota_texto=f"Tirar da base o único acórdão de {brl(E['maximo'])} derruba "
                   f"a média em {brl(E['queda_sem_maximo'])}. Um caso em "
                   f"{E['n']} move o resumo de todos os outros, e é por isso que "
                   f"a média sozinha engana em valor de indenização.",
    )


def s09_mediana(prs, lays):
    return slide_medida(
        prs, lays,
        titulo="Mediana",
        definicao="Ordene os valores do menor para o maior e pegue o do meio. "
                  "Metade dos casos fica abaixo dela, metade acima.",
        rotulo="mediana das indenizações", valor=brl(E["mediana"]), cor=TURQUESA,
        calculo=f"Com {E['n']} valores ordenados, é o valor da posição "
                f"{(E['n'] + 1) // 2}.",
        codigo='df["valor_indenizacao"].median()',
        figura="mediana.png",
        nota_titulo=f"A média é {(E['media'] / E['mediana'] - 1) * 100:.0f}% "
                    f"maior que a mediana",
        nota_texto="A mediana só conta quantos casos estão de cada lado, e não "
                   "o quanto eles valem. Trocar o maior acórdão por um de um "
                   "milhão não a move um centavo. Em valor monetário no "
                   "Direito, que quase sempre tem cauda longa à direita, é ela "
                   "que responde “quanto costuma ser”.",
        largo=True,
    )


def s09b_mediana_na_mao(prs, lays):
    """A mecanica da mediana em sete numeros, antes do histograma de 303.

    Sao sete valores reais da base, os sete primeiros distintos. Numero impar
    de proposito: com par nao existe "o do meio" e a regra muda, o que vira o
    comentario do rodape.
    """
    slide = slide_com_titulo(prs, lays, "A mediana feita na mão", EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.92), FAIXA, Inches(0.40))
    escrever(tb.text_frame, [
        {"texto": "Sete acórdãos da base, para ver a conta acontecer. Ordenar "
                  "primeiro não é detalhe: é a conta inteira.",
         "tamanho": 14, "cor": CINZA_ESCURO, "depois": 0},
    ])

    valores = E["sete"]
    ordenados = sorted(valores)
    meio = len(ordenados) // 2
    largura = (FAIXA - Inches(0.12) * 6) // 7

    linhas = [
        (Inches(2.44), "como vieram da tabela", valores, None, CINZA_CLARO),
        (Inches(3.54), "ordenados do menor para o maior", ordenados, None,
         CINZA_CLARO),
        (Inches(4.64), "o do meio é a mediana", ordenados, meio, TURQUESA),
    ]
    for y, rotulo, quais, destaque, cor in linhas:
        tb = texto_livre(slide, MARGEM, y, FAIXA, Inches(0.26))
        escrever(tb.text_frame, [
            {"texto": rotulo, "tamanho": 11.5, "cor": CINZA_ESCURO, "depois": 0},
        ])
        for i, valor in enumerate(quais):
            x = MARGEM + i * (largura + Inches(0.12))
            marcado = destaque is not None and i == destaque
            caixa(slide, x, y + Inches(0.28), largura, Inches(0.54),
                  preenchimento=cor if marcado else QUASE_BRANCO,
                  borda=None if marcado else CINZA_CLARO)
            tb = texto_livre(slide, x, y + Inches(0.36), largura, Inches(0.40))
            escrever(tb.text_frame, [
                {"texto": brl(valor), "fonte": DISPLAY,
                 "tamanho": 16 if marcado else 15, "bold": True,
                 "alinhamento": PP_ALIGN.CENTER, "depois": 0},
            ])
        if destaque is not None:
            trio = largura * 3 + Inches(0.12) * 2
            for i, texto in ((0, "3 abaixo"), (destaque + 1, "3 acima")):
                tb = texto_livre(slide, MARGEM + i * (largura + Inches(0.12)),
                                 y + Inches(0.86), trio, Inches(0.26))
                escrever(tb.text_frame, [
                    {"texto": texto, "tamanho": 11, "cor": CINZA_ESCURO,
                     "alinhamento": PP_ALIGN.CENTER, "depois": 0},
                ])

    nota(slide, Inches(5.78),
         f"Com sete valores a mediana é o quarto; com {E['n']}, é o "
         f"{(E['n'] + 1) // 2}º",
         "Se o número de casos for par não existe “o do meio”, e a mediana "
         "passa a ser a média dos dois centrais. Repare também que a média "
         "destes sete não é nenhum deles, enquanto a mediana é sempre um valor "
         "que existe na base.",
         TURQUESA, altura=Inches(0.90))
    return slide


def s10_moda(prs, lays):
    pct = pctbr(E["n_moda"] / E["n"] * 100)
    return slide_medida(
        prs, lays,
        titulo="Moda",
        definicao="O valor que mais se repete na coluna. Serve sobretudo em "
                  "variável categórica, onde média e mediana não existem, e "
                  "pouco em numérica, onde os valores raramente se repetem.",
        rotulo="valor mais frequente", valor=brl(E["moda"]), cor=ROXO,
        calculo=f"Aparece em {E['n_moda']} dos {E['n']} acórdãos, {pct} do "
                f"total.",
        codigo='df["valor_indenizacao"].value_counts()',
        figura="moda.png",
        nota_titulo="Aqui a moda diz algo justamente porque não devia dizer nada",
        nota_texto=f"Se cada indenização fosse calculada caso a caso, dois "
                   f"acórdãos dificilmente cairiam no mesmo centavo. Mas {pct} "
                   f"caem exatamente em {brl(E['moda'])}, e os oito valores "
                   f"mais frequentes são todos redondos. O valor é arbitrado "
                   f"por ancoragem em números redondos, e não por cálculo. "
                   f"Nenhuma outra medida de posição enxerga isso.",
        largo=False,
    )


def s11_minimo(prs, lays):
    return slide_medida(
        prs, lays,
        titulo="Mínimo",
        definicao="O menor valor observado. Junto com o máximo, delimita onde a "
                  "sua variável de fato vive.",
        rotulo="menor valor da base", valor=brl(E["minimo"], 2), cor=LARANJA,
        calculo="O primeiro valor depois de ordenar a coluna.",
        codigo='df["valor_indenizacao"].min()',
        figura="minimo.png",
        nota_titulo="O mínimo é onde o erro de extração aparece primeiro",
        nota_texto=[
            f"{brl(E['minimo'], 2)} não é uma indenização por dano moral. "
            f"Naquele acórdão é o desconto mensal indevido no benefício "
            f"previdenciário, e a regra que usamos pegou o primeiro valor em "
            f"reais da ementa. Olhar o mínimo e o máximo antes de qualquer "
            f"outra conta é o jeito mais barato de descobrir que a coluna "
            f"mediu outra coisa.",
            f"Faltante também não é zero. {E['n_faltantes']} acórdãos não "
            f"trazem valor e as contas de hoje ignoram essas linhas. Preencher "
            f"zero derrubaria a média de {brl(E['media'])} para "
            f"{brl(E['media_com_zero'])} sem que nenhum juiz arbitrasse nada a "
            f"menos.",
            "No Projeto 1 vale o mesmo: em branco quando a informação não "
            "está na sentença, “não se aplica” quando a pergunta não se coloca, "
            "e nunca zero.",
        ],
        largo=False,
        altura_nota=Inches(3.20),
    )


def s12_maximo(prs, lays):
    return slide_medida(
        prs, lays,
        titulo="Máximo",
        definicao="O maior valor observado. Em base judicial, costuma ser um "
                  "caso só, e costuma ser ele que aparece na manchete.",
        rotulo="maior valor da base", valor=brl(E["maximo"]), cor=VERMELHO,
        calculo="O último valor depois de ordenar a coluna.",
        codigo='df["valor_indenizacao"].max()',
        figura="maximo.png",
        nota_titulo="Um caso, e ele mexe na média de todos",
        nota_texto=f"{brl(E['maximo'])} é {E['maximo'] / E['mediana']:.0f} vezes "
                   f"a mediana. Retirar esse único acórdão baixa a média em "
                   f"{brl(E['queda_sem_maximo'])} e deixa a mediana exatamente "
                   f"onde estava. Máximo alto não é erro: é a informação de que "
                   f"a distribuição tem cauda, e de que a média vai precisar de "
                   f"companhia.",
        largo=False,
    )


def s13_quartis(prs, lays):
    """Os quartis em outra coluna da mesma tabela, e sem boxplot.

    Duas mudancas em relacao a primeira versao. A coluna e o tamanho da ementa
    porque em valor_indenizacao 33% dos acordaos valem exatamente R$ 5.000 e o
    Q1 cai em cima da mediana: os tres cortes nao aparecem separados, e coletar
    mais acordaos nao resolve, porque a concentracao e do fenomeno e nao da
    amostra. E o grafico e histograma com regua, e nao boxplot, que a turma so
    vai encontrar mais para a frente.
    """
    q1, q2, q3 = E["q1_palavras"], E["mediana_palavras"], E["q3_palavras"]
    return slide_medida(
        prs, lays,
        titulo="Quartis",
        definicao="Três cortes que partem a base ordenada em quatro pedaços "
                  "com o mesmo número de casos: Q1 deixa 25% abaixo, Q2 é a "
                  "mediana, Q3 deixa 75%. O exemplo é outra coluna da tabela.",
        rotulo="Q1 · mediana · Q3, em palavras",
        valor=f"{num(q1, 2)} · {num(q2, 1)} · {num(q3)}",
        cor=AMARELO, tamanho_valor=19,
        calculo="A metade central das ementas tem entre "
                f"{num(q1, 2)} e {num(q3)} palavras.",
        codigo='df["n_palavras_ementa"].quantile([0.25, 0.5, 0.75])',
        figura="quartis.png",
        nota_titulo="Por que não usamos aqui a coluna do valor",
        nota_texto=f"Nela os três cortes não apareceriam separados: "
                   f"{E['n_moda']} dos {E['n']} acórdãos valem exatamente "
                   f"{brl(E['moda'])}, então Q1 e mediana caem os dois nesse "
                   f"valor. Não é falta de dados, é o próprio fenômeno. Repare "
                   f"também nas casas decimais: quando o corte cai entre duas "
                   f"observações, o pandas interpola, e por isso aparece meia "
                   f"palavra.",
        largo=True,
    )


def s14_proporcao(prs, lays):
    p = E["prop_dano_moral"]
    return slide_medida(
        prs, lays,
        titulo="Proporção",
        definicao="A parte dividida pelo total. É a medida de posição das "
                  "binárias, serve para cada categoria de uma categórica, e é a "
                  "que mais aparece em trabalho empírico de Direito.",
        rotulo="reconhecem dano moral",
        valor=pctbr(p * 100, 1), cor=TURQUESA,
        calculo=f"{E['n_dano_moral']} dividido por {E['n_acordaos']} dá "
                f"{num(p, 3)}, ou {pctbr(p * 100, 1)}.",
        codigo='df["tem_dano_moral"].mean()',
        figura="proporcao.png",
        nota_titulo="Serve para binária e serve para categórica",
        nota_texto=f"Em coluna de verdadeiro e falso, verdadeiro vale 1 e falso "
                   f"vale 0: somar dá quantos são verdadeiro, e dividir pelo "
                   f"total dá a fração. É por isso que .mean() devolve a "
                   f"proporção. Em coluna de categorias a conta é a mesma, uma "
                   f"proporção por categoria: {E['n_sp']} dos "
                   f"{E['n_acordaos']} acórdãos são da comarca de São Paulo, "
                   f"ou {pctbr(E['prop_sp'] * 100, 1)}, que sai de "
                   f"value_counts(normalize=True).",
        largo=True,
    )


def s15_secao_variabilidade(prs, lays):
    return slide_secao(prs, lays, "Medidas de variabilidade")


def s16_amplitude_iqr(prs, lays):
    slide = slide_com_titulo(prs, lays,
                             "Amplitude e intervalo interquartil", EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.98), FAIXA, Inches(0.40))
    escrever(tb.text_frame, [
        {"texto": "As duas medem largura e saem em reais. O que muda é quantos "
                  "casos cada uma escuta.",
         "tamanho": 14, "cor": CINZA_ESCURO, "depois": 0},
    ])

    colunas = [
        (CINZA, "Amplitude", brl(E["amplitude"], 2), "máximo menos mínimo",
         "Depende de exatamente dois acórdãos, os dois mais atípicos da base."),
        (AMARELO, "Intervalo interquartil", brl(E["iqr"]), "Q3 menos Q1",
         "A largura da metade central. Trocar o maior valor por outro dez vezes "
         "maior não a muda em nada."),
    ]
    largura = Inches(5.59)
    vao = Inches(0.32)
    y = Inches(2.50)
    altura = Inches(1.56)
    for i, (cor, titulo, valor, formula, texto) in enumerate(colunas):
        x = MARGEM + i * (largura + vao)
        caixa(slide, x, y, largura, altura, preenchimento=QUASE_BRANCO)
        caixa(slide, x, y, largura, Inches(0.09), preenchimento=cor)
        tb = texto_livre(slide, x + Inches(0.28), y + Inches(0.26),
                         largura - Inches(0.56), altura - Inches(0.44))
        escrever(tb.text_frame, [
            {"texto": f"{titulo}  ·  {formula}", "tamanho": 12,
             "cor": CINZA_ESCURO, "depois": 2},
            {"texto": valor, "fonte": DISPLAY, "tamanho": 24, "bold": True,
             "depois": 6},
            {"texto": texto, "tamanho": 12.5, "entrelinhas": 1.12, "depois": 0},
        ])

    imagem_ajustada(slide, os.path.join(FIGURAS, "amplitude_iqr.png"),
                    MARGEM, Inches(4.08), FAIXA, Inches(1.78), moldura=False)

    nota(slide, Inches(5.90),
         f"A base inteira ocupa {brl(E['amplitude'], 2)}, e a metade dela cabe "
         f"em {brl(E['iqr'])}",
         "Quando esses dois números são tão diferentes, é sinal de cauda longa: "
         "a maioria dos casos está espremida e uns poucos estão muito longe.",
         AMARELO, altura=Inches(0.80))
    return slide


def s17_variancia_desvio(prs, lays):
    slide = slide_com_titulo(prs, lays, "Variância e desvio padrão", EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.98), FAIXA, Inches(0.40))
    escrever(tb.text_frame, [
        {"texto": "As duas medem o afastamento típico em relação à média. São a "
                  "mesma informação em duas unidades diferentes.",
         "tamanho": 14, "cor": CINZA_ESCURO, "depois": 0},
    ])

    colunas = [
        (ROXO, "Variância", f"{E['variancia']:,.0f}".replace(",", "."),
         "média dos quadrados dos desvios",
         "Sai em reais ao quadrado, unidade que não existe no mundo, e por isso "
         "quase nunca se reporta."),
        (LARANJA, "Desvio padrão", brl(E["desvio"], 2), "raiz da variância",
         "A raiz devolve a unidade original, em reais, e é isso que o torna "
         "reportável."),
    ]
    largura = Inches(5.59)
    vao = Inches(0.32)
    y = Inches(2.50)
    altura = Inches(1.56)
    for i, (cor, titulo, valor, formula, texto) in enumerate(colunas):
        x = MARGEM + i * (largura + vao)
        caixa(slide, x, y, largura, altura, preenchimento=QUASE_BRANCO)
        caixa(slide, x, y, largura, Inches(0.09), preenchimento=cor)
        tb = texto_livre(slide, x + Inches(0.28), y + Inches(0.26),
                         largura - Inches(0.56), altura - Inches(0.44))
        escrever(tb.text_frame, [
            {"texto": f"{titulo}  ·  {formula}", "tamanho": 12,
             "cor": CINZA_ESCURO, "depois": 2},
            {"texto": valor, "fonte": DISPLAY, "tamanho": 24, "bold": True,
             "depois": 6},
            {"texto": texto, "tamanho": 12.5, "entrelinhas": 1.12, "depois": 0},
        ])

    imagem_ajustada(slide, os.path.join(FIGURAS, "desvio.png"),
                    MARGEM, Inches(4.08), FAIXA, Inches(1.78), moldura=False)

    nota(slide, Inches(5.90),
         f"Um desvio padrão de {brl(E['desvio'])} contra uma média de "
         f"{brl(E['media'])}",
         "Quando o desvio padrão passa da média, a faixa habitual de um desvio "
         "para baixo cai no negativo, e valor de indenização não é negativo. "
         "Isso é o aviso de que a distribuição é assimétrica e de que mediana e "
         "IIQ descrevem melhor essa coluna.",
         LARANJA, altura=Inches(0.80))
    return slide


def s18_tudo_junto(prs, lays):
    slide = slide_com_titulo(prs, lays,
                             "A mesma coluna, todas as estatísticas", EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.88), FAIXA, Inches(0.52))
    escrever(tb.text_frame, [
        {"texto": f"valor_indenizacao, os {E['n']} acórdãos que têm valor: as "
                  f"contas abaixo ignoram os {E['n_faltantes']} sem valor. "
                  f"Dizer sobre quantos casos a estatística foi calculada faz "
                  f"parte de reportá-la.",
         "tamanho": 13.5, "cor": CINZA_ESCURO, "entrelinhas": 1.1, "depois": 0},
    ])

    posicao = [
        ("mínimo", brl(E["minimo"], 2)),
        ("Q1", brl(E["q1"])),
        ("mediana", brl(E["mediana"])),
        ("moda", f"{brl(E['moda'])}  ({E['n_moda']} casos)"),
        ("média", brl(E["media"], 2)),
        ("Q3", brl(E["q3"])),
        ("máximo", brl(E["maximo"])),
        ("proporção acima de R$ 10 mil",
         f"{E['prop_acima_10k'] * 100:.1f}%".replace(".", ",") +
         f"  ({E['n_acima_10k']} casos)"),
    ]
    variabilidade = [
        ("amplitude", brl(E["amplitude"], 2)),
        ("intervalo interquartil", brl(E["iqr"])),
        ("variância", f"{E['variancia']:,.0f}".replace(",", ".") + " reais²"),
        ("desvio padrão", brl(E["desvio"], 2)),
    ]

    def bloco(x, largura, titulo, cor, itens):
        tb = texto_livre(slide, x, Inches(2.44), largura, Inches(0.26))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 11.5, "cor": CINZA_ESCURO, "depois": 0},
        ])
        tabela_shape = slide.shapes.add_table(
            len(itens), 3, x, Inches(2.62), largura,
            Inches(0.28 * len(itens))
        )
        tabela = tabela_shape.table
        tabela_sem_estilo(tabela)
        larguras = (Inches(0.10), largura - Inches(2.36), Inches(2.26))
        for lg, i in zip(larguras, range(3)):
            tabela.columns[i].width = lg
        for r, (nome, valor) in enumerate(itens):
            tabela.rows[r].height = Inches(0.28)
            zebra = QUASE_BRANCO if r % 2 == 0 else None
            celula(tabela, r, 0, "", fundo=cor if r == 0 else zebra)
            celula(tabela, r, 1, nome, tamanho=11.5, fundo=zebra)
            celula(tabela, r, 2, valor, tamanho=11.5, bold=True, fundo=zebra)

    bloco(MARGEM, Inches(5.59), "POSIÇÃO", TURQUESA, posicao)
    bloco(MARGEM + Inches(5.91), Inches(5.59), "VARIABILIDADE", LARANJA,
          variabilidade)

    # a tabela da variabilidade tem quatro linhas contra oito da posicao, e a
    # sobra do lado direito e onde cabe a conta que a pergunta de pesquisa
    # realmente pede
    x_comparar = MARGEM + Inches(5.91)
    # a tabela de cima renderiza um pouco mais alta que o pedido; a folga
    # evita que a caixa encoste na ultima linha dela
    y_comparar = Inches(3.88)
    altura_comparar = Inches(1.12)
    caixa(slide, x_comparar, y_comparar, Inches(5.59), altura_comparar,
          preenchimento=QUASE_BRANCO)
    caixa(slide, x_comparar, y_comparar, Inches(0.09), altura_comparar,
          preenchimento=VERMELHO)
    por_assunto = E["mediana_por_assunto"]
    maior, menor = por_assunto.index[0], por_assunto.index[-1]
    tb = texto_livre(slide, x_comparar + Inches(0.26), y_comparar + Inches(0.10),
                     Inches(5.10), altura_comparar - Inches(0.18))
    escrever(tb.text_frame, [
        {"texto": "A conta que a pergunta pede é uma comparação",
         "tamanho": 12, "bold": True, "depois": 2},
        {"texto": f"Quase sempre é a mesma medida separada por grupo: mediana "
                  f"por assunto dá {brl(por_assunto.iloc[0])} em “{maior}” "
                  f"contra {brl(por_assunto.iloc[-1])} em “{menor}”. É uma "
                  f"linha de pandas, e é a aula 4; hoje basta dizer a conta.",
         "tamanho": 10.5, "cor": CINZA_ESCURO, "entrelinhas": 1.06, "depois": 3},
        {"texto": 'df.groupby("assunto")["valor_indenizacao"].median()',
         "fonte": "Consolas", "tamanho": 10, "cor": CINZA_ESCURO, "depois": 0},
    ])

    imagem_ajustada(slide, os.path.join(FIGURAS, "resumo.png"),
                    MARGEM, Inches(5.06), FAIXA, Inches(1.64), moldura=False)
    return slide


def s19_mapa_mental(prs, lays):
    """Mapa mental desenhado em formas do PowerPoint.

    A arvore comeca na variavel e passa pelo tipo antes de chegar nas medidas,
    e nao o contrario: a pergunta que o aluno vai se fazer no Projeto 1 e "que
    conta eu posso fazer com esta coluna", e a resposta depende do tipo.

    Nasceu como mindmap em mermaid (fonte no fim deste arquivo), mas o mermaid
    entra no slide como imagem, com tipografia e paleta proprias. Em formas,
    fica vetorial, na fonte da marca, e editavel na hora da aula.
    """
    slide = slide_com_titulo(prs, lays, "O mapa da aula em um slide", EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.90), FAIXA, Inches(0.32))
    escrever(tb.text_frame, [
        {"texto": "Comece pela coluna, descubra o tipo, e o tipo diz quais "
                  "contas estão disponíveis.",
         "tamanho": 13.5, "cor": CINZA_ESCURO, "depois": 0},
    ])

    fio = Inches(0.03)
    x_tronco = MARGEM + Inches(1.98)
    x_tipo = MARGEM + Inches(2.20)
    largura_tipo = Inches(2.10)
    x_linha = x_tipo + largura_tipo + Inches(0.22)
    largura_linha = FAIXA - (x_linha - MARGEM)

    # (cor do tipo, nome, subtitulo, [(cor da familia, familia, medidas)])
    tipos = [
        (TURQUESA, "numérica", "contínua ou discreta", [
            (TURQUESA, "posição",
             "média · mediana · moda · mínimo · máximo · quartis"),
            (LARANJA, "variabilidade",
             "amplitude · intervalo interquartil · variância · desvio padrão"),
        ]),
        (ROXO, "categórica", "nominal, ordinal ou binária", [
            (TURQUESA, "posição",
             "moda · proporção  ·  se for ordinal, também mediana e quartis"),
            (LARANJA, "variabilidade",
             "nenhuma das quatro: o que se faz é contar cada categoria"),
        ]),
        (CINZA, "identificador", "número do processo", [
            (CINZA, "nenhuma",
             "serve para achar o caso depois, não para calcular"),
        ]),
    ]

    altura_linha = Inches(0.62)
    passo_linha = Inches(0.70)
    vao_tipo = Inches(0.30)
    altura_minima = Inches(0.76)

    # posiciona os blocos de cima para baixo e guarda o centro de cada um
    y = Inches(2.34)
    centros = []
    for cor, nome, subtitulo, familias in tipos:
        altura_bloco = max(
            passo_linha * len(familias) - (passo_linha - altura_linha),
            altura_minima,
        )
        centro = y + altura_bloco // 2
        centros.append(centro)

        caixa(slide, x_tipo, y, largura_tipo, altura_bloco,
              preenchimento=QUASE_BRANCO, borda=CINZA_CLARO)
        caixa(slide, x_tipo, y, Inches(0.07), altura_bloco, preenchimento=cor)
        tb = texto_livre(slide, x_tipo + Inches(0.22), y + Inches(0.10),
                         largura_tipo - Inches(0.36), altura_bloco - Inches(0.20))
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        escrever(tb.text_frame, [
            {"texto": nome, "fonte": DISPLAY, "tamanho": 16, "bold": True,
             "depois": 2},
            {"texto": subtitulo, "tamanho": 10.5, "cor": CINZA_ESCURO,
             "depois": 0},
        ])
        caixa(slide, x_tipo + largura_tipo, centro, x_linha - x_tipo - largura_tipo,
              fio, preenchimento=cor)

        for i, (cor_familia, familia, medidas) in enumerate(familias):
            ly = y + i * passo_linha
            caixa(slide, x_linha, ly, largura_linha, altura_linha,
                  preenchimento=QUASE_BRANCO, borda=CINZA_CLARO)
            caixa(slide, x_linha, ly, Inches(0.05), altura_linha,
                  preenchimento=cor_familia)
            tb = texto_livre(slide, x_linha + Inches(0.18), ly + Inches(0.06),
                             largura_linha - Inches(0.30), altura_linha - Inches(0.12))
            escrever(tb.text_frame, [
                {"texto": familia, "tamanho": 10.5, "cor": CINZA_ESCURO,
                 "depois": 1},
                {"texto": medidas, "tamanho": 12, "bold": True, "depois": 0},
            ])
        y = y + altura_bloco + vao_tipo

    # raiz e espinha, agora que os centros sao conhecidos
    centro_raiz = (centros[0] + centros[-1]) // 2
    altura_raiz = Inches(1.10)
    caixa(slide, MARGEM, centro_raiz - altura_raiz // 2, Inches(1.76),
          altura_raiz, preenchimento=PRETO)
    tb = texto_livre(slide, MARGEM + Inches(0.18),
                     centro_raiz - altura_raiz // 2 + Inches(0.18),
                     Inches(1.40), Inches(0.80))
    escrever(tb.text_frame, [
        {"texto": "UMA COLUNA", "tamanho": 9.5, "cor": CINZA, "depois": 4},
        {"texto": "que tipo de variável é?", "fonte": DISPLAY, "tamanho": 13,
         "bold": True, "cor": BRANCO, "entrelinhas": 1.08, "depois": 0},
    ])
    caixa(slide, MARGEM + Inches(1.76), centro_raiz, x_tronco - MARGEM - Inches(1.76),
          fio, preenchimento=CINZA)
    caixa(slide, x_tronco, centros[0], fio, centros[-1] - centros[0],
          preenchimento=CINZA)
    for centro in centros:
        caixa(slide, x_tronco, centro, x_tipo - x_tronco, fio,
              preenchimento=CINZA)

    return slide


def s20_projeto(prs, lays):
    slide = slide_com_titulo(prs, lays, "Projeto 1, na segunda hora de hoje",
                             EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(1.94), FAIXA, Inches(0.40))
    escrever(tb.text_frame, [
        {"texto": "Uma hora, em grupo. Você vai fazer o caminho inverso do "
                  "desta aula: partir de sentenças em texto e decidir quais "
                  "colunas precisam existir para que essas contas sejam "
                  "possíveis.",
         "tamanho": 14, "cor": CINZA_ESCURO, "entrelinhas": 1.12, "depois": 0},
    ])

    etapas = [
        (VERMELHO, "Escolher o tema e explorar as sentenças",
         "cinco temas, cada um com uma pergunta de pesquisa e 10 sentenças "
         "reais de primeiro grau do TJSP. Três ou quatro leituras bastam."),
        (LARANJA, "Definir as variáveis",
         "nome, tipo, papel na pergunta e de onde tirar. É a definição "
         "operacional: o texto que faria duas pessoas lerem a mesma sentença e "
         "registrarem o mesmo valor."),
        (AMARELO, "Rodar o modelo e dizer que estatísticas calcularia",
         "o modelo lê cada sentença e preenche as suas colunas. Depois você diz "
         "que contas faria, e elas precisam caber nos tipos que você declarou."),
        (TURQUESA, "Avaliar, ajustar e entregar",
         "avaliações pela rubrica, que está visível desde o começo. Entrega no "
         "BlackBoard: a tabela em CSV e a documentação em JSON."),
    ]
    lista_com_barras(slide, etapas, topo=Inches(2.44), altura=Inches(0.92),
                     passo=Inches(0.98), tamanho_titulo=16, tamanho_detalhe=12.5)

    tb = texto_livre(slide, MARGEM, Inches(6.40), FAIXA, Inches(0.28))
    escrever(tb.text_frame, [
        {"texto": "A rubrica avalia a qualidade das escolhas, e não a "
                  "quantidade de variáveis. Lista curta e bem definida vale "
                  "mais do que lista longa.",
         "tamanho": 12, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def main():
    prs = deck_limpo()
    lays = layouts(prs)

    for construir in (
        s01_capa,
        s02_secao_estatistica, s03_o_que_e, s04_coluna_vira_numero,
        s05_duas_perguntas, s06_mapa_tipos,
        s07_secao_posicao, s08_media, s09_mediana, s09b_mediana_na_mao,
        s10_moda, s11_minimo,
        s12_maximo, s13_quartis, s14_proporcao,
        s15_secao_variabilidade, s16_amplitude_iqr, s17_variancia_desvio,
        s18_tudo_junto,
        s19_mapa_mental, s20_projeto,
    ):
        construir(prs, lays)

    n = gravar(prs, SAIDA, titulo=TITULO_DECK)
    print(f"{n} slides em {SAIDA}")


# Fonte do mapa mental do slide 19, caso ele precise sair daqui (README, site,
# notebook). No deck ele e desenhado em formas, para ficar na tipografia da
# marca e editavel na hora da aula.
MERMAID = """
mindmap
  root((uma coluna:<br/>que tipo de variável é?))
    numérica<br/>contínua ou discreta
      posição
        média · mediana · moda · mínimo · máximo · quartis
      variabilidade
        amplitude · intervalo interquartil · variância · desvio padrão
    categórica<br/>nominal, ordinal ou binária
      posição
        moda · proporção
        se for ordinal, também mediana e quartis
      variabilidade
        nenhuma das quatro: o que se faz é contar cada categoria
    identificador<br/>número do processo
      nenhuma: serve para achar o caso, não para calcular
"""


if __name__ == "__main__":
    import sys

    sys.stdout.reconfigure(encoding="utf-8")
    main()

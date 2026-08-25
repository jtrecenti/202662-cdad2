"""
Monta os slides da Aula 2: tipos de variaveis.

A aula e dada no notebook (notebooks/aula02_*.ipynb); estes slides existem como
referencia e para projetar o diagrama de tipos e a tabela de papeis.

Os exemplos saem do Projeto 1 (apps/ex01-variaveis), tema "Furto e roubo", para
que o aluno reencontre em 18/08 exatamente o vocabulario visto aqui.

Uso:
    python build_aula02.py
"""

from __future__ import annotations

import os

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from insper import (
    AMARELO, BASE, BRANCO, CINZA, CINZA_CLARO, CINZA_ESCURO, DISPLAY, FAIXA,
    LARANJA, MARGEM, PRETO, QUASE_BRANCO, ROSA, ROXO, TEXTO, TOPO, TURQUESA,
    VERDE, VERMELHO, caixa, celula, deck_limpo, escrever, gravar, layouts,
    lista_com_barras, slide_capa, slide_com_titulo, slide_secao, tabela_sem_estilo,
    texto_livre,
)

AQUI = os.path.dirname(os.path.abspath(__file__))
SAIDA = os.path.join(AQUI, "aula02.pptx")

EYEBROW = "Ciência de Dados Aplicada ao Direito II  ·  Aula 2  ·  Tipos de variáveis"
TITULO_DECK = "Aula 2: tipos de variáveis"

# a cor de cada ramo do diagrama, reusada nos slides seguintes
COR_NUMERICA = TURQUESA
COR_CATEGORICA = ROXO


# ---------------------------------------------------------------- slides


def s01_capa(prs, lays):
    return slide_capa(
        prs, lays,
        subtitulo="Aula 2: tipos de variáveis",
        subtema="Da variável aleatória à coluna da tabela",
    )


def s02_secao_va(prs, lays):
    return slide_secao(prs, lays, "Antes de tipar: o que é uma variável?")


def s03_variavel_aleatoria(prs, lays):
    slide = slide_com_titulo(prs, lays, "Variável aleatória", EYEBROW)

    caixa(slide, MARGEM, Inches(2.02), FAIXA, Inches(1.32),
          preenchimento=QUASE_BRANCO)
    caixa(slide, MARGEM, Inches(2.02), Inches(0.09), Inches(1.32),
          preenchimento=VERMELHO)
    tb = texto_livre(slide, Inches(1.48), Inches(2.22), Inches(10.9), Inches(1.0))
    escrever(tb.text_frame, [
        {"texto": "Uma regra que associa um valor a cada caso do estudo, antes "
                  "de você olhar o caso.",
         "tamanho": 19, "bold": True, "depois": 5},
        {"texto": "Aleatória não quer dizer imprevisível: quer dizer que o valor "
                  "muda de caso para caso, e você não sabe qual é até medir.",
         "tamanho": 13.5, "cor": CINZA_ESCURO, "depois": 0},
    ])

    tb = texto_livre(slide, MARGEM, Inches(3.60), FAIXA, Inches(0.32))
    escrever(tb.text_frame, [
        {"texto": "Sorteie uma sentença de furto do TJSP. Antes de abrir o arquivo:",
         "tamanho": 15, "cor": CINZA_ESCURO, "depois": 0},
    ])

    exemplos = [
        (COR_NUMERICA, "X = pena aplicada, em anos",
         "você sabe que vai sair um número entre 1 e 8, mas não sabe qual"),
        (COR_CATEGORICA, "Y = houve confissão espontânea",
         "você sabe que vai sair sim ou não, mas não sabe qual"),
        (COR_CATEGORICA, "Z = regime inicial fixado",
         "você sabe que vai sair aberto, semiaberto ou fechado"),
    ]
    lista_com_barras(slide, exemplos, topo=Inches(4.02), altura=Inches(0.74),
                     passo=Inches(0.80))

    tb = texto_livre(slide, MARGEM, Inches(6.46), FAIXA, Inches(0.32))
    escrever(tb.text_frame, [
        {"texto": "X, Y e Z são as variáveis. O que aparece quando você abre "
                  "a sentença é o valor observado delas naquele caso.",
         "tamanho": 13, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def s04_variavel_valor_estatistica(prs, lays):
    slide = slide_com_titulo(
        prs, lays, "Variável, valor observado e estatística", EYEBROW
    )

    tb = texto_livre(slide, MARGEM, Inches(2.02), FAIXA, Inches(0.35))
    escrever(tb.text_frame, [
        {"texto": "Três coisas diferentes, que a conversa do dia a dia mistura.",
         "tamanho": 15, "cor": CINZA_ESCURO, "depois": 0},
    ])

    blocos = [
        ("Variável", VERMELHO, "o que se mede",
         "“pena aplicada, em anos”\n\nExiste antes dos dados. É a coluna da tabela."),
        ("Valor observado", LARANJA, "o valor de um caso",
         "na sentença 0001234-56, a pena foi 5,5 anos\n\nÉ uma célula da tabela."),
        ("Estatística", TURQUESA, "resumo de muitos valores observados",
         "a pena mediana das 120 sentenças foi 5 anos\n\nÉ um número calculado da coluna inteira."),
    ]
    largura = Inches(3.62)
    vao = Inches(0.32)
    y = Inches(2.55)
    altura = Inches(3.05)
    for i, (titulo, cor, subtitulo, corpo) in enumerate(blocos):
        x = MARGEM + i * (largura + vao)
        caixa(slide, x, y, largura, altura, preenchimento=QUASE_BRANCO)
        caixa(slide, x, y, largura, Inches(0.09), preenchimento=cor)
        tb = texto_livre(slide, x + Inches(0.26), y + Inches(0.32),
                         largura - Inches(0.52), altura - Inches(0.5))
        escrever(tb.text_frame, [
            {"texto": titulo, "fonte": DISPLAY, "tamanho": 24, "bold": True,
             "depois": 2},
            {"texto": subtitulo, "tamanho": 13, "cor": CINZA_ESCURO, "depois": 12},
            {"texto": corpo, "tamanho": 14, "entrelinhas": 1.15, "depois": 0},
        ])

    caixa(slide, MARGEM, Inches(5.92), FAIXA, Inches(0.72), borda=CINZA_CLARO)
    tb = texto_livre(slide, Inches(1.45), Inches(6.06), Inches(10.9), Inches(0.48))
    escrever(tb.text_frame, [
        {"texto": "O tipo é propriedade da variável, e não do valor observado. "
                  "Ele determina quais estatísticas podem ser calculadas a "
                  "partir daquela coluna.",
         "tamanho": 14, "cor": CINZA_ESCURO, "entrelinhas": 1.15, "depois": 0},
    ])
    return slide


def s05_secao_tipos(prs, lays):
    return slide_secao(prs, lays, "Tipos de variáveis")


def s06_diagrama(prs, lays):
    """O diagrama de tipos. É o slide que a aula projeta e volta a projetar."""
    slide = slide_com_titulo(prs, lays, "Os cinco tipos", EYEBROW)

    def bloco(x, y, cx, cy, texto, cor, cor_texto, tamanho):
        sh = caixa(slide, Inches(x), Inches(y), Inches(cx), Inches(cy),
                   preenchimento=cor)
        sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        escrever(sh.text_frame, [
            {"texto": texto, "fonte": DISPLAY, "tamanho": tamanho, "bold": True,
             "cor": cor_texto, "alinhamento": PP_ALIGN.CENTER, "depois": 0},
        ])
        return sh

    def linha_v(x, y0, y1, cor=CINZA):
        caixa(slide, Inches(x - 0.01), Inches(y0), Inches(0.02),
              Inches(y1 - y0), preenchimento=cor)

    def linha_h(x0, x1, y, cor=CINZA):
        caixa(slide, Inches(x0), Inches(y - 0.01), Inches(x1 - x0),
              Inches(0.02), preenchimento=cor)

    # raiz
    bloco(5.40, 2.15, 2.50, 0.62, "VARIÁVEL", VERMELHO, BRANCO, 18)
    linha_v(6.65, 2.77, 3.00)
    linha_h(3.90, 9.40, 3.00)
    linha_v(3.90, 3.00, 3.25)
    linha_v(9.40, 3.00, 3.25)

    # ramos
    bloco(2.20, 3.25, 3.40, 0.62, "NUMÉRICA", COR_NUMERICA, PRETO, 17)
    bloco(7.70, 3.25, 3.40, 0.62, "CATEGÓRICA", COR_CATEGORICA, BRANCO, 17)

    linha_v(3.90, 3.87, 4.10)
    linha_h(2.925, 4.875, 4.10)
    linha_v(9.40, 3.87, 4.10)
    linha_h(7.45, 11.35, 4.10)

    folhas = [
        (2.075, "discreta", COR_NUMERICA, PRETO,
         "conta\n\nnº de réus\nnº de testemunhas"),
        (4.025, "contínua", COR_NUMERICA, PRETO,
         "mede numa escala\n\npena em anos\nvalor em reais"),
        (6.60, "nominal", COR_CATEGORICA, BRANCO,
         "rótulo sem ordem\n\ncomarca\nfundamento"),
        (8.55, "ordinal", COR_CATEGORICA, BRANCO,
         "rótulo com ordem\n\nregime: aberto,\nsemiaberto, fechado"),
        (10.50, "binária", COR_CATEGORICA, BRANCO,
         "só dois valores\n\nhouve confissão?\nsim ou não"),
    ]
    for x, nome, cor, cor_texto, detalhe in folhas:
        linha_v(x + 0.85, 4.10, 4.35)
        bloco(x, 4.35, 1.70, 0.58, nome, cor, cor_texto, 15)
        tb = texto_livre(slide, Inches(x - 0.10), Inches(5.02), Inches(1.90),
                         Inches(1.15))
        escrever(tb.text_frame, [
            {"texto": detalhe, "tamanho": 10.5, "cor": CINZA_ESCURO,
             "entrelinhas": 1.12, "alinhamento": PP_ALIGN.CENTER, "depois": 0},
        ])

    tb = texto_livre(slide, MARGEM, Inches(6.30), FAIXA, Inches(0.4))
    escrever(tb.text_frame, [
        {"texto": "Fora da árvore: o identificador (número do processo) e o "
                  "texto livre. Nenhum dos dois entra em conta; eles servem "
                  "para achar o caso e para construir as outras variáveis.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "entrelinhas": 1.1, "depois": 0},
    ])
    return slide


def s07_numerica(prs, lays):
    slide = slide_com_titulo(prs, lays, "Numérica: discreta e contínua", EYEBROW)

    largura = Inches(5.59)
    vao = Inches(0.32)
    y = Inches(2.05)
    altura = Inches(3.55)

    colunas = [
        ("Discreta", "resultado de contar",
         [
             "número de réus no processo",
             "número de testemunhas ouvidas",
             "número de assuntos do processo",
             "quantidade de invólucros apreendidos",
         ],
         "Entre 2 e 3 réus não existe nada. Some sempre em passos inteiros."),
        ("Contínua", "resultado de medir numa escala",
         [
             "pena aplicada, em anos",
             "valor da indenização, em reais",
             "massa da droga, em gramas",
             "tempo até a sentença, em dias",
         ],
         "Entre 5,0 e 5,5 anos existe 5,25. O limite é a precisão do "
         "instrumento, não a variável."),
    ]
    for i, (titulo, subtitulo, itens, nota) in enumerate(colunas):
        x = MARGEM + i * (largura + vao)
        caixa(slide, x, y, largura, altura, preenchimento=QUASE_BRANCO)
        caixa(slide, x, y, largura, Inches(0.09), preenchimento=COR_NUMERICA)
        tb = texto_livre(slide, x + Inches(0.28), y + Inches(0.32),
                         largura - Inches(0.56), altura - Inches(0.5))
        linhas = [
            {"texto": titulo, "fonte": DISPLAY, "tamanho": 24, "bold": True,
             "depois": 2},
            {"texto": subtitulo, "tamanho": 13, "cor": CINZA_ESCURO, "depois": 12},
        ]
        linhas += [
            {"texto": f"·  {item}", "tamanho": 14, "depois": 5} for item in itens
        ]
        linhas.append({"texto": nota, "tamanho": 12.5, "cor": CINZA_ESCURO,
                       "antes": 10, "entrelinhas": 1.1, "depois": 0})
        escrever(tb.text_frame, linhas)

    caixa(slide, MARGEM, Inches(5.78), FAIXA, Inches(0.90), borda=CINZA_CLARO)
    tb = texto_livre(slide, Inches(1.45), Inches(5.94), Inches(10.9), Inches(0.62))
    escrever(tb.text_frame, [
        {"texto": "Quando a distinção importa",
         "tamanho": 14, "bold": True, "depois": 3},
        {"texto": "Discreta com muitos valores possíveis, como dias de "
                  "tramitação, é tratada como contínua sem prejuízo. Já uma "
                  "contagem de 0, 1 e 2 não deve ser tratada como escala.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def s08_categorica(prs, lays):
    slide = slide_com_titulo(
        prs, lays, "Categórica: nominal, ordinal e binária", EYEBROW
    )

    largura = Inches(3.62)
    vao = Inches(0.32)
    y = Inches(2.05)
    altura = Inches(3.60)

    colunas = [
        ("Nominal", "categorias sem ordem",
         "comarca\nfundamento invocado\nclasse processual\ntipo de prova",
         "Não dá para dizer que São Paulo é maior que Santos."),
        ("Ordinal", "categorias com ordem",
         "regime de pena\ngrau de jurisdição\nfaixa de valor\nnível de sigilo",
         "A ordem existe, mas a distância entre duas categorias não é medível."),
        ("Binária", "só dois valores",
         "houve confissão?\nréu é primário?\nrecurso foi provido?\nhouve dano moral?",
         "Caso especial de nominal. A média de uma binária é a proporção."),
    ]
    for i, (titulo, subtitulo, itens, nota) in enumerate(colunas):
        x = MARGEM + i * (largura + vao)
        caixa(slide, x, y, largura, altura, preenchimento=QUASE_BRANCO)
        caixa(slide, x, y, largura, Inches(0.09), preenchimento=COR_CATEGORICA)
        tb = texto_livre(slide, x + Inches(0.26), y + Inches(0.32),
                         largura - Inches(0.52), altura - Inches(0.5))
        escrever(tb.text_frame, [
            {"texto": titulo, "fonte": DISPLAY, "tamanho": 24, "bold": True,
             "depois": 2},
            {"texto": subtitulo, "tamanho": 13, "cor": CINZA_ESCURO, "depois": 12},
            {"texto": itens, "tamanho": 14, "entrelinhas": 1.2, "depois": 12},
            {"texto": nota, "tamanho": 12.5, "cor": CINZA_ESCURO,
             "entrelinhas": 1.1, "depois": 0},
        ])

    caixa(slide, MARGEM, Inches(5.78), FAIXA, Inches(0.90), borda=CINZA_CLARO)
    tb = texto_livre(slide, Inches(1.45), Inches(5.94), Inches(10.9), Inches(0.62))
    escrever(tb.text_frame, [
        {"texto": "Requisito das categorias",
         "tamanho": 14, "bold": True, "depois": 3},
        {"texto": "As categorias devem ser mutuamente exclusivas e exaustivas: "
                  "cada caso pertence a exatamente uma delas. Quando não são, a "
                  "variável não pode ser tabulada nem comparada entre grupos.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def s09_numero_que_nao_e_numero(prs, lays):
    slide = slide_com_titulo(prs, lays, "Número que não é número", EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(2.02), FAIXA, Inches(0.35))
    escrever(tb.text_frame, [
        {"texto": "O erro mais comum da disciplina: o pandas lê a coluna como "
                  "número, e o número não mede nada.",
         "tamanho": 15, "cor": CINZA_ESCURO, "depois": 0},
    ])

    dados = [
        ("0001234-56.2024.8.26.0100", "número do processo", "identificador",
         "É a chave do caso, e não uma medida."),
        ("3550308", "código do município (IBGE)", "identificador",
         "São dígitos que nomeiam São Paulo."),
        ("13", "número da câmara julgadora", "identificador",
         "A 13ª não é maior que a 2ª, e o número só nomeia o órgão."),
        ("2", "polo: 1 ativo, 2 passivo", "categórica nominal",
         "Aqui o número agrupa, em vez de nomear."),
        ("1116", "código da classe processual", "identificador",
         "Rótulo do CNJ para execução fiscal."),
        ("3", "número de réus no polo passivo", "numérica discreta",
         "Esta é contagem, e soma faz sentido."),
    ]

    tabela_shape = slide.shapes.add_table(
        len(dados) + 1, 4, MARGEM, Inches(2.50), FAIXA, Inches(3.27)
    )
    tabela = tabela_shape.table
    tabela_sem_estilo(tabela)
    for largura, i in zip(
        (Inches(3.05), Inches(2.85), Inches(2.30), Inches(3.30)), range(4)
    ):
        tabela.columns[i].width = largura
    tabela.rows[0].height = Inches(0.30)
    for r in range(1, len(dados) + 1):
        tabela.rows[r].height = Inches(0.42)

    for c, titulo in enumerate(("o que aparece", "o que é", "tipo", "por quê")):
        celula(tabela, 0, c, titulo, tamanho=11, bold=True, cor=BRANCO, fundo=PRETO)
    for r, (valor, oque, tipo, porque) in enumerate(dados, start=1):
        zebra = QUASE_BRANCO if r % 2 == 0 else None
        eh_numero = tipo.startswith("numérica")
        celula(tabela, r, 0, valor, tamanho=11, fundo=zebra)
        celula(tabela, r, 1, oque, tamanho=11, fundo=zebra)
        # turquesa sobre branco nao passa em contraste; o destaque vem do peso
        celula(tabela, r, 2, tipo, tamanho=11, bold=True, fundo=zebra,
               cor=PRETO if eh_numero else CINZA_ESCURO)
        celula(tabela, r, 3, porque, tamanho=11, cor=CINZA_ESCURO, fundo=zebra)

    caixa(slide, MARGEM, Inches(5.62), FAIXA, Inches(1.02),
          preenchimento=QUASE_BRANCO)
    caixa(slide, MARGEM, Inches(5.62), Inches(0.09), Inches(1.02),
          preenchimento=VERMELHO)
    tb = texto_livre(slide, Inches(1.48), Inches(5.80), Inches(10.9), Inches(0.75))
    escrever(tb.text_frame, [
        {"texto": "Como decidir", "tamanho": 15, "bold": True, "depois": 3},
        {"texto": "Verifique se somar dois valores da variável produz algo com "
                  "sentido. Se não produz, a variável não é numérica, por mais "
                  "que seja escrita com dígitos.",
         "tamanho": 14, "cor": CINZA_ESCURO, "entrelinhas": 1.1, "depois": 0},
    ])
    return slide


def s10_secao_papeis(prs, lays):
    return slide_secao(prs, lays, "Papéis das variáveis na pergunta")


def s11_papeis(prs, lays):
    slide = slide_com_titulo(prs, lays, "Os quatro papéis", EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(2.02), FAIXA, Inches(0.35))
    escrever(tb.text_frame, [
        {"texto": "O papel não é propriedade da variável: é a posição dela na "
                  "sua pergunta. A mesma coluna muda de papel quando a pergunta "
                  "muda.",
         "tamanho": 14, "cor": CINZA_ESCURO, "depois": 0},
    ])

    papeis = [
        (CINZA_ESCURO, "Identifica o caso",
         "acha o processo depois e permite conferir e juntar com outra base. "
         "Não entra na comparação."),
        (VERMELHO, "O que eu quero explicar",
         "o desfecho da pergunta, aquilo que você quer entender por que varia. "
         "Em estatística: variável resposta, ou dependente."),
        (TURQUESA, "O que pode explicar",
         "aquilo que talvez esteja associado ao desfecho. Em estatística: "
         "variável explicativa, ou independente."),
        (AMARELO, "Define quais casos entram",
         "não responde à pergunta, mas delimita o universo: só sentenças de "
         "mérito, só condenações, só um período."),
    ]
    lista_com_barras(slide, papeis, topo=Inches(2.52), altura=Inches(0.88),
                     passo=Inches(0.95), tamanho_titulo=17, tamanho_detalhe=13)

    caixa(slide, MARGEM, Inches(6.36), FAIXA, Inches(0.31), borda=CINZA_CLARO)
    tb = texto_livre(slide, Inches(1.45), Inches(6.40), Inches(10.9), Inches(0.25))
    escrever(tb.text_frame, [
        {"texto": "Sem uma variável no papel de desfecho não há o que comparar. "
                  "A rubrica do Projeto 1 verifica isso.",
         "tamanho": 12, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def s12_exemplo(prs, lays):
    slide = slide_com_titulo(prs, lays, "Da sentença para a tabela", EYEBROW)

    caixa(slide, MARGEM, Inches(1.98), FAIXA, Inches(0.88),
          preenchimento=QUASE_BRANCO)
    caixa(slide, MARGEM, Inches(1.98), Inches(0.09), Inches(0.88),
          preenchimento=VERMELHO)
    tb = texto_livre(slide, Inches(1.48), Inches(2.10), Inches(10.9), Inches(0.66))
    escrever(tb.text_frame, [
        {"texto": "Pergunta de pesquisa",
         "tamanho": 11.5, "cor": CINZA_ESCURO, "depois": 3},
        {"texto": "Nas ações de despejo por falta de pagamento, o tamanho do "
                  "débito está associado à concessão da liminar de desocupação?",
         "tamanho": 15, "bold": True, "depois": 0},
    ])

    linhas = [
        ("numero_processo", "identificador", None, "identifica o caso",
         "cabeçalho da sentença"),
        ("tipo_imovel", "categórica nominal", COR_CATEGORICA,
         "o que pode explicar", "residencial ou comercial"),
        ("meses_de_atraso", "numérica discreta", COR_NUMERICA,
         "o que pode explicar", "quantos aluguéis vencidos"),
        ("valor_do_debito", "numérica contínua", COR_NUMERICA,
         "o que pode explicar", "valor em reais na inicial"),
        ("faixa_do_debito", "categórica ordinal", COR_CATEGORICA,
         "o que pode explicar", "derivada do valor, em três faixas"),
        ("liminar_concedida", "categórica binária", COR_CATEGORICA,
         "o que quero explicar", "sim ou não, no despacho inicial"),
        ("houve_purgacao", "categórica binária", COR_CATEGORICA,
         "define quais casos entram", "só entram os que não purgaram a mora"),
    ]

    tabela_shape = slide.shapes.add_table(
        len(linhas) + 1, 5, MARGEM, Inches(3.06), FAIXA, Inches(2.82)
    )
    tabela = tabela_shape.table
    tabela_sem_estilo(tabela)
    for largura, i in zip(
        (Inches(0.13), Inches(2.55), Inches(2.35), Inches(2.85), Inches(3.62)),
        range(5),
    ):
        tabela.columns[i].width = largura
    tabela.rows[0].height = Inches(0.28)
    for r in range(1, len(linhas) + 1):
        tabela.rows[r].height = Inches(0.32)

    for c, titulo in enumerate(("", "variável", "tipo", "papel", "de onde tirar")):
        celula(tabela, 0, c, titulo, tamanho=11, bold=True, cor=BRANCO, fundo=PRETO)
    for r, (nome, tipo, cor, papel, origem) in enumerate(linhas, start=1):
        zebra = QUASE_BRANCO if r % 2 == 0 else None
        celula(tabela, r, 0, "", fundo=cor or CINZA)
        celula(tabela, r, 1, nome, tamanho=11, bold=True, fundo=zebra)
        celula(tabela, r, 2, tipo, tamanho=11, fundo=zebra, cor=CINZA_ESCURO)
        celula(tabela, r, 3, papel, tamanho=11, fundo=zebra)
        celula(tabela, r, 4, origem, tamanho=11, cor=CINZA_ESCURO, fundo=zebra)

    tb = texto_livre(slide, MARGEM, Inches(6.12), FAIXA, Inches(0.52))
    escrever(tb.text_frame, [
        {"texto": "Uma tabela com esta estrutura, nome, tipo, papel e de onde "
                  "tirar, é o que o Projeto 1 pede em 18/08. O tipo declarado "
                  "define o formato do valor que o modelo de linguagem devolve: "
                  "numérica sai como número, categórica sai como uma das "
                  "categorias listadas.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "entrelinhas": 1.12, "depois": 0},
    ])
    return slide


def s13_tipo_decide_a_conta(prs, lays):
    slide = slide_com_titulo(prs, lays, "O tipo decide a conta", EYEBROW)

    dados = [
        ("numérica contínua", COR_NUMERICA, "média, mediana, quantis",
         "desvio padrão, IQR", "reportar só a média quando a distribuição é torta"),
        ("numérica discreta", COR_NUMERICA, "média, mediana, moda",
         "desvio padrão, IQR", "estranhar que a média dê 1,7 réu"),
        ("categórica ordinal", COR_CATEGORICA, "mediana, moda",
         "amplitude de postos", "média das categorias"),
        ("categórica nominal", COR_CATEGORICA, "moda", "nenhuma clássica",
         "média, mediana, desvio padrão"),
        ("categórica binária", COR_CATEGORICA, "proporção (= média)",
         "raiz de p(1 - p)", "tratar como se fosse contínua"),
        ("identificador", CINZA, "nenhuma", "nenhuma", "qualquer conta"),
    ]

    tabela_shape = slide.shapes.add_table(
        len(dados) + 1, 5, MARGEM, Inches(2.15), FAIXA, Inches(3.05)
    )
    tabela = tabela_shape.table
    tabela_sem_estilo(tabela)
    for largura, i in zip(
        (Inches(0.13), Inches(2.62), Inches(2.70), Inches(2.35), Inches(3.70)),
        range(5),
    ):
        tabela.columns[i].width = largura
    tabela.rows[0].height = Inches(0.30)
    for r in range(1, len(dados) + 1):
        tabela.rows[r].height = Inches(0.42)

    cabecalho = ("", "tipo", "posição", "dispersão", "o que não fazer")
    for c, titulo in enumerate(cabecalho):
        celula(tabela, 0, c, titulo, tamanho=11, bold=True, cor=BRANCO, fundo=PRETO)
    for r, (tipo, cor, posicao, dispersao, nao) in enumerate(dados, start=1):
        zebra = QUASE_BRANCO if r % 2 == 0 else None
        celula(tabela, r, 0, "", fundo=cor)
        celula(tabela, r, 1, tipo, tamanho=11, bold=True, fundo=zebra)
        celula(tabela, r, 2, posicao, tamanho=11, fundo=zebra)
        celula(tabela, r, 3, dispersao, tamanho=11, fundo=zebra, cor=CINZA_ESCURO)
        celula(tabela, r, 4, nao, tamanho=11, cor=VERMELHO, fundo=zebra)

    caixa(slide, MARGEM, Inches(5.48), FAIXA, Inches(1.16),
          preenchimento=QUASE_BRANCO)
    caixa(slide, MARGEM, Inches(5.48), Inches(0.09), Inches(1.16),
          preenchimento=TURQUESA)
    tb = texto_livre(slide, Inches(1.48), Inches(5.66), Inches(10.9), Inches(0.9))
    escrever(tb.text_frame, [
        {"texto": "A média de uma binária é a proporção",
         "tamanho": 15, "bold": True, "depois": 3},
        {"texto": "Guardada como sim/não, a variável vale 1 e 0. Somar dá quantos "
                  "sim; dividir pelo total dá a fração de sim. Por isso "
                  "df[\"houve_confissao\"].mean() devolve a proporção de "
                  "confissões, e a mesma conta funciona dentro de um groupby.",
         "tamanho": 13, "cor": CINZA_ESCURO, "entrelinhas": 1.12, "depois": 0},
    ])
    return slide


def s14_proximos_passos(prs, lays):
    slide = slide_com_titulo(prs, lays, "O que vem agora", EYEBROW)

    itens = [
        (TURQUESA, "Hoje, no notebook",
         "ler uma base do DataJud, dizer o tipo de cada coluna e converter o que "
         "o pandas entendeu errado, inclusive com pd.Categorical."),
        (LARANJA, "Aula 3, em 18/08",
         "medidas de posição e de dispersão, e definição de variáveis a partir "
         "do texto da decisão."),
        (VERMELHO, "Projeto 1, em 18/08",
         "em grupo: escolher um tema, definir as variáveis com tipo e papel, e "
         "rodar o modelo de linguagem para preencher a tabela."),
        (ROXO, "No resto da disciplina",
         "a escolha do gráfico, do teste de hipótese e do modelo depende do tipo "
         "das variáveis envolvidas."),
    ]
    lista_com_barras(slide, itens, topo=Inches(2.35), altura=Inches(0.95),
                     passo=Inches(1.05), tamanho_titulo=17, tamanho_detalhe=13)
    return slide


# ---------------------------------------------------------------- main


def main():
    prs = deck_limpo()
    lays = layouts(prs)

    for construir in (
        s01_capa, s02_secao_va, s03_variavel_aleatoria,
        s04_variavel_valor_estatistica, s05_secao_tipos, s06_diagrama,
        s07_numerica, s08_categorica, s09_numero_que_nao_e_numero,
        s10_secao_papeis, s11_papeis, s12_exemplo, s13_tipo_decide_a_conta,
        s14_proximos_passos,
    ):
        construir(prs, lays)

    gravar(prs, SAIDA, titulo=TITULO_DECK)


if __name__ == "__main__":
    import sys

    sys.stdout.reconfigure(encoding="utf-8")
    main()

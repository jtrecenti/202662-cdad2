"""Monta os slides da Aula 5: gramatica de dados e gramatica de graficos.

O deck tem duas metades com funcoes diferentes:

  1. os 25 minutos expositivos (slides 2 a 13), que sao a unica exposicao da
     aula: o pipeline do pandas com os cinco verbos, a definicao de grafico
     estatistico, a gramatica e um grafico montado camada por camada;
  2. os slides operacionais da gincana (14 em diante), que ficam projetados
     durante a dinamica: uma pergunta por rodada, e a resposta logo depois.

As rodadas nao sao escritas aqui: vem de `atividades/gincana.py`, que e a fonte
unica da dinamica e cujos gabaritos rodam de verdade sobre a base.

Sai em tres versoes, do mesmo codigo:

    aula05.pptx          fundo claro, completo (professor)
    aula05_dark.pptx     fundo escuro, completo (professor)
    aula05_alunos.pptx   SEM os gabaritos das rodadas, para distribuir antes
                         da aula. Os aquecimentos continuam resolvidos.

A versao escura nao e um filtro: e a mesma estrutura com outra paleta, montada
pelo dicionario TEMA. Nenhuma cor e escrita direto nas funcoes de slide.

Uso:
    python build_aula05.py            as tres versoes
    python build_aula05.py --claro    so a clara (mais a dos alunos)
    python build_aula05.py --escuro   so a escura (mais a dos alunos)
    python build_aula05.py --alunos   so a dos alunos
"""

from __future__ import annotations

import os
import sys

from pptx.util import Inches, Pt

sys.path.insert(0, os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "..", "atividades"))

from gincana import (AQUECIMENTO, TUDO, codigo_rodada, eh_aquecimento,  # noqa: E402
                     n_cartas, rotulo)

from insper import (  # noqa: E402
    BRANCO, CINZA, CINZA_CLARO, CINZA_ESCURO, FAIXA, LARANJA, MARGEM, PRETO,
    QUASE_BRANCO, ROXO, TURQUESA, VERDE, VERMELHO, caixa, celula, deck_limpo,
    escrever, gravar, imagem_ajustada, layouts, slide_capa, slide_com_titulo,
    tabela_sem_estilo, texto_livre,
)
from pptx.dml.color import RGBColor  # noqa: E402

AQUI = os.path.dirname(os.path.abspath(__file__))
FIGURAS = os.path.join(AQUI, "..", "atividades", "_saida", "gincana", "figuras")
LOGO_BRANCO = os.path.join(AQUI, "assets", "insper_branco.png")

EYEBROW = ("Ciência de Dados Aplicada ao Direito II  ·  Aula 5  ·  "
           "Pipeline e gramática de gráficos")

MONO = "Consolas"

# posicao do logotipo no layout "Slide padrao 1", em EMU
LOGO = (10521031, 6249245, 1173011, 531627)


# ---------------------------------------------------------------------- tema
#
# Todo slide pinta por papel, nunca por cor literal. Trocar o dicionario troca
# o deck inteiro, e e o que faz a versao escura existir sem duplicar codigo.

GRAFITE = RGBColor(0x14, 0x16, 0x1A)
GRAFITE_CLARO = RGBColor(0x22, 0x25, 0x2B)
CINZA_TEXTO = RGBColor(0xC9, 0xCC, 0xD2)
VERMELHO_CLARO = RGBColor(0xFF, 0x4D, 0x4D)
TURQUESA_CLARO = RGBColor(0x5C, 0xE0, 0xB8)
ROXO_CLARO = RGBColor(0xB4, 0x7D, 0xE0)
LARANJA_CLARO = RGBColor(0xFF, 0xB2, 0x6B)
VERDE_CLARO = RGBColor(0xA9, 0xE0, 0x6E)

CLARO = {
    "escuro": False,
    "fundo": None,
    "titulo": PRETO,
    "texto": PRETO,
    "fraco": CINZA_ESCURO,
    "tenue": CINZA,
    "superficie": QUASE_BRANCO,
    "superficie_texto": PRETO,
    "superficie_fraco": CINZA_ESCURO,
    "inverso": PRETO,
    "inverso_texto": BRANCO,
    "marca": VERMELHO,
    "acento1": TURQUESA,
    "acento2": ROXO,
    "acento3": LARANJA,
    "ok": VERDE,
    "etiqueta_texto": BRANCO,
    "codigo_fundo": QUASE_BRANCO,
    "codigo_texto": PRETO,
    "figura": "",
}

ESCURO = {
    "escuro": True,
    "fundo": GRAFITE,
    "titulo": BRANCO,
    "texto": BRANCO,
    "fraco": CINZA_TEXTO,
    "tenue": CINZA_ESCURO,
    "superficie": GRAFITE_CLARO,
    "superficie_texto": BRANCO,
    "superficie_fraco": CINZA_TEXTO,
    "inverso": BRANCO,
    "inverso_texto": GRAFITE,
    "marca": VERMELHO_CLARO,
    "acento1": TURQUESA_CLARO,
    "acento2": ROXO_CLARO,
    "acento3": LARANJA_CLARO,
    "ok": VERDE_CLARO,
    "etiqueta_texto": GRAFITE,
    "codigo_fundo": GRAFITE_CLARO,
    "codigo_texto": BRANCO,
    "figura": "_dark",
}

T = CLARO


def fig(nome: str) -> str:
    """O caminho da figura, na variante do tema."""
    return os.path.join(FIGURAS, f"{nome}{T['figura']}.png")


# ------------------------------------------------------------------ utilidades

def novo_slide(prs, lays, titulo, eyebrow):
    """Slide com titulo, ja com o fundo e as cores do tema aplicados."""
    slide = slide_com_titulo(prs, lays, titulo, eyebrow)

    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if not ph.text_frame.paragraphs[0].runs:
            continue
        run = ph.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = T["titulo"] if idx == 0 else T["fraco"]

    if T["escuro"]:
        pintar_fundo(slide)
    return slide


def pintar_fundo(slide):
    """Fundo escuro atras de tudo, mais o logotipo em negativo.

    O fundo do slide fica *atras* das formas do layout, e o logotipo do layout
    e preto: no escuro ele sumiria. Por isso ele e recoberto por uma versao em
    branco, na mesma posicao.
    """
    fundo = slide.background
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = T["fundo"]

    x, y, cx, cy = LOGO
    slide.shapes.add_picture(LOGO_BRANCO, x, y, cx, cy)


def bloco_codigo(slide, x, y, cx, cy, linhas, *, tamanho=14, barra=None):
    """Caixa com codigo em monoespacada, e uma barra colorida a esquerda."""
    caixa(slide, x, y, cx, cy, preenchimento=T["codigo_fundo"])
    if barra is not None:
        caixa(slide, x, y, Inches(0.07), cy, preenchimento=barra)
    tb = texto_livre(slide, x + Inches(0.26), y + Inches(0.14),
                     cx - Inches(0.42), cy - Inches(0.22))
    escrever(tb.text_frame, linhas, fonte=MONO, tamanho=tamanho,
             cor=T["codigo_texto"], entrelinhas=1.16, espaco_depois=0)
    return tb


def legenda(slide, x, y, cx, texto, *, cor=None, tamanho=13):
    tb = texto_livre(slide, x, y, cx, Inches(0.3))
    escrever(tb.text_frame, [texto], tamanho=tamanho,
             cor=cor or T["fraco"], espaco_depois=0)
    return tb


def etiqueta(slide, x, y, cx, cy, texto, cor, *, tamanho=11.5):
    sh = caixa(slide, x, y, cx, cy, preenchimento=cor)
    escrever(sh.text_frame, [texto], tamanho=tamanho, cor=T["etiqueta_texto"],
             bold=True, espaco_depois=0)
    return sh


def painel(slide, x, y, cx, cy, linhas, *, inverso=False, recuo=0.30):
    """Caixa de destaque, clara ou invertida, com paragrafos ja formatados."""
    fundo = T["inverso"] if inverso else T["superficie"]
    base = T["inverso_texto"] if inverso else T["superficie_texto"]
    caixa(slide, x, y, cx, cy, preenchimento=fundo)
    tb = texto_livre(slide, x + Inches(recuo), y + Inches(0.16),
                     cx - Inches(recuo * 2), cy - Inches(0.24))
    pronto = []
    for item in linhas:
        cfg = dict(item)
        cfg.setdefault("cor", base)
        pronto.append(cfg)
    escrever(tb.text_frame, pronto)
    return tb


def tabela(slide, x, y, cx, cy, cabecalho, linhas, larguras, alt_linha):
    shape = slide.shapes.add_table(len(linhas) + 1, len(cabecalho), x, y, cx, cy)
    tab = shape.table
    tabela_sem_estilo(tab)
    for larg, i in zip(larguras, range(len(cabecalho))):
        tab.columns[i].width = larg
    tab.rows[0].height = Inches(0.32)
    for r in range(1, len(linhas) + 1):
        tab.rows[r].height = alt_linha
    for c, texto in enumerate(cabecalho):
        celula(tab, 0, c, texto, tamanho=11, bold=True,
               cor=T["inverso_texto"], fundo=T["inverso"])
    return tab


def zebra(r):
    return T["superficie"] if r % 2 == 0 else None


# ------------------------------------------------------------ parte expositiva

def s01_capa(prs, lays):
    slide = slide_capa(
        prs, lays,
        subtitulo="Aula 5: o pipeline e a gramática de gráficos",
        subtema="E a Gincana do Pipeline",
    )
    if T["escuro"]:
        fundo = slide.background
        fundo.fill.solid()
        fundo.fill.fore_color.rgb = T["fundo"]
        for ph in list(slide.placeholders):
            for p in ph.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = BRANCO
    return slide


def s02_plano(prs, lays):
    slide = novo_slide(prs, lays, "O plano de hoje", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "A aula 4 ficou nas medidas resumo e sobrou a manipulação de dados. "
            "Hoje ela sai, junto com os gráficos, e quase toda a aula é de pé.")

    itens = [
        ("25 min", "Slides", "os cinco verbos do pandas e a gramática de "
         "gráficos. É a única exposição da aula.", T["marca"]),
        ("55 min", "Gincana do Pipeline", "três aquecimentos e seis rodadas em "
         "grupo, com cartas de papel. Todo grupo que acertar pontua.",
         T["acento1"]),
        ("15 min", "Notebook", "vocês digitam os pipelines que montaram no "
         "papel e veem rodar.", T["acento2"]),
        ("8 min", "Fechamento", "prêmios e o que cai no Projeto 02, quinta.",
         T["fraco"]),
    ]
    y = Inches(2.62)
    for tempo, titulo, texto, cor in itens:
        etiqueta(slide, MARGEM, y, Inches(1.05), Inches(0.42), tempo, cor)
        tb = texto_livre(slide, MARGEM + Inches(1.30), y - Inches(0.02),
                         FAIXA - Inches(1.30), Inches(0.9))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 17, "bold": True, "depois": 1,
             "cor": T["texto"]},
            {"texto": texto, "tamanho": 13, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(1.02)
    return slide


def s03_pergunta(prs, lays):
    slide = novo_slide(prs, lays, "A pergunta de hoje", EYEBROW)

    painel(slide, MARGEM, Inches(2.05), FAIXA, Inches(1.05), [
        {"texto": "Nas apelações criminais do TJSP, a proporção de acórdãos "
                  "que mencionam reincidência varia conforme o regime inicial "
                  "fixado?", "tamanho": 19, "depois": 0},
    ], recuo=0.42)
    caixa(slide, MARGEM, Inches(2.05), Inches(0.09), Inches(1.05),
          preenchimento=T["marca"])

    legenda(slide, MARGEM, Inches(3.45), FAIXA,
            "A base tem 475 acórdãos. Três colunas resolvem a pergunta:")

    linhas = [
        ("regime_inicial", "aberto, semiaberto ou fechado, lido da ementa",
         "vazio em 142 acórdãos"),
        ("houve_reincidencia", "verdadeiro ou falso, lido da ementa", "completo"),
        ("pena_anos", "a pena em anos, lida da ementa",
         "vazio em 55%, e com valores implausíveis"),
    ]
    tab = tabela(slide, MARGEM, Inches(3.90), FAIXA, Inches(1.75),
                 ["coluna", "o que é", "o estado dela"], linhas,
                 [Inches(3.20), Inches(4.90), Inches(3.40)], Inches(0.47))
    for r, (nome, oque, estado) in enumerate(linhas, start=1):
        z = zebra(r)
        celula(tab, r, 0, nome, tamanho=12.5, bold=True, cor=T["texto"], fundo=z)
        celula(tab, r, 1, oque, tamanho=12.5, cor=T["texto"], fundo=z)
        celula(tab, r, 2, estado, tamanho=12.5, cor=T["fraco"], fundo=z)

    legenda(slide, MARGEM, Inches(5.85), FAIXA,
            "Guardem a terceira linha. A pena foi lida pegando o primeiro número "
            "seguido de \"anos\", e às vezes o número errado. Isso aparece já no "
            "aquecimento.", cor=T["marca"])
    return slide


def s04_jeito_antigo(prs, lays):
    slide = novo_slide(prs, lays, "Do jeito da aula 3", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "Uma variável nova a cada operação. Funciona, e responde a pergunta:")

    bloco_codigo(slide, MARGEM, Inches(2.45), Inches(7.30), Inches(2.65), [
        'apelacoes  = criminal[criminal["classe"] == "Apelação Criminal"]',
        'com_regime = apelacoes.dropna(subset=["regime_inicial"])',
        'fechado    = com_regime[com_regime["regime_inicial"] == "fechado"]',
        'semiaberto = com_regime[com_regime["regime_inicial"] == "semiaberto"]',
        'aberto     = com_regime[com_regime["regime_inicial"] == "aberto"]',
        "",
        "pd.Series({",
        '    "fechado":    fechado["houve_reincidencia"].mean(),',
        '    "semiaberto": semiaberto["houve_reincidencia"].mean(),',
        '    "aberto":     aberto["houve_reincidencia"].mean(),',
        "})",
    ], tamanho=12, barra=T["tenue"])

    x = MARGEM + Inches(7.72)
    problemas = [
        ("Seis variáveis", "que existem só para chegar num resultado"),
        ("Nomes que não dizem nada", "com_regime vai ser reaproveitado por "
         "engano daqui a três células"),
        ("Não escala", "um quarto regime obriga a escrever mais uma linha, e a "
         "lembrar de incluí-la"),
    ]
    y = Inches(2.45)
    for titulo, texto in problemas:
        caixa(slide, x, y, Inches(0.07), Inches(0.78),
              preenchimento=T["marca"])
        tb = texto_livre(slide, x + Inches(0.24), y - Inches(0.02),
                         Inches(3.30), Inches(0.85))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 14, "bold": True, "depois": 1,
             "cor": T["texto"]},
            {"texto": texto, "tamanho": 11.5, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(0.92)
    return slide


def s05_pipeline(prs, lays):
    slide = novo_slide(prs, lays, "A mesma coisa, encadeada", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "Nenhuma variável intermediária, e a ordem das operações é a ordem "
            "das linhas. Leia de cima para baixo:")

    bloco_codigo(slide, MARGEM, Inches(2.50), Inches(6.90), Inches(2.30), [
        "(",
        "    criminal",
        "    .query(\"classe == 'Apelação Criminal'\")",
        '    .dropna(subset=["regime_inicial"])',
        '    .groupby("regime_inicial")',
        '    .agg(proporcao=("houve_reincidencia", "mean"))',
        ")",
    ], tamanho=14, barra=T["acento1"])

    x = MARGEM + Inches(7.35)
    leitura = [
        "pegue a tabela criminal,",
        "fique só com as apelações,",
        "descarte quem não tem regime,",
        "faça as pilhas por regime,",
        "e calcule a proporção de cada pilha.",
    ]
    tb = texto_livre(slide, x, Inches(2.72), Inches(3.70), Inches(2.0))
    escrever(tb.text_frame, [
        {"texto": l, "tamanho": 14, "cor": T["fraco"], "depois": 7}
        for l in leitura
    ])

    painel(slide, MARGEM, Inches(5.15), FAIXA, Inches(0.98), [
        {"texto": "Por que os parênteses", "tamanho": 13, "bold": True,
         "depois": 2},
        {"texto": "Dentro de um par de parênteses o python deixa você quebrar a "
                  "linha à vontade. Sem eles, criminal seguido de quebra de "
                  "linha e .query(...) é erro de sintaxe. Os parênteses existem "
                  "só para você pôr uma operação por linha.",
         "tamanho": 12.5, "cor": T["superficie_fraco"], "depois": 0},
    ])
    return slide


def s06_verbos(prs, lays):
    slide = novo_slide(prs, lays, "Os cinco verbos", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Cinco operações resolvem quase toda análise descritiva. As cartas "
            "da gincana são estas cinco, mais duas de apoio.")

    # os detalhes cabem em uma linha de propósito: a tabela do PowerPoint
    # cresce para acomodar o texto, e duas linhas em qualquer célula empurram
    # a última operação para fora do slide.
    linhas = [
        ('.query("coluna == \'valor\'")', "escolher linhas",
         "a condição vai escrita como texto", True),
        ('[["a", "b"]]', "escolher colunas",
         "dois pares de colchetes; um só devolve a coluna solta", True),
        ('.sort_values("a", ascending=False)', "ordenar",
         "ascending=False põe o maior primeiro", True),
        ('.groupby("a").agg(nome=("b", "mean"))', "agregar por grupo",
         "uma linha por grupo, e o .reset_index() devolve a coluna", True),
        ('.assign(nova=tabela["b"] > 8)', "criar uma coluna",
         "vem antes de qualquer operação que use a coluna nova", True),
        ('.dropna(subset=["a"])', "apoio: descartar vazios",
         "só as linhas em que aquela coluna está vazia", False),
        (".head(n)", "apoio: as n primeiras",
         "da tabela como ela está naquele ponto", False),
    ]

    tab = tabela(slide, MARGEM, Inches(2.44), FAIXA, Inches(3.30),
                 ["", "operação", "o que faz", "o detalhe que pega"], linhas,
                 [Inches(0.14), Inches(4.05), Inches(2.35), Inches(4.96)],
                 Inches(0.425))
    for r, (op, oque, detalhe, verbo) in enumerate(linhas, start=1):
        z = zebra(r)
        celula(tab, r, 0, "", fundo=T["acento1"] if verbo else T["tenue"])
        celula(tab, r, 1, op, tamanho=12, bold=True, cor=T["texto"], fundo=z)
        celula(tab, r, 2, oque, tamanho=12, cor=T["texto"], fundo=z)
        celula(tab, r, 3, detalhe, tamanho=11.5, cor=T["fraco"], fundo=z)

    painel(slide, MARGEM, Inches(6.05), FAIXA, Inches(0.72), [
        {"texto": "A coluna criada pelo .assign() só existe da linha seguinte "
                  "em diante. Por isso ele vem antes de quem for usar ela.",
         "tamanho": 13.5, "bold": True, "depois": 0},
    ])
    return slide


def s07_ordem(prs, lays):
    slide = novo_slide(prs, lays, "A ordem importa", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "As duas células têm exatamente as mesmas operações. Só trocam duas "
            "linhas de lugar.")

    largura = Inches(5.55)
    x2 = MARGEM + largura + Inches(0.40)

    etiqueta(slide, MARGEM, Inches(2.48), Inches(2.85), Inches(0.36),
             "ordena e depois corta", T["ok"])
    bloco_codigo(slide, MARGEM, Inches(2.96), largura, Inches(1.70), [
        "(",
        "    criminal",
        '    .sort_values("pena_anos", ascending=False)',
        "    .head(5)",
        ")",
    ], tamanho=13, barra=T["ok"])
    legenda(slide, MARGEM, Inches(4.82), largura,
            "As cinco maiores penas da base. É o que a pergunta pedia.")

    etiqueta(slide, x2, Inches(2.48), Inches(2.85), Inches(0.36),
             "corta e depois ordena", T["marca"])
    bloco_codigo(slide, x2, Inches(2.96), largura, Inches(1.70), [
        "(",
        "    criminal",
        "    .head(5)",
        '    .sort_values("pena_anos", ascending=False)',
        ")",
    ], tamanho=13, barra=T["marca"])
    legenda(slide, x2, Inches(4.82), largura,
            "Cinco acórdãos quaisquer, os cinco primeiros da tabela, ordenados "
            "entre si. Nada a ver com as maiores penas.")

    painel(slide, MARGEM, Inches(5.65), FAIXA, Inches(0.72), [
        {"texto": "O .head(5) não sabe o que você queria. Ele pega as cinco "
                  "primeiras linhas da tabela como ela está naquele ponto.",
         "tamanho": 14, "bold": True, "depois": 0},
    ])
    return slide


def s08_definicao(prs, lays):
    slide = novo_slide(prs, lays, "O que é um gráfico", EYEBROW)

    legenda(slide, MARGEM, Inches(1.92), FAIXA,
            "Antes de desenhar qualquer coisa, uma definição. Ela cabe numa "
            "frase, e a frase organiza tudo o que vem depois.")

    painel(slide, MARGEM, Inches(2.40), FAIXA, Inches(1.42), [
        {"texto": "Um gráfico estatístico é um mapeamento de variáveis "
                  "(colunas) em aspectos estéticos de formas geométricas.",
         "tamanho": 24, "bold": True, "depois": 0, "entrelinhas": 1.18},
    ], inverso=True, recuo=0.42)

    legenda(slide, MARGEM, Inches(4.10), FAIXA,
            "Quatro palavras fazem o trabalho:")

    partes = [
        ("mapeamento", "uma ligação: cada valor da coluna vira um valor "
         "visual", T["marca"]),
        ("variáveis (colunas)", "o que sai da tabela, e nada mais",
         T["acento1"]),
        ("aspectos estéticos", "posição, altura, cor, tamanho, forma",
         T["acento3"]),
        ("formas geométricas", "barra, ponto, linha, caixa", T["acento2"]),
    ]
    largura = Inches(2.72)
    x = MARGEM
    for nome, texto, cor in partes:
        caixa(slide, x, Inches(4.58), largura, Inches(0.08), preenchimento=cor)
        tb = texto_livre(slide, x, Inches(4.78), largura, Inches(1.40))
        escrever(tb.text_frame, [
            {"texto": nome, "tamanho": 15, "bold": True, "depois": 4,
             "cor": T["texto"]},
            {"texto": texto, "tamanho": 12.5, "cor": T["fraco"], "depois": 0},
        ])
        x += largura + Inches(0.20)

    legenda(slide, MARGEM, Inches(6.42), FAIXA,
            "Note o que a definição NÃO diz: nada sobre que biblioteca, que "
            "cor bonita ou que tipo de gráfico. Ela diz o que precisa ser "
            "decidido.", cor=T["marca"])
    return slide


def s09_gramatica(prs, lays):
    slide = novo_slide(prs, lays, "A gramática de gráficos", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "A definição vira código quase palavra por palavra. São três "
            "camadas, somadas com +, e todas as três são obrigatórias.")

    caixa(slide, MARGEM, Inches(2.46), FAIXA, Inches(0.86),
          preenchimento=T["inverso"])
    tb = texto_livre(slide, MARGEM + Inches(0.35), Inches(2.66),
                     FAIXA - Inches(0.7), Inches(0.55))
    escrever(tb.text_frame, [
        {"texto": "gráfico  =  dados  +  estética  +  geometria",
         "fonte": MONO, "tamanho": 22, "bold": True,
         "cor": T["inverso_texto"], "depois": 0},
    ])

    elementos = [
        ("dados", "de que tabela saem as variáveis", "ggplot(penas)",
         T["acento1"]),
        ("estética", "que variável vira que aspecto visual",
         'aes(x="regime")', T["acento3"]),
        ("geometria", "que forma aparece na tela", "geom_bar()", T["acento2"]),
    ]

    y = Inches(3.70)
    for nome, oque, exemplo, cor in elementos:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(0.62), preenchimento=cor)
        tb = texto_livre(slide, MARGEM + Inches(0.28), y - Inches(0.02),
                         Inches(2.00), Inches(0.45))
        escrever(tb.text_frame, [
            {"texto": nome, "tamanho": 17, "bold": True, "cor": T["texto"],
             "depois": 0}])
        tb = texto_livre(slide, MARGEM + Inches(2.40), y + Inches(0.04),
                         Inches(5.25), Inches(0.45))
        escrever(tb.text_frame, [
            {"texto": oque, "tamanho": 14, "cor": T["fraco"], "depois": 0}])
        tb = texto_livre(slide, MARGEM + Inches(7.90), y + Inches(0.04),
                         Inches(3.55), Inches(0.45))
        escrever(tb.text_frame, [
            {"texto": exemplo, "fonte": MONO, "tamanho": 14,
             "cor": T["texto"], "depois": 0}])
        y += Inches(0.76)

    painel(slide, MARGEM, Inches(6.10), FAIXA, Inches(0.78), [
        {"texto": "É a mesma ideia dos cinco verbos: lá, operações somadas com "
                  "ponto. Aqui, camadas somadas com mais.",
         "tamanho": 14, "bold": True, "depois": 0},
    ])
    return slide


def s10_camadas(prs, lays):
    slide = novo_slide(prs, lays, "Uma camada de cada vez", EYEBROW)

    passos = [
        ("1. os dados", "ggplot(penas)",
         "Um retângulo vazio. Já é um gráfico válido: só não dissemos que "
         "variável vira o quê.", "cam1", T["acento1"]),
        ("2. a estética", 'ggplot(penas)\n+ aes(x="regime")',
         "O mapeamento aconteceu: regime virou posição no eixo. Ainda não há "
         "forma nenhuma.", "cam2", T["acento3"]),
        ("3. a geometria", 'ggplot(penas)\n+ aes(x="regime")\n+ geom_bar()',
         "A forma apareceu. geom_bar() conta as linhas de cada categoria "
         "sozinho: não existe groupby antes.", "cam3", T["acento2"]),
    ]

    largura = Inches(3.62)
    x = MARGEM
    for titulo, cod, texto, nome, cor in passos:
        etiqueta(slide, x, Inches(1.92), Inches(1.55), Inches(0.34), titulo, cor)
        bloco_codigo(slide, x, Inches(2.38), largura, Inches(0.98),
                     cod.split("\n"), tamanho=11.5, barra=cor)
        imagem_ajustada(slide, fig(nome), x, Inches(3.50), largura,
                        Inches(2.10), moldura=False)
        tb = texto_livre(slide, x, Inches(5.72), largura, Inches(0.95))
        escrever(tb.text_frame, [
            {"texto": texto, "tamanho": 12, "cor": T["fraco"], "depois": 0},
        ])
        x += largura + Inches(0.32)
    return slide


def s11_pintar(prs, lays):
    slide = novo_slide(prs, lays, "Pintar não é mapear", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "A distinção mais importante da aula, e a que mais gera erro. "
            "Depende só de o argumento estar dentro ou fora do aes().")

    largura = Inches(5.55)
    x2 = MARGEM + largura + Inches(0.40)

    etiqueta(slide, MARGEM, Inches(2.42), Inches(3.75), Inches(0.36),
             "fill FORA do aes(): é tinta", T["marca"])
    bloco_codigo(slide, MARGEM, Inches(2.90), largura, Inches(1.28), [
        "ggplot(penas)",
        '+ aes(x="regime")',
        '+ geom_bar(fill="#E50505")',
    ], tamanho=13, barra=T["marca"])
    imagem_ajustada(slide, fig("s11_tinta"), MARGEM, Inches(4.32), largura,
                    Inches(1.95), moldura=False)
    legenda(slide, MARGEM, Inches(6.38), largura,
            "Vale para todas as barras, não representa nada, e não gera "
            "legenda.", tamanho=12)

    etiqueta(slide, x2, Inches(2.42), Inches(4.05), Inches(0.36),
             "fill DENTRO do aes(): é variável", T["acento2"])
    bloco_codigo(slide, x2, Inches(2.90), largura, Inches(1.28), [
        "ggplot(penas)",
        '+ aes(x="regime", fill="houve_reincidencia")',
        '+ geom_bar(position="dodge")',
    ], tamanho=13, barra=T["acento2"])
    imagem_ajustada(slide, fig("s11_variavel"), x2, Inches(4.32), largura,
                    Inches(1.95),
                    moldura=False)
    legenda(slide, x2, Inches(6.38), largura,
            "Cada regime vira duas barras e aparece legenda: agora a cor mapeia "
            "uma variável.", tamanho=12)
    return slide


def s12_regra_aes(prs, lays):
    slide = novo_slide(prs, lays, "A regra vale para toda estética", EYEBROW)

    painel(slide, MARGEM, Inches(2.10), FAIXA, Inches(1.30), [
        {"texto": "Dentro do aes(), o argumento recebe o NOME DE UMA COLUNA.",
         "tamanho": 19, "bold": True, "depois": 6},
        {"texto": "Fora do aes(), o argumento recebe um VALOR FIXO.",
         "tamanho": 19, "bold": True, "depois": 0},
    ], inverso=True, recuo=0.42)

    legenda(slide, MARGEM, Inches(3.66), FAIXA,
            "Trocar os dois de lugar é o erro mais comum de quem começa, e o "
            "sintoma é sempre um destes dois:")

    sintomas = [
        ("Apareceu uma legenda que você não queria",
         'você pôs fill="azul" dentro do aes(), e o plotnine criou uma variável '
         'com um valor só, chamada "azul"'),
        ("Sumiu a legenda que você queria",
         'você pôs fill="regime" fora do aes(), e o plotnine tentou pintar tudo '
         'de uma cor chamada "regime", que não existe'),
    ]
    y = Inches(4.20)
    for titulo, texto in sintomas:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(0.85),
              preenchimento=T["marca"])
        tb = texto_livre(slide, MARGEM + Inches(0.28), y - Inches(0.02),
                         FAIXA - Inches(0.4), Inches(0.9))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 16, "bold": True, "depois": 2,
             "cor": T["texto"]},
            {"texto": texto, "tamanho": 13, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(1.05)

    legenda(slide, MARGEM, Inches(6.40), FAIXA,
            "Os aspectos estéticos mais usados: x, y, fill (preenchimento), "
            "color (traço), size (tamanho), alpha (transparência).")
    return slide


def s13_geometria(prs, lays):
    slide = novo_slide(prs, lays, "Geometrias e tipos de variáveis", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "Hoje só variáveis sozinhas. Quinta-feira, duas de cada vez.")

    linhas = [
        ("uma categórica", "quantos casos em cada categoria", "geom_bar()",
         "conta sozinho: a altura sai da contagem", T["acento1"]),
        ("uma categórica, altura já calculada", "mostrar um valor por categoria",
         "geom_col()", "usa a coluna que você mapeou em y", T["acento1"]),
        ("uma numérica", "como os valores se distribuem",
         "geom_histogram(bins=)", "o número de caixas muda a leitura",
         T["acento3"]),
        ("uma numérica", "a mesma distribuição, alisada", "geom_density()",
         "sem o degrau das caixas", T["acento3"]),
        ("uma numérica", "mediana, quartis e pontos fora", "geom_boxplot()",
         "resume, e por isso esconde a forma", T["acento3"]),
    ]

    tab = tabela(slide, MARGEM, Inches(2.50), FAIXA, Inches(3.35),
                 ["", "a variável", "a pergunta", "a geometria", "o detalhe"],
                 linhas, [Inches(0.14), Inches(3.05), Inches(3.15),
                          Inches(2.40), Inches(2.76)], Inches(0.60))
    for r, (var, perg, geom, det, cor) in enumerate(linhas, start=1):
        z = zebra(r)
        celula(tab, r, 0, "", fundo=cor)
        celula(tab, r, 1, var, tamanho=12, bold=True, cor=T["texto"], fundo=z)
        celula(tab, r, 2, perg, tamanho=12, cor=T["texto"], fundo=z)
        celula(tab, r, 3, geom, tamanho=12, bold=True, cor=T["texto"], fundo=z)
        celula(tab, r, 4, det, tamanho=11.5, cor=T["fraco"], fundo=z)

    painel(slide, MARGEM, Inches(6.05), FAIXA, Inches(0.72), [
        {"texto": "geom_bar() conta. geom_col() usa a altura que você deu. "
                  "Confundir os dois é a pegadinha da última rodada.",
         "tamanho": 14, "bold": True, "depois": 0},
    ])
    return slide


# ---------------------------------------------------------------- a gincana

def s14_regras(prs, lays):
    slide = novo_slide(prs, lays, "Gincana do Pipeline", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Cada mesa tem um baralho de 42 cartas, um tabuleiro e uma folha de "
            "resposta. Na frente da carta está o código; no verso, o que ele "
            "faz. Virar a carta é permitido, e é o ponto da atividade.")

    passos = [
        ("1", "Três aquecimentos", "de duas ou três cartas, que não pontuam. "
         "São para pegar a mecânica."),
        ("2", "Eu projeto o problema", "e digo quantas cartas a resposta tem."),
        ("3", "A mesa monta no tabuleiro",
         "uma carta por linha, na ordem em que o python executa."),
        ("4", "Levanta a placa", "e escreve a ordem na folha de resposta."),
        ("5", "Escreve a linha do descarte",
         "uma carta que ficou de fora, e por que ela não serve."),
    ]
    y = Inches(2.72)
    for num, titulo, texto in passos:
        etiqueta(slide, MARGEM, y, Inches(0.42), Inches(0.42), num, T["marca"])
        tb = texto_livre(slide, MARGEM + Inches(0.68), y - Inches(0.02),
                         Inches(5.60), Inches(0.8))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 14.5, "bold": True, "depois": 1,
             "cor": T["texto"]},
            {"texto": texto, "tamanho": 12, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(0.78)

    x = MARGEM + Inches(6.70)
    painel(slide, x, Inches(2.60), Inches(4.80), Inches(2.45), [
        {"texto": "Como se pontua", "tamanho": 15, "bold": True, "depois": 8},
        {"texto": "3 pontos  ·  ordem certa, dentro do tempo. Vale igual para "
                  "o primeiro e para o último que acertar.",
         "tamanho": 12.5, "cor": T["superficie_fraco"], "depois": 6},
        {"texto": "+1 ponto  ·  a linha do descarte, se ela disser por quê.",
         "tamanho": 12.5, "cor": T["superficie_fraco"], "depois": 6},
        {"texto": "+1 ponto  ·  para os três primeiros a levantar a placa com a "
                  "ordem certa.",
         "tamanho": 12.5, "cor": T["superficie_fraco"], "depois": 0},
    ])

    painel(slide, x, Inches(5.22), Inches(4.80), Inches(1.20), [
        {"texto": "Dois prêmios", "tamanho": 14, "bold": True, "depois": 4},
        {"texto": "A lane com mais pontos somados, e um grupo sorteado entre "
                  "todos os que pontuaram em quatro rodadas ou mais.",
         "tamanho": 12, "depois": 0},
    ], inverso=True)

    legenda(slide, MARGEM, Inches(6.42), Inches(6.30),
            "O relator muda a cada rodada: é quem explica se a mesa for "
            "chamada.", cor=T["marca"])
    return slide


def slide_pergunta(prs, lays, r):
    slide = novo_slide(prs, lays, f"{rotulo(r)}: {r['titulo']}", EYEBROW)

    aquece = eh_aquecimento(r)
    etiqueta(slide, MARGEM, Inches(1.92), Inches(1.62), Inches(0.38),
             f"{n_cartas(r)} cartas", T["marca"])
    # no aquecimento o minuto não importa e o que importa é não valer ponto
    etiqueta(slide, MARGEM + Inches(1.82), Inches(1.92), Inches(1.85),
             Inches(0.38),
             "não pontua" if aquece else f"{r['minutos']} min", T["fraco"])

    tem_figura = bool(r.get("figura"))
    largura_texto = Inches(5.30) if tem_figura else FAIXA

    tb = texto_livre(slide, MARGEM, Inches(2.60), largura_texto, Inches(2.2))
    escrever(tb.text_frame, [
        {"texto": r["pergunta"], "tamanho": 18 if tem_figura else 21,
         "depois": 0, "entrelinhas": 1.22, "cor": T["texto"]},
    ])

    if tem_figura:
        imagem_ajustada(slide, fig(r["figura"]), MARGEM + Inches(5.75),
                        Inches(2.20), Inches(5.75), Inches(4.10),
                        moldura=False)

    naipe = {"pandas": "cartas turquesa", "gráfico": "cartas roxas",
             "os dois": "as duas cores"}[r["naipe"]]
    legenda(slide, MARGEM, Inches(6.45), largura_texto,
            f"Procurem só nas {naipe}.")
    return slide


def slide_resposta(prs, lays, r):
    slide = novo_slide(prs, lays, f"{rotulo(r)}: a resposta", EYEBROW)

    seq = "  >  ".join(r["gabarito"])
    if "gabarito2" in r:
        seq += "        //        " + "  >  ".join(r["gabarito2"])

    caixa(slide, MARGEM, Inches(1.92), FAIXA, Inches(0.60),
          preenchimento=T["marca"])
    tb = texto_livre(slide, MARGEM + Inches(0.30), Inches(2.05),
                     FAIXA - Inches(0.6), Inches(0.4))
    escrever(tb.text_frame, [
        {"texto": seq, "fonte": MONO, "tamanho": 16, "bold": True,
         "cor": T["etiqueta_texto"], "depois": 0},
    ])

    linhas = codigo_rodada(r).split("\n")
    altura = Inches(0.26) * len(linhas) + Inches(0.30)
    barra = T["acento1"] if r["naipe"] == "pandas" else T["acento2"]
    bloco_codigo(slide, MARGEM, Inches(2.72), Inches(6.75), altura, linhas,
                 tamanho=12.5, barra=barra)

    x = MARGEM + Inches(7.20)
    if eh_aquecimento(r):
        painel(slide, x, Inches(2.72), Inches(4.30), Inches(3.40), [
            {"texto": "O QUE O AQUECIMENTO ENSINA", "tamanho": 10,
             "bold": True, "cor": T["marca"], "depois": 5},
            {"texto": r["ensina"], "tamanho": 11.5, "depois": 0},
        ], recuo=0.26)
        return slide

    painel(slide, x, Inches(2.72), Inches(4.30), Inches(1.95), [
        {"texto": "O DISTRATOR QUE IMPORTA", "tamanho": 10, "bold": True,
         "cor": T["marca"], "depois": 5},
        {"texto": r["distrator"], "tamanho": 12.5, "bold": True, "depois": 3},
        {"texto": r["porque"], "tamanho": 11.5,
         "cor": T["superficie_fraco"], "depois": 0},
    ], recuo=0.26)

    painel(slide, x, Inches(4.85), Inches(4.30), Inches(1.55), [
        {"texto": "O QUE ESTA RODADA ENSINA", "tamanho": 10, "bold": True,
         "depois": 5},
        {"texto": r["ensina"], "tamanho": 11.5, "depois": 0},
    ], inverso=True, recuo=0.26)
    return slide


# ----------------------------------------------------------------- fechamento

def s_fechamento(prs, lays):
    slide = novo_slide(prs, lays, "Quinta-feira", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "Aula 6, dia 27/08, das 16h30 às 18h30.")

    blocos = [
        ("Primeira hora", "Gráficos de duas variáveis",
         "Duas numéricas, uma numérica e uma categórica, duas categóricas. "
         "É a mesma gramática de hoje, com mais uma variável mapeada.",
         T["acento2"]),
        ("Segunda hora", "Projeto 02, individual, valendo nota",
         "Você recebe blocos de pandas e de plotnine fora de ordem, e organiza "
         "na ordem em que devem rodar. É exatamente a gincana de hoje, sozinho "
         "e no computador. Entrega no mesmo dia.", T["marca"]),
    ]
    y = Inches(2.60)
    for tempo, titulo, texto, cor in blocos:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(1.55), preenchimento=cor)
        tb = texto_livre(slide, MARGEM + Inches(0.32), y - Inches(0.02),
                         FAIXA - Inches(0.5), Inches(1.6))
        escrever(tb.text_frame, [
            {"texto": tempo, "tamanho": 11, "bold": True, "cor": cor,
             "depois": 3},
            {"texto": titulo, "tamanho": 19, "bold": True, "depois": 5,
             "cor": T["texto"]},
            {"texto": texto, "tamanho": 14, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(1.85)

    painel(slide, MARGEM, Inches(6.20), FAIXA, Inches(0.75), [
        {"texto": "Para chegar pronto: rode o notebook da aula 5 inteiro, "
                  "principalmente a parte que você não digitou em sala.",
         "tamanho": 14, "bold": True, "depois": 0},
    ])
    return slide


# ---------------------------------------------------------------------- main

def montar(escuro: bool, respostas: bool = True) -> str:
    """Monta um deck. `respostas=False` deixa de fora os gabaritos das rodadas.

    A versão sem gabarito é a que vai para os alunos antes da aula. Os
    aquecimentos continuam resolvidos: eles não pontuam, e existem justamente
    para a turma chegar sabendo como a mecânica funciona.
    """
    global T
    T = ESCURO if escuro else CLARO

    prs = deck_limpo()
    lays = layouts(prs)

    for construir in (s01_capa, s02_plano, s03_pergunta, s04_jeito_antigo,
                      s05_pipeline, s06_verbos, s07_ordem, s08_definicao,
                      s09_gramatica, s10_camadas, s11_pintar, s12_regra_aes,
                      s13_geometria, s14_regras):
        construir(prs, lays)

    for r in TUDO:
        slide_pergunta(prs, lays, r)
        if respostas or eh_aquecimento(r):
            slide_resposta(prs, lays, r)

    s_fechamento(prs, lays)

    if not respostas:
        nome = "aula05_alunos.pptx"
    else:
        nome = "aula05_dark.pptx" if escuro else "aula05.pptx"
    caminho = os.path.join(AQUI, nome)
    n = gravar(prs, caminho,
               titulo="Aula 5: pipeline e gramática de gráficos")
    print(f"{nome}  ({n} slides)")
    return caminho


def main():
    so_claro = "--claro" in sys.argv
    so_escuro = "--escuro" in sys.argv
    so_alunos = "--alunos" in sys.argv
    if so_alunos:
        montar(escuro=False, respostas=False)
        return
    if not so_escuro:
        montar(escuro=False)
    if not so_claro:
        montar(escuro=True)
    # a versão que os alunos recebem antes da aula, sem os gabaritos
    montar(escuro=False, respostas=False)


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    main()

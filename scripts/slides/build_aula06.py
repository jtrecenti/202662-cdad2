"""Monta os slides da Aula 6: graficos de duas variaveis.

O deck e curto de proposito: sao 18 minutos de exposicao, porque a segunda hora
da aula inteira e o Projeto 02. Os slides existem para instalar uma frase, "o
par de tipos escolhe a geometria", e para mostrar as tres formas do `position`
lado a lado, que e a unica ideia da aula que nao cabe numa tabela.

As figuras vem de `graficos_aula06.py`, e saem da mesma tabela do notebook.

Uso:
    python graficos_aula06.py && python build_aula06.py
"""

from __future__ import annotations

import os
import sys

from pptx.dml.color import RGBColor
from pptx.util import Inches

from insper import (
    BRANCO, CINZA, CINZA_CLARO, CINZA_ESCURO, FAIXA, LARANJA, MARGEM, PRETO,
    QUASE_BRANCO, ROXO, TURQUESA, VERDE, VERMELHO, caixa, celula, deck_limpo,
    escrever, gravar, imagem_ajustada, layouts, slide_capa, slide_com_titulo,
    tabela_sem_estilo, texto_livre,
)

AQUI = os.path.dirname(os.path.abspath(__file__))
FIGURAS = os.path.join(AQUI, "assets", "aula06")
SAIDA = os.path.join(AQUI, "aula06.pptx")

EYEBROW = ("Ciência de Dados Aplicada ao Direito II  ·  Aula 6  ·  "
           "Gráficos de duas variáveis")

MONO = "Consolas"


# --------------------------------------------------------------------- tema
#
# O deck sai em duas versoes a partir do mesmo codigo. A escura existe porque
# em sala com projetor fraco o fundo claro estoura; o que muda e so o
# dicionario abaixo, trocado antes de montar.

GRAFITE = RGBColor(0x14, 0x16, 0x1A)
GRAFITE_CLARO = RGBColor(0x22, 0x25, 0x2B)
CINZA_TEXTO = RGBColor(0xC9, 0xCC, 0xD2)
VERMELHO_CLARO = RGBColor(0xFF, 0x6B, 0x6B)
TURQUESA_CLARO = RGBColor(0x5A, 0xD8, 0xA6)
ROXO_CLARO = RGBColor(0xB4, 0x7D, 0xE0)
LARANJA_CLARO = RGBColor(0xFF, 0xB2, 0x6B)

CLARO = {
    "escuro": False,
    "fundo": None,
    "titulo": PRETO,
    "texto": PRETO,
    "fraco": CINZA_ESCURO,
    "neutro": CINZA_ESCURO,
    "superficie": QUASE_BRANCO,
    "inverso": PRETO,
    "inverso_texto": BRANCO,
    "etiqueta_texto": BRANCO,
    "marca": VERMELHO,
    "acento1": TURQUESA,
    "acento2": ROXO,
    "acento3": LARANJA,
    "figura": "",
}

ESCURO = {
    "escuro": True,
    "fundo": GRAFITE,
    "titulo": BRANCO,
    "texto": BRANCO,
    "fraco": CINZA_TEXTO,
    "neutro": CINZA_TEXTO,
    "superficie": GRAFITE_CLARO,
    "inverso": BRANCO,
    "inverso_texto": GRAFITE,
    "etiqueta_texto": GRAFITE,
    "marca": VERMELHO_CLARO,
    "acento1": TURQUESA_CLARO,
    "acento2": ROXO_CLARO,
    "acento3": LARANJA_CLARO,
    "figura": "_dark",
}

T = CLARO

LOGO_BRANCO = os.path.join(AQUI, "assets", "insper_branco.png")
# a posicao exata do logotipo no layout, lida do proprio arquivo: chutar
# deixa as duas versoes desalinhadas e as duas aparecem
LOGO = (Inches(11.5059), Inches(6.8343), Inches(1.2828), Inches(0.5814))


def fig(nome: str) -> str:
    """O caminho da figura, na variante do tema."""
    return os.path.join(FIGURAS, f"{nome}{T['figura']}.png")


def novo_slide(prs, lays, titulo, eyebrow):
    """Slide com titulo, ja com o fundo e as cores do tema aplicados."""
    slide = slide_com_titulo(prs, lays, titulo, eyebrow)
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if not ph.text_frame.paragraphs[0].runs:
            continue
        run = ph.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = T["titulo"] if idx == 0 else T["fraco"]
    if T["escuro"]:
        pintar_fundo(slide)
    return slide


def pintar_fundo(slide):
    """Fundo escuro atras de tudo, mais o logotipo em negativo.

    O fundo do slide fica ATRAS das formas do layout, e o logotipo do layout e
    preto: no escuro ele sumiria. Por isso ele e recoberto por uma versao em
    branco, na mesma posicao.
    """
    fundo = slide.background
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = T["fundo"]
    x, y, cx, cy = LOGO
    slide.shapes.add_picture(LOGO_BRANCO, x, y, cx, cy)


# ------------------------------------------------------------------ utilidades

def bloco_codigo(slide, x, y, cx, cy, linhas, *, tamanho=13, barra=None):
    caixa(slide, x, y, cx, cy, preenchimento=T["superficie"])
    if barra is not None:
        caixa(slide, x, y, Inches(0.07), cy, preenchimento=barra)
    tb = texto_livre(slide, x + Inches(0.26), y + Inches(0.13),
                     cx - Inches(0.42), cy - Inches(0.22))
    escrever(tb.text_frame, linhas, fonte=MONO, tamanho=tamanho, cor=T["texto"],
             entrelinhas=1.16, espaco_depois=0)
    return tb


def legenda(slide, x, y, cx, texto, *, cor=T["fraco"], tamanho=13):
    tb = texto_livre(slide, x, y, cx, Inches(0.3))
    escrever(tb.text_frame, [texto], tamanho=tamanho, cor=cor, espaco_depois=0)
    return tb


def painel(slide, x, y, cx, cy, linhas, *, inverso=False, recuo=0.30):
    """Caixa de destaque, clara ou invertida, com paragrafos ja formatados."""
    fundo = PRETO if inverso else QUASE_BRANCO
    base = BRANCO if inverso else PRETO
    caixa(slide, x, y, cx, cy, preenchimento=fundo)
    tb = texto_livre(slide, x + Inches(recuo), y + Inches(0.16),
                     cx - Inches(recuo * 2), cy - Inches(0.24))
    escrever(tb.text_frame, [{**item, "cor": item.get("cor", base)}
                             for item in linhas])
    return tb


def tabela(slide, x, y, cx, cy, cabecalho, linhas, larguras, alt_linha):
    """Tabela no estilo do deck: cabecalho preto e altura fixa por linha."""
    shape = slide.shapes.add_table(len(linhas) + 1, len(cabecalho), x, y, cx, cy)
    tab = shape.table
    tabela_sem_estilo(tab)
    for larg, i in zip(larguras, range(len(cabecalho))):
        tab.columns[i].width = larg
    tab.rows[0].height = Inches(0.32)
    for r in range(1, len(linhas) + 1):
        tab.rows[r].height = alt_linha
    for c, texto in enumerate(cabecalho):
        celula(tab, 0, c, texto, tamanho=11, bold=True, cor=T["etiqueta_texto"], fundo=T["inverso"])
    return tab


def etiqueta(slide, x, y, cx, cy, texto, cor, *, tamanho=11.5):
    sh = caixa(slide, x, y, cx, cy, preenchimento=cor)
    escrever(sh.text_frame, [texto], tamanho=tamanho, cor=T["etiqueta_texto"], bold=True,
             espaco_depois=0)
    return sh


# ---------------------------------------------------------------------- slides

def s01_capa(prs, lays):
    return slide_capa(
        prs, lays,
        subtitulo="Aula 6: gráficos de duas variáveis",
        subtema="O par de tipos escolhe a geometria",
    )


def s02_plano(prs, lays):
    slide = novo_slide(prs, lays, "O plano de hoje", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "Aula curta de propósito: a segunda hora inteira é o projeto "
            "aplicado, em grupo e valendo nota.")

    itens = [
        ("25 min", "Slides", "o que ficou da gincana, e de uma variável "
         "para duas.", T["marca"]),
        ("25 min", "Notebook", "as geometrias novas e as três operações de "
         "pandas que o projeto usa.", T["acento2"]),
        ("60 min", "Projeto 02", "no app da disciplina. Você organiza blocos de "
         "pandas e de plotnine na ordem em que devem rodar.", T["acento1"]),
        ("5 min", "Fechamento", "entrega e o que vem na aula 7.", CINZA_ESCURO),
    ]
    y = Inches(2.62)
    for tempo, titulo, texto, cor in itens:
        etiqueta(slide, MARGEM, y, Inches(1.05), Inches(0.42), tempo, cor)
        tb = texto_livre(slide, MARGEM + Inches(1.30), y - Inches(0.02),
                         FAIXA - Inches(1.30), Inches(0.9))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 17, "bold": True, "depois": 1},
            {"texto": texto, "tamanho": 13, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(1.02)
    return slide


def s02b_retomada_pipeline(prs, lays):
    """Retomada da aula 5: a sequencia de operacoes."""
    slide = novo_slide(prs, lays, "De onde viemos: o encadeamento",
                             EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Uma análise é uma sequência de operações somadas com ponto. Cada "
            "uma recebe a tabela que a anterior devolveu.")

    linhas = [
        ('.query("coluna == \'valor\'")', "escolher linhas"),
        ('[["a", "b"]]', "escolher colunas: dois pares de colchetes"),
        ('.sort_values("a", ascending=False)', "ordenar"),
        ('.groupby("a").agg(nome=("b", "mean"))', "uma linha por grupo"),
        ('.assign(nova=tabela["b"] > 8)', "criar uma coluna"),
    ]
    y = Inches(2.50)
    for op, oque in linhas:
        caixa(slide, MARGEM, y, Inches(0.08), Inches(0.40),
              preenchimento=T["acento1"])
        tb = texto_livre(slide, MARGEM + Inches(0.24), y - Inches(0.05),
                         Inches(4.10), Inches(0.46))
        escrever(tb.text_frame, [{"texto": op, "fonte": MONO, "tamanho": 12,
                                  "bold": True, "depois": 0}])
        tb = texto_livre(slide, MARGEM + Inches(4.40), y - Inches(0.05),
                         Inches(2.40), Inches(0.46))
        escrever(tb.text_frame, [{"texto": oque, "tamanho": 12,
                                  "cor": T["fraco"], "depois": 0}])
        y += Inches(0.52)

    etiqueta(slide, MARGEM + Inches(6.85), Inches(2.42), Inches(3.20),
             Inches(0.40), "e sempre entre parênteses", CINZA_ESCURO)
    bloco_codigo(slide, MARGEM + Inches(6.85), Inches(2.94), Inches(4.65),
                 Inches(2.10), [
        "(",
        "    criminal",
        '    .query("eh_trafico")',
        '    .sort_values("pena_anos", ascending=False)',
        '    [["processo", "comarca", "pena_anos"]]',
        "    .head(5)",
        ")",
    ], tamanho=11, barra=T["acento1"])
    legenda(slide, MARGEM + Inches(6.85), Inches(5.20), Inches(4.65),
            "Os parênteses são o que permitem quebrar a linha e ler o "
            "encadeamento de cima para baixo, uma operação por linha.",
            tamanho=12)

    painel(slide, MARGEM, Inches(5.35), Inches(6.30), Inches(1.00), [
        {"texto": "A ordem é o conteúdo.", "tamanho": 15, "bold": True,
         "depois": 3},
        {"texto": "Filtrar antes ou depois de agrupar dá números diferentes, e "
                  "uma operação só usa coluna que já existe naquele ponto.",
         "tamanho": 12.5, "cor": T["fraco"], "depois": 0},
    ])
    return slide


def s02c_retomada_grafico(prs, lays):
    """Retomada da aula 5: a gramatica de graficos."""
    slide = novo_slide(prs, lays, "De onde viemos: a gramática",
                             EYEBROW)

    painel(slide, MARGEM, Inches(1.94), FAIXA, Inches(0.80), [
        {"texto": "Um gráfico estatístico é um mapeamento de variáveis "
                  "(colunas) em aspectos estéticos de formas geométricas.",
         "tamanho": 17, "bold": True, "depois": 0},
    ], inverso=True, recuo=0.42)

    camadas = [
        ("1. os dados", "de qual tabela o gráfico sai", "ggplot(penas)",
         T["acento1"]),
        ("2. a estética", "qual coluna vira qual aspecto do desenho",
         '+ aes(x="regime")', T["acento3"]),
        ("3. a geometria", "que forma ocupa essas posições", "+ geom_bar()",
         T["acento2"]),
    ]
    y = Inches(3.05)
    for titulo, oque, codigo, cor in camadas:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(0.78), preenchimento=cor)
        tb = texto_livre(slide, MARGEM + Inches(0.28), y - Inches(0.04),
                         Inches(3.30), Inches(0.84))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 14, "bold": True, "depois": 2},
            {"texto": oque, "tamanho": 12, "cor": T["fraco"], "depois": 0},
        ])
        bloco_codigo(slide, MARGEM + Inches(3.80), y - Inches(0.02),
                     Inches(3.60), Inches(0.62), [codigo], tamanho=12,
                     barra=cor)
        y += Inches(0.95)

    imagem_ajustada(slide, fig("retomada_bar"),
                    MARGEM + Inches(7.70), Inches(3.00), Inches(3.80),
                    Inches(2.40), moldura=False)

    painel(slide, MARGEM, Inches(6.00), FAIXA, Inches(0.85), [
        {"texto": "Dentro do aes(), nome de coluna. Fora do aes(), valor fixo.",
         "tamanho": 15, "bold": True, "depois": 3},
        {"texto": "É a distinção que mais gera erro, e ela vale para x, y, "
                  "fill, color, size e alpha.",
         "tamanho": 12.5, "cor": T["fraco"], "depois": 0},
    ])
    return slide


def s09b_catalogo_uni(prs, lays):
    """Um exemplo desenhado para cada grafico de uma variavel."""
    slide = novo_slide(prs, lays, "Uma variável: os três casos", EYEBROW)

    legenda(slide, MARGEM, Inches(1.90), FAIXA,
            "O tipo da variável escolhe a geometria. Estes são todos os casos "
            "de uma variável sozinha.")

    casos = [
        ("categórica · frequência", T["acento1"], "cat_freq",
         ['ggplot(penas, aes(x="regime"))', "+ geom_bar()"],
         "geom_bar() conta as linhas sozinho: não existe coluna y."),
        ("categórica · frequência relativa", T["acento1"], "cat_prop",
         ['ggplot(freq, aes(x="regime", y="prop"))', "+ geom_col()"],
         "A divisão é feita antes, no pandas. A altura já vem pronta na "
         "tabela, então o geom é geom_col()."),
        ("numérica · distribuição", T["acento2"], "num_hist",
         ['ggplot(penas, aes(x="pena_anos"))', "+ geom_histogram(bins=20)"],
         "Numérica não tem categoria para contar: o histograma corta em faixas "
         "e conta cada faixa."),
    ]

    largura = Inches(3.70)
    passo = largura + Inches(0.20)
    for i, (rotulo, cor, figura, codigo, nota) in enumerate(casos):
        x = MARGEM + i * passo
        etiqueta(slide, x, Inches(2.42), largura, Inches(0.38), rotulo, cor)
        bloco_codigo(slide, x, Inches(2.92), largura, Inches(0.78), codigo,
                     tamanho=10, barra=cor)
        imagem_ajustada(slide, fig(figura),
                        x, Inches(3.82), largura, Inches(2.20), moldura=False)
        tb = texto_livre(slide, x, Inches(6.10), largura, Inches(0.90))
        escrever(tb.text_frame, [{"texto": nota, "tamanho": 11,
                                  "cor": T["fraco"], "depois": 0}])
    return slide


def s09c_catalogo_bi(prs, lays):
    """Um exemplo desenhado para cada grafico de duas variaveis."""
    slide = novo_slide(prs, lays, "Duas variáveis: os quatro casos",
                             EYEBROW)

    legenda(slide, MARGEM, Inches(1.90), FAIXA,
            "Agora o par de tipos escolhe a geometria. E quando o eixo x é "
            "tempo, o par muda de nome: vira série.")

    casos = [
        ("categórica × categórica", T["acento2"], "fill",
         ['aes(x="regime", fill="houve_reincidencia")',
          '+ geom_bar(position="fill")'],
         "Barras repartidas. Em proporção com fill, em contagem com dodge."),
        ("categórica × numérica", T["acento3"], "boxplot",
         ['aes(x="regime", y="pena_anos")', "+ geom_boxplot() + coord_flip()"],
         "Uma caixa por categoria, comparando a distribuição inteira."),
        ("numérica × numérica", T["acento1"], "pontos_smooth",
         ['aes(x="n_palavras_ementa", y="pena_anos")',
          "+ geom_point() + geom_smooth(method=\"lm\")"],
         "Um ponto por linha, e a reta resume a direção da nuvem."),
        ("numérica × tempo", T["marca"], "serie",
         ['aes(x="mes", y="n")', "+ geom_line() + geom_point()"],
         "Com o tempo no eixo x, a linha liga os pontos e mostra a trajetória."),
    ]

    largura = Inches(2.75)
    passo = largura + Inches(0.18)
    for i, (rotulo, cor, figura, codigo, nota) in enumerate(casos):
        x = MARGEM + i * passo
        etiqueta(slide, x, Inches(2.42), largura, Inches(0.38), rotulo, cor,
                 tamanho=10.5)
        bloco_codigo(slide, x, Inches(2.92), largura, Inches(0.86), codigo,
                     tamanho=8.5, barra=cor)
        imagem_ajustada(slide, fig(figura),
                        x, Inches(3.90), largura, Inches(2.10), moldura=False)
        tb = texto_livre(slide, x, Inches(6.08), largura, Inches(0.95))
        escrever(tb.text_frame, [{"texto": nota, "tamanho": 10.5,
                                  "cor": T["fraco"], "depois": 0}])
    return slide


def s09b_univariados(prs, lays):
    """Retomada: uma variavel sozinha, por tipo."""
    slide = novo_slide(prs, lays, "Uma variável: qual gráfico", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "Os mesmos casos, agora em tabela, para consultar depois. "
            "Entram também duas geometrias que só aparecem no notebook.")

    linhas = [
        ("categórica", "quantos casos em cada categoria", "geom_bar()",
         "conta sozinho: não existe coluna y", T["acento1"]),
        ("categórica, altura já calculada", "um valor por categoria",
         "geom_col()", "usa a coluna que você mapeou em y", T["acento1"]),
        ("numérica", "como os valores se distribuem", "geom_histogram(bins=)",
         "o número de faixas muda a leitura", T["acento2"]),
        ("numérica", "mediana, quartis e pontos fora", "geom_boxplot()",
         "resume, e por isso esconde a forma", T["acento2"]),
    ]

    tab = tabela(slide, MARGEM, Inches(2.55), FAIXA, Inches(3.00),
                 ["", "a variável", "a pergunta", "a geometria", "o detalhe"],
                 linhas, [Inches(0.14), Inches(3.05), Inches(3.15),
                          Inches(2.40), Inches(2.76)], Inches(0.66))
    for r, (var, perg, geom, det, cor) in enumerate(linhas, start=1):
        z = T["superficie"] if r % 2 else None
        celula(tab, r, 0, "", fundo=cor)
        celula(tab, r, 1, var, tamanho=12, bold=True, cor=T["texto"], fundo=z)
        celula(tab, r, 2, perg, tamanho=12, cor=T["texto"], fundo=z)
        celula(tab, r, 3, geom, tamanho=12, bold=True, cor=T["texto"], fundo=z)
        celula(tab, r, 4, det, tamanho=11.5, cor=T["fraco"], fundo=z)

    painel(slide, MARGEM, Inches(6.05), FAIXA, Inches(0.85), [
        {"texto": "A pergunta a fazer antes de escolher: a altura é uma "
                  "contagem, ou uma conta que eu já fiz?", "tamanho": 15,
         "bold": True, "depois": 3},
        {"texto": "Contagem é geom_bar(). Conta pronta é geom_col(). Todo o "
                  "resto decorre do tipo da variável.",
         "tamanho": 12.5, "cor": T["fraco"], "depois": 0},
    ])
    return slide


def s09c_datatoviz(prs, lays):
    """O data-to-viz, para quando a duvida aparecer fora da aula."""
    slide = novo_slide(prs, lays, "From data to viz", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "As tabelas destes slides cobrem o que a disciplina usa. Para o "
            "resto, existe uma árvore de decisão pronta: você diz que tipo de "
            "dado tem, e ela leva ao gráfico.")

    imagem_ajustada(slide, fig("data_to_viz"),
                    MARGEM, Inches(2.50), Inches(8.10), Inches(3.95),
                    moldura=True)

    x2 = MARGEM + Inches(8.45)
    etiqueta(slide, x2, Inches(2.50), Inches(3.05), Inches(0.40),
             "data-to-viz.com", T["marca"])
    tb = texto_livre(slide, x2, Inches(3.10), Inches(3.05), Inches(3.20))
    escrever(tb.text_frame, [
        {"texto": "Escolha o tipo de dado no alto (numérico, categórico, os "
                  "dois, mapa, rede, série temporal) e siga a árvore.",
         "tamanho": 12.5, "depois": 8},
        {"texto": "Cada gráfico do fim leva a uma página com o código, e a uma "
                  "lista do que costuma dar errado naquele tipo.",
         "tamanho": 12.5, "depois": 8},
        {"texto": "A seção CAVEATS é a melhor parte: é um catálogo de gráficos "
                  "que enganam, com o motivo.",
         "tamanho": 12.5, "cor": T["fraco"], "depois": 0},
    ])
    return slide


def s03_assign(prs, lays):
    """A rodada 3 da gincana, que a aula 5 nao alcancou."""
    slide = novo_slide(prs, lays, "A coluna que não estava na tabela",
                             EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Ficou faltando uma rodada da gincana, e ela é a mais usada hoje: "
            "a coluna que responde à pergunta nem sempre existe na base.")

    etiqueta(slide, MARGEM, Inches(2.42), Inches(4.10), Inches(0.40),
             "a coluna nasce no meio do caminho", T["acento1"])
    bloco_codigo(slide, MARGEM, Inches(2.90), Inches(6.30), Inches(1.95), [
        "(",
        "    criminal",
        '    .assign(eh_capital=criminal["comarca"] == "São Paulo")',
        '    .groupby("eh_capital")',
        '    .agg(n=("processo", "size"))',
        "    .reset_index()",
        ")",
    ], tamanho=12, barra=T["acento1"])

    itens = [
        ("A base tem comarca, e não tem eh_capital.",
         "Nenhuma coluna responde “capital ou interior”."),
        (".assign() cria a coluna sem quebrar o encadeamento.",
         "Por isso ele cabe no meio do pipeline."),
        ("Ele vem ANTES do .groupby() que usa a coluna.",
         "A coluna precisa existir para ser agrupada. É a única ordem que não "
         "pode trocar."),
    ]
    y = Inches(4.98)
    for titulo, texto in itens:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(0.72), preenchimento=T["acento1"])
        tb = texto_livre(slide, MARGEM + Inches(0.28), y - Inches(0.06),
                         Inches(6.20), Inches(0.82))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 12.5, "bold": True, "depois": 2},
            {"texto": texto, "tamanho": 11, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(0.84)

    # a tabela que ESTE pipeline produz, e não uma figura qualquer: o slide
    # existe para mostrar o que o código da esquerda devolve
    etiqueta(slide, MARGEM + Inches(6.85), Inches(2.42), Inches(2.60),
             Inches(0.40), "o que isso devolve", T["acento1"])
    linhas = [("False", "378"), ("True", "97")]
    tab = tabela(slide, MARGEM + Inches(6.85), Inches(2.96), Inches(4.00),
                 Inches(1.24), ["eh_capital", "n"], linhas,
                 [Inches(2.55), Inches(1.45)], Inches(0.44))
    for r, (v, n) in enumerate(linhas, start=1):
        z = T["superficie"] if r % 2 else None
        celula(tab, r, 0, v, tamanho=13, bold=True, cor=T["texto"], fundo=z)
        celula(tab, r, 1, n, tamanho=13, cor=T["texto"], fundo=z)

    legenda(slide, MARGEM + Inches(6.85), Inches(4.50), Inches(4.65),
            "Duas linhas, porque a coluna nova só tem dois valores. O "
            ".reset_index() é o que traz eh_capital para dentro da tabela: sem "
            "ele, ela vira rótulo de linha e some do alcance do aes().",
            tamanho=12)
    return slide


def s04_proporcao_a(prs, lays):
    """O mesmo grafico por dois caminhos.

    A proporcao de acordaos por regime desenha exatamente a mesma coisa que o
    geom_bar, so que com o eixo y em fracao. E isso que isola a diferenca entre
    as duas geometrias: nao e o grafico que muda, e quem faz a conta. O slide
    seguinte troca a altura por outra coisa, e ai o desenho muda.
    """
    slide = novo_slide(prs, lays, "Contagem vs proporção", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "O mesmo gráfico por dois caminhos. Repare que o desenho é "
            "idêntico: só o eixo y muda de contagem para fração.")

    etiqueta(slide, MARGEM, Inches(2.42), Inches(5.30), Inches(0.40),
             "quantos? a geometria conta: geom_bar()", T["neutro"])
    bloco_codigo(slide, MARGEM, Inches(2.94), Inches(5.30), Inches(1.95), [
        "(",
        '    ggplot(penas, aes(x="regime"))',
        "    + geom_bar()",
        ")",
    ], tamanho=11.5, barra=T["neutro"])
    imagem_ajustada(slide, fig("retomada_bar"),
                    MARGEM, Inches(5.05), Inches(5.30), Inches(2.15),
                    moldura=False)

    x2 = MARGEM + Inches(5.75)
    etiqueta(slide, x2, Inches(2.42), Inches(5.75), Inches(0.40),
             "que fração? a conta é sua: geom_col()", T["marca"])
    bloco_codigo(slide, x2, Inches(2.94), Inches(5.75), Inches(1.95), [
        "resumo = (",
        '    penas.groupby("regime", as_index=False)',
        '    .agg(n=("processo", "size"))',
        ")",
        'resumo["prop"] = resumo["n"] / resumo["n"].sum()',
        "(",
        '    ggplot(resumo, aes(x="regime", y="prop"))',
        "    + geom_col()",
        ")",
    ], tamanho=9.5, barra=T["marca"])
    imagem_ajustada(slide, fig("prop_acordaos"),
                    x2, Inches(5.05), Inches(5.30), Inches(2.15),
                    moldura=False)
    return slide


def s04_proporcao(prs, lays):
    """A rodada 6: media de verdadeiro/falso e proporcao, e geom_col."""
    slide = novo_slide(prs, lays, "A altura pode ser qualquer conta", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Agora a coluna y é outra coisa, e o desenho muda junto. O "
            "geom_col() não serve só para refazer a contagem: ele desenha "
            "qualquer altura que você tenha calculado.")

    # A ideia de cada lado vai na etiqueta, e o codigo inteiro no bloco: era o
    # ggplot da direita que estava escrito em prosa, embaixo, onde ninguem le.
    # esquerda: contar
    etiqueta(slide, MARGEM, Inches(2.42), Inches(5.30), Inches(0.40),
             "quantos? a geometria conta: geom_bar()", CINZA_ESCURO)
    bloco_codigo(slide, MARGEM, Inches(2.94), Inches(5.30), Inches(1.95), [
        "(",
        '    ggplot(penas, aes(x="regime"))',
        "    + geom_bar()",
        ")",
    ], tamanho=11.5, barra=T["neutro"])
    imagem_ajustada(slide, fig("retomada_bar"),
                    MARGEM, Inches(5.05), Inches(5.30), Inches(2.15),
                    moldura=False)

    # direita: medir
    x2 = MARGEM + Inches(5.75)
    etiqueta(slide, x2, Inches(2.42), Inches(5.75), Inches(0.40),
             "que proporção é reincidente? geom_col()", T["marca"])
    bloco_codigo(slide, x2, Inches(2.94), Inches(5.75), Inches(1.95), [
        "resumo = (",
        '    penas.groupby("regime", as_index=False)',
        '    .agg(prop=("houve_reincidencia", "mean"))',
        ")",
        "(",
        '    ggplot(resumo, aes(x="regime", y="prop"))',
        "    + geom_col()",
        ")",
    ], tamanho=10, barra=T["marca"])
    imagem_ajustada(slide, fig("retomada_col"),
                    x2, Inches(5.05), Inches(5.30), Inches(2.15),
                    moldura=False)
    return slide


def s03_ideia(prs, lays):
    slide = novo_slide(prs, lays, "De uma variável para duas", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Na aula 5 todo gráfico tinha uma coluna só. Hoje entram duas, e a "
            "gramática não muda: o aes() passa a receber dois nomes.")

    caixa(slide, MARGEM, Inches(2.50), FAIXA, Inches(1.15), preenchimento=T["inverso"])
    tb = texto_livre(slide, MARGEM + Inches(0.40), Inches(2.76),
                     FAIXA - Inches(0.8), Inches(0.85))
    escrever(tb.text_frame, [
        {"texto": "O par de tipos escolhe a geometria.",
         "tamanho": 26, "bold": True, "cor": T["inverso_texto"], "depois": 5},
        {"texto": "Antes de escrever qualquer coisa, diga em voz alta o tipo "
                  "das duas variáveis.",
         "tamanho": 14, "cor": T["inverso_texto"], "depois": 0},
    ])

    pares = [
        ("numérica  ×  numérica", "quando uma cresce, o que a outra faz",
         "geom_point()", T["acento1"]),
        ("numérica  ×  categórica", "a distribuição muda entre as categorias",
         "geom_boxplot()", T["acento3"]),
        ("categórica  ×  categórica", "a composição muda entre as categorias",
         'geom_bar(position="fill")', T["acento2"]),
    ]
    y = Inches(4.05)
    for par, pergunta, geom, cor in pares:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(0.72), preenchimento=cor)
        tb = texto_livre(slide, MARGEM + Inches(0.28), y - Inches(0.02),
                         Inches(3.60), Inches(0.4))
        escrever(tb.text_frame, [
            {"texto": par, "tamanho": 16, "bold": True, "depois": 0}])
        tb = texto_livre(slide, MARGEM + Inches(4.00), y + Inches(0.02),
                         Inches(4.20), Inches(0.4))
        escrever(tb.text_frame, [
            {"texto": pergunta, "tamanho": 13, "cor": T["fraco"],
             "depois": 0}])
        tb = texto_livre(slide, MARGEM + Inches(8.40), y + Inches(0.02),
                         Inches(3.10), Inches(0.4))
        escrever(tb.text_frame, [
            {"texto": geom, "fonte": MONO, "tamanho": 13, "depois": 0}])
        y += Inches(0.82)
    return slide


def s04_forma_curta(prs, lays):
    slide = novo_slide(prs, lays, "Uma forma mais curta de escrever",
                             EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "O aes() pode ir dentro do ggplot(), como segundo argumento. "
            "As duas células produzem exatamente o mesmo gráfico.")

    largura = Inches(5.55)
    x2 = MARGEM + largura + Inches(0.40)

    etiqueta(slide, MARGEM, Inches(2.55), Inches(3.10), Inches(0.36),
             "a forma da aula 5", CINZA_ESCURO)
    bloco_codigo(slide, MARGEM, Inches(3.03), largura, Inches(1.45), [
        "(",
        "    ggplot(penas)",
        '    + aes(x="regime")',
        "    + geom_bar()",
        ")",
    ], barra=CINZA)
    legenda(slide, MARGEM, Inches(4.62), largura,
            "Separa as camadas, e é melhor para aprender.", tamanho=12)

    etiqueta(slide, x2, Inches(2.55), Inches(3.55), Inches(0.36),
             "a forma curta", T["marca"])
    bloco_codigo(slide, x2, Inches(3.03), largura, Inches(1.45), [
        "(",
        '    ggplot(penas, aes(x="regime"))',
        "    + geom_bar()",
        ")",
    ], barra=T["marca"])
    legenda(slide, x2, Inches(4.62), largura,
            "É a mais comum na prática, e a que você vai achar na internet e "
            "nos exercícios.", tamanho=12)

    caixa(slide, MARGEM, Inches(5.35), FAIXA, Inches(1.05),
          preenchimento=T["superficie"])
    tb = texto_livre(slide, MARGEM + Inches(0.30), Inches(5.52),
                     FAIXA - Inches(0.6), Inches(0.8))
    escrever(tb.text_frame, [
        {"texto": "A regra da aula 5 continua valendo, sem exceção.",
         "tamanho": 15, "bold": True, "depois": 3},
        {"texto": "Nome de coluna dentro do aes(), valor fixo fora dele. "
                  "Mudou onde o aes() está escrito, não o que ele faz.",
         "tamanho": 13, "cor": T["fraco"], "depois": 0},
    ])
    return slide


def s05b_histograma(prs, lays):
    """Fecha a parte de uma variavel: a numerica sozinha."""
    slide = novo_slide(prs, lays, "Uma numérica: o histograma", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Categórica tem categorias para contar. Numérica não tem: o que se "
            "faz é cortar o intervalo em faixas e contar cada faixa.")

    bloco_codigo(slide, MARGEM, Inches(2.48), Inches(6.30), Inches(1.35), [
        "(",
        '    ggplot(penas, aes(x="pena_anos"))',
        "    + geom_histogram(bins=20)",
        ")",
    ], tamanho=12, barra=T["acento2"])

    partes = [
        ("cada barra", "uma faixa de valores, não uma categoria"),
        ("a altura", "quantos casos caem naquela faixa"),
        ("bins=", "quantas faixas. Muda o que o gráfico deixa ver"),
        ("não há y", "a altura é contada pela geometria"),
    ]
    y = Inches(4.25)
    for nome, oque in partes:
        tb = texto_livre(slide, MARGEM, y, Inches(1.90), Inches(0.32))
        escrever(tb.text_frame, [
            {"texto": nome, "tamanho": 13, "bold": True, "depois": 0}])
        tb = texto_livre(slide, MARGEM + Inches(2.00), y, Inches(4.30),
                         Inches(0.32))
        escrever(tb.text_frame, [
            {"texto": oque, "tamanho": 13, "cor": T["fraco"], "depois": 0}])
        y += Inches(0.42)

    imagem_ajustada(slide, fig("histograma"),
                    MARGEM + Inches(6.85), Inches(2.40), Inches(4.65),
                    Inches(3.40), moldura=False)
    legenda(slide, MARGEM + Inches(6.85), Inches(5.95), Inches(4.65),
            "Rode com bins=5 e com bins=60. Não existe número certo: existe o "
            "que responde à sua pergunta.", tamanho=12)
    return slide


def s05_pontos(prs, lays):
    """Duas numericas. Vem antes do boxplot: e a leitura mais direta de duas
    variaveis, uma em cada eixo, sem forma nova para aprender."""
    slide = novo_slide(prs, lays, "Duas numéricas: pontos", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Uma variável em cada eixo, e um ponto por linha da tabela. É o "
            "gráfico mais direto de duas variáveis.")

    bloco_codigo(slide, MARGEM, Inches(2.48), Inches(6.30), Inches(1.55), [
        "(",
        '    ggplot(penas, aes(x="n_palavras_ementa", y="pena_anos"))',
        "    + geom_point(alpha=0.4)",
        '    + geom_smooth(method="lm")',
        ")",
    ], tamanho=12, barra=T["acento1"])

    partes = [
        ("cada ponto", "um acórdão, com as duas medidas dele"),
        ("alpha=0.4", "transparência. Fora do aes(): é valor fixo"),
        ("geom_smooth", "a reta que resume a direção da nuvem"),
        ("a leitura", "a reta é quase plana, e isso é uma resposta"),
    ]
    y = Inches(4.45)
    for nome, oque in partes:
        tb = texto_livre(slide, MARGEM, y, Inches(1.90), Inches(0.32))
        escrever(tb.text_frame, [
            {"texto": nome, "tamanho": 13, "bold": True, "depois": 0}])
        tb = texto_livre(slide, MARGEM + Inches(2.00), y, Inches(4.30),
                         Inches(0.32))
        escrever(tb.text_frame, [
            {"texto": oque, "tamanho": 13, "cor": T["fraco"], "depois": 0}])
        y += Inches(0.42)

    imagem_ajustada(slide, fig("pontos_smooth"),
                    MARGEM + Inches(6.85), Inches(2.40), Inches(4.65),
                    Inches(3.40), moldura=False)
    legenda(slide, MARGEM + Inches(6.85), Inches(5.95), Inches(4.65),
            "O tamanho da ementa não diz quase nada sobre a pena. Gráfico que "
            "não mostra relação também responde à pergunta.", tamanho=12)
    return slide

def s06_nao_diz(prs, lays):
    slide = novo_slide(prs, lays, "O que a reta não diz", EYEBROW)

    legenda(slide, MARGEM, Inches(1.98), FAIXA,
            "Duas leituras que o gráfico anterior não autoriza, e que vão "
            "aparecer no relatório de vocês se ninguém avisar.")

    itens = [
        ("Reta plana não prova que não há relação",
         "Prova que não há relação linear nestes 197 acórdãos, que são os que "
         "tinham regime e pena legíveis na ementa. Mais da metade da base "
         "ficou de fora, e ela não saiu por sorteio."),
        ("Reta inclinada não provaria causa",
         "Se ementas longas viessem com penas altas, o mais provável seria que "
         "caso grave gera ementa longa e pena alta. A reta mede associação, e "
         "associação não é causa."),
    ]
    y = Inches(2.60)
    for titulo, texto in itens:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(1.35),
              preenchimento=T["marca"])
        tb = texto_livre(slide, MARGEM + Inches(0.30), y - Inches(0.02),
                         FAIXA - Inches(0.5), Inches(1.4))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 18, "bold": True, "depois": 4},
            {"texto": texto, "tamanho": 13.5, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(1.70)

    caixa(slide, MARGEM, Inches(6.00), FAIXA, Inches(0.80),
          preenchimento=T["superficie"])
    tb = texto_livre(slide, MARGEM + Inches(0.30), Inches(6.18),
                     FAIXA - Inches(0.6), Inches(0.55))
    escrever(tb.text_frame, [
        {"texto": "No exercício 3 do notebook vocês escrevem as duas linhas: "
                  "o que o gráfico mostra, e o que ele não permite concluir.",
         "tamanho": 14, "bold": True, "depois": 0},
    ])
    return slide


def s07_boxplot(prs, lays):
    slide = novo_slide(prs, lays, "Numérica e categórica: caixas",
                             EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Uma caixa por categoria. Cada caixa é o resumo da aula 3, "
            "desenhado.")

    bloco_codigo(slide, MARGEM, Inches(2.48), Inches(6.30), Inches(1.55), [
        "(",
        '    ggplot(penas, aes(x="regime", y="pena_anos"))',
        "    + geom_boxplot()",
        "    + coord_flip()",
        ")",
    ], tamanho=12, barra=T["acento3"])

    partes = [
        ("a linha do meio", "a mediana"),
        ("a caixa", "do primeiro ao terceiro quartil"),
        ("os fios", "até os valores ainda típicos"),
        ("os pontos soltos", "os distantes do resto"),
    ]
    y = Inches(4.25)
    for nome, oque in partes:
        tb = texto_livre(slide, MARGEM, y, Inches(2.30), Inches(0.32))
        escrever(tb.text_frame, [
            {"texto": nome, "tamanho": 13, "bold": True, "depois": 0}])
        tb = texto_livre(slide, MARGEM + Inches(2.40), y, Inches(3.90),
                         Inches(0.32))
        escrever(tb.text_frame, [
            {"texto": oque, "tamanho": 13, "cor": T["fraco"], "depois": 0}])
        y += Inches(0.42)

    imagem_ajustada(slide, fig("boxplot"),
                    MARGEM + Inches(6.85), Inches(2.30), Inches(4.65),
                    Inches(3.60), moldura=False)
    return slide


def s08a_position_codigo(prs, lays):
    """Os tres codigos completos, nas mesmas colunas em que o slide seguinte
    poe os tres graficos. Ver o codigo e depois o desenho no mesmo lugar e o
    que deixa a comparacao imediata."""
    slide = novo_slide(prs, lays, "Duas categóricas: o position",
                             EYEBROW)

    legenda(slide, MARGEM, Inches(1.90), FAIXA,
            "As mesmas duas variáveis, a mesma geometria, três gráficos "
            "diferentes. Só muda o argumento position.")

    casos = [
        ('position="stack"', T["acento1"], "stack",
         "o padrão: quantos casos, empilhados"),
        ('position="fill"', T["marca"], "fill",
         "todas as barras com altura 1: proporção"),
        ('position="dodge"', T["acento3"], "dodge",
         "lado a lado: contagem, sem empilhar"),
    ]

    largura = Inches(3.70)
    passo = largura + Inches(0.20)
    for i, (rotulo, cor, valor, oque) in enumerate(casos):
        x = MARGEM + i * passo
        etiqueta(slide, x, Inches(2.45), largura, Inches(0.40), rotulo, cor)
        bloco_codigo(slide, x, Inches(3.00), largura, Inches(2.05), [
            "(",
            "    ggplot(",
            "        penas,",
            '        aes(x="regime",',
            '            fill="houve_reincidencia")',
            "    )",
            f'    + geom_bar(position="{valor}")',
            ")",
        ], tamanho=10, barra=cor)
        tb = texto_livre(slide, x, Inches(5.20), largura, Inches(0.70))
        escrever(tb.text_frame, [{"texto": oque, "tamanho": 12.5,
                                  "cor": T["fraco"], "depois": 0}])
    return slide

def s08b_position_resultado(prs, lays):
    """Os tres resultados lado a lado."""
    slide = novo_slide(prs, lays, "Os três, lado a lado", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "A resposta é o fill: ele é o único em que a comparação entre "
            "regimes não é atrapalhada pelo tamanho de cada grupo.")

    largura = Inches(3.65)
    for i, (nome, rotulo, cor) in enumerate([
        ("stack", 'position="stack"', T["acento1"]),
        ("fill", 'position="fill"', T["marca"]),
        ("dodge", 'position="dodge"', T["acento3"]),
    ]):
        x = MARGEM + i * (largura + Inches(0.28))
        etiqueta(slide, x, Inches(2.50), largura, Inches(0.40), rotulo, cor)
        imagem_ajustada(slide, fig(nome),
                        x, Inches(3.02), largura, Inches(2.85), moldura=False)

    painel(slide, MARGEM, Inches(6.10), FAIXA, Inches(0.85), [
        {"texto": 'position="fill" esconde quanta gente tem em cada barra.',
         "tamanho": 15, "bold": True, "depois": 3},
        {"texto": "Uma barra com 4 casos e outra com 400 ficam do mesmo "
                  "tamanho. Quando a proporção interessa, ele é o certo; "
                  "quando o tamanho do grupo importa, ele mente por omissão.",
         "tamanho": 12.5, "cor": T["fraco"], "depois": 0},
    ])
    return slide

def s09_preparar(prs, lays):
    slide = novo_slide(prs, lays, "Preparar a tabela antes do gráfico",
                             EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Quando a altura da barra é uma conta, o pandas faz a conta e o "
            "plotnine só desenha. São três operações, e elas voltam sempre.")

    itens = [
        (".assign(nova=expressao)", "criar uma coluna sem quebrar o "
         "encadeamento", T["acento1"]),
        (".groupby(col, as_index=False)", "agrupar já saindo com a coluna "
         "dentro da tabela, sem precisar do .reset_index()", T["acento1"]),
        ('reorder("categoria", "valor")', "ordenar as barras, escrito dentro "
         "do aes(). Ordenar a tabela NÃO ordena as barras.", T["acento2"]),
    ]
    y = Inches(2.50)
    for op, texto, cor in itens:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(0.90), preenchimento=cor)
        tb = texto_livre(slide, MARGEM + Inches(0.28), y - Inches(0.02),
                         Inches(6.20), Inches(0.95))
        escrever(tb.text_frame, [
            {"texto": op, "fonte": MONO, "tamanho": 14, "bold": True,
             "depois": 3},
            {"texto": texto, "tamanho": 12.5, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(1.05)

    bloco_codigo(slide, MARGEM, Inches(5.75), Inches(6.30), Inches(0.98), [
        'ggplot(resumo, aes(x="reorder(comarca, pena_mediana)", '
        'y="pena_mediana"))',
        "+ geom_col(fill=\"#E50505\") + coord_flip()",
    ], tamanho=10.5, barra=T["acento2"])

    imagem_ajustada(slide, fig("reorder"),
                    MARGEM + Inches(6.85), Inches(2.20), Inches(4.65),
                    Inches(4.35), moldura=False)
    return slide


def s10_tabela(prs, lays):
    slide = novo_slide(prs, lays, "Que gráfico para que par", EYEBROW)

    linhas = [
        ("numérica × numérica", "quando uma cresce, o que a outra faz",
         "geom_point() + geom_smooth()", T["acento1"]),
        ("numérica × numérica, x é tempo", "como evolui", "geom_line()",
         T["acento1"]),
        ("numérica × categórica", "a distribuição muda entre as categorias",
         "geom_boxplot()", T["acento3"]),
        ("numérica × categórica, já resumida", "comparar um valor por categoria",
         "geom_col()", T["acento3"]),
        ("categórica × categórica", "a composição muda entre as categorias",
         'geom_bar(position="fill")', T["acento2"]),
        ("categórica × categórica, contagem", "quantos casos em cada combinação",
         'geom_bar(position="dodge")', T["acento2"]),
    ]

    tabela_shape = slide.shapes.add_table(
        len(linhas) + 1, 4, MARGEM, Inches(2.15), FAIXA, Inches(3.95))
    tabela = tabela_shape.table
    tabela_sem_estilo(tabela)
    for larg, i in zip([Inches(0.14), Inches(3.85), Inches(3.75),
                        Inches(3.76)], range(4)):
        tabela.columns[i].width = larg
    tabela.rows[0].height = Inches(0.32)
    for r in range(1, len(linhas) + 1):
        tabela.rows[r].height = Inches(0.60)

    for c, t in enumerate(["", "as duas variáveis", "a pergunta",
                           "a geometria"]):
        celula(tabela, 0, c, t, tamanho=11, bold=True, cor=T["etiqueta_texto"], fundo=T["inverso"])
    for r, (par, perg, geom, cor) in enumerate(linhas, start=1):
        zebra = T["superficie"] if r % 2 == 0 else None
        celula(tabela, r, 0, "", fundo=cor)
        celula(tabela, r, 1, par, tamanho=12, bold=True, fundo=zebra)
        celula(tabela, r, 2, perg, tamanho=12, cor=T["fraco"], fundo=zebra)
        celula(tabela, r, 3, geom, tamanho=12, bold=True, fundo=zebra)

    caixa(slide, MARGEM, Inches(6.30), FAIXA, Inches(0.72),
          preenchimento=T["superficie"])
    tb = texto_livre(slide, MARGEM + Inches(0.30), Inches(6.45),
                     FAIXA - Inches(0.6), Inches(0.5))
    escrever(tb.text_frame, [
        {"texto": "A terceira variável, quando existe, entra em color ou fill "
                  "dentro do aes(), ou em facet_wrap() quando ficar carregado.",
         "tamanho": 13.5, "bold": True, "depois": 0},
    ])
    return slide


def s11_projeto(prs, lays):
    slide = novo_slide(prs, lays, "Projeto 02", EYEBROW)

    legenda(slide, MARGEM, Inches(1.94), FAIXA,
            "Em grupo, valendo nota, uma hora. É a Gincana do Pipeline de "
            "terça-feira, no computador.")

    passos = [
        ("Seis desafios", "cada um com uma pergunta e a imagem do gráfico "
         "que vocês precisam produzir. Dá para avançar e voltar."),
        ("Valem 11 pontos", "o primeiro vale 1 e os outros 2. Cada tentativa "
         "errada desconta 0,25, e a nota para em 10."),
        ("Blocos fora de ordem", "de pandas e de plotnine. Você escolhe quais "
         "usar e em que ordem."),
        ("Tem bloco que não serve", "e o app explica por que, quando você usa. "
         "É o distrator da gincana."),
        ("Entrega hoje", "no fim da aula. Não há prazo no dia seguinte."),
    ]
    y = Inches(2.55)
    for titulo, texto in passos:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(0.82),
              preenchimento=T["marca"])
        tb = texto_livre(slide, MARGEM + Inches(0.28), y - Inches(0.02),
                         Inches(6.10), Inches(0.88))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 15, "bold": True, "depois": 2},
            {"texto": texto, "tamanho": 12.5, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(0.98)

    return slide


def s12_fechamento(prs, lays):
    slide = novo_slide(prs, lays, "Antes de sair", EYEBROW)

    itens = [
        ("Submeta pelo BlackBoard", "A entrega é hoje, no fim da aula. O app "
         "mostra o que foi registrado na sua sessão.", T["marca"]),
        ("Rode o notebook inteiro em casa", "Principalmente os exercícios 1 a "
         "3, que são os que pedem a leitura escrita do gráfico.", T["acento2"]),
        ("Aula 7", "Probabilidade. Volta o papel e a caneta, e some o "
         "computador.", CINZA_ESCURO),
    ]
    y = Inches(2.40)
    for titulo, texto, cor in itens:
        caixa(slide, MARGEM, y, Inches(0.09), Inches(1.15), preenchimento=cor)
        tb = texto_livre(slide, MARGEM + Inches(0.32), y - Inches(0.02),
                         FAIXA - Inches(0.5), Inches(1.2))
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 19, "bold": True, "depois": 4},
            {"texto": texto, "tamanho": 14, "cor": T["fraco"], "depois": 0},
        ])
        y += Inches(1.45)
    return slide


# ---------------------------------------------------------------------- main

def montar(escuro: bool) -> str:
    global T
    T = ESCURO if escuro else CLARO

    prs = deck_limpo()
    lays = layouts(prs)

    # O boxplot vem DEPOIS das barras: a aula sobe de uma variavel para duas
    # pelo caminho mais curto, que e acrescentar cor a um grafico de barras que
    # a turma ja sabe ler. A caixa e forma nova, e entra depois.
    #
    # Saiu o slide de `reorder()`: nao ha tempo para ele, e por isso ele
    # tambem saiu do Projeto 02 e virou material extra no notebook.
    for construir in (s01_capa,
                      s02b_retomada_pipeline, s02c_retomada_grafico,
                      s03_assign, s04_proporcao_a, s04_proporcao,
                      s05b_histograma,
                      s03_ideia, s04_forma_curta,
                      s08a_position_codigo, s08b_position_resultado,
                      s05_pontos, s07_boxplot,
                      s09b_catalogo_uni, s09c_catalogo_bi,
                      s09b_univariados, s10_tabela, s09c_datatoviz):
        construir(prs, lays)

    saida = SAIDA.replace(".pptx", "_dark.pptx") if escuro else SAIDA
    n = gravar(prs, saida, titulo="Aula 6: gráficos de duas variáveis")
    print(f"{saida}  ({n} slides)")
    return saida


def main():
    montar(escuro=False)
    montar(escuro=True)


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    main()

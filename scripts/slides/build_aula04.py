"""
Monta os slides da Aula 4: encadear operacoes.

O deck nao ensina: a aula e o notebook (notebooks/aula04_*.ipynb). Estes slides
servem para duas coisas so, e ficam projetados enquanto o aluno programa:

  1. lembrar as funcoes vistas ate aqui, num lugar so;
  2. mostrar como se le a documentacao, com capturas do site do pandas.

Uso:
    python build_aula04.py
"""

from __future__ import annotations

import os

from pptx.util import Inches

from insper import (
    BRANCO, CINZA, CINZA_CLARO, CINZA_ESCURO, DISPLAY, FAIXA, LARANJA, MARGEM,
    PRETO, QUASE_BRANCO, ROXO, TEXTO, TURQUESA, VERMELHO, caixa, celula,
    deck_limpo, escrever, gravar, imagem_ajustada, layouts, slide_capa,
    slide_com_titulo, tabela_sem_estilo, texto_livre,
)

AQUI = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(AQUI, "assets")
SAIDA = os.path.join(AQUI, "aula04.pptx")

EYEBROW = "Ciência de Dados Aplicada ao Direito II  ·  Aula 4  ·  Encadear operações"
TITULO_DECK = "Aula 4: encadear operações"

COR_LER = TURQUESA
COR_RESUMIR = LARANJA
COR_ENCADEAR = ROXO


def tabela_de_funcoes(slide, linhas, topo, altura, larguras, cabecalho):
    tabela_shape = slide.shapes.add_table(
        len(linhas) + 1, 4, MARGEM, topo, FAIXA, altura
    )
    tabela = tabela_shape.table
    tabela_sem_estilo(tabela)
    for largura, i in zip(larguras, range(4)):
        tabela.columns[i].width = largura
    tabela.rows[0].height = Inches(0.26)
    alt = int((altura - Inches(0.26)) / len(linhas))
    for r in range(1, len(linhas) + 1):
        tabela.rows[r].height = alt

    for c, titulo in enumerate(cabecalho):
        celula(tabela, 0, c, titulo, tamanho=10.5, bold=True, cor=BRANCO, fundo=PRETO)
    for r, (cor, funcao, oque, quando) in enumerate(linhas, start=1):
        zebra = QUASE_BRANCO if r % 2 == 0 else None
        celula(tabela, r, 0, "", fundo=cor)
        celula(tabela, r, 1, funcao, tamanho=11, bold=True, fundo=zebra)
        celula(tabela, r, 2, oque, tamanho=11, fundo=zebra)
        celula(tabela, r, 3, quando, tamanho=11, cor=CINZA_ESCURO, fundo=zebra)
    return tabela


def s01_capa(prs, lays):
    return slide_capa(
        prs, lays,
        subtitulo="Aula 4: encadear operações",
        subtema="E como ler a documentação de uma função",
    )


def s02_funcoes_aula_2_3(prs, lays):
    slide = slide_com_titulo(prs, lays, "O que você já sabe fazer", EYEBROW)

    tb = texto_livre(slide, MARGEM, Inches(2.00), FAIXA, Inches(0.30))
    escrever(tb.text_frame, [
        {"texto": "Das aulas 2 e 3. Vale deixar este slide aberto enquanto você "
                  "programa.",
         "tamanho": 13, "cor": CINZA_ESCURO, "depois": 0},
    ])

    linhas = [
        (COR_LER, "pd.read_csv(caminho)", "lê um arquivo e devolve a tabela",
         "aula 2"),
        (COR_LER, ".info()", "colunas, faltantes e dtype de cada uma", "aula 2"),
        (COR_LER, ".nunique()", "quantos valores distintos por coluna", "aula 2"),
        (COR_LER, ".astype(\"string\")", "converte a coluna, aqui para texto", "aula 2"),
        (COR_LER, "pd.to_datetime(coluna)", "de texto para data", "aula 2"),
        (COR_LER, ".dt.year, .dt.month, .dt.days", "pedaços de uma data ou duração",
         "aula 2"),
        (COR_LER, "pd.Categorical(coluna, categories=, ordered=)",
         "declara as categorias e a ordem", "aula 2"),
        (COR_LER, ".cat.add_categories() e .fillna()", "abre uma categoria e preenche",
         "aula 2"),
        (COR_LER, "pd.cut(coluna, bins=, labels=)", "de numérica para faixas",
         "aula 2"),
        (COR_RESUMIR, ".describe()", "resumo rápido de uma coluna numérica", "aula 3"),
        (COR_RESUMIR, ".mean(), .median(), .mode()", "medidas de posição", "aula 3"),
        (COR_RESUMIR, ".std(ddof=), .var(), .quantile()", "medidas de dispersão",
         "aula 3"),
        (COR_RESUMIR, ".value_counts(normalize=)", "contagem ou proporção por categoria",
         "aula 3"),
        (COR_RESUMIR, "df[condicao]", "índice lógico, o filtro mais explícito",
         "aula 3"),
    ]
    tabela_de_funcoes(
        slide, linhas, Inches(2.38), Inches(4.28),
        (Inches(0.13), Inches(4.55), Inches(5.42), Inches(1.40)),
        ("", "função", "o que faz", "onde apareceu"),
    )
    return slide


def s03_funcoes_aula_4(prs, lays):
    slide = slide_com_titulo(prs, lays, "O que entra hoje", EYEBROW)

    linhas = [
        (COR_ENCADEAR, ".query(\"coluna == 'valor'\")", "escolhe linhas por condição",
         "filtrar"),
        (COR_ENCADEAR, ".dropna(subset=[...])", "descarta linhas sem valor",
         "filtrar"),
        (COR_ENCADEAR, ".assign(nova=lambda d: ...)", "cria coluna", "transformar"),
        (COR_ENCADEAR, "[[\"a\", \"b\"]]", "escolhe colunas", "transformar"),
        (COR_ENCADEAR, ".sort_values(\"a\", ascending=False)", "ordena", "transformar"),
        (COR_ENCADEAR, ".groupby(\"a\").agg(s=(\"b\", \"mean\"))",
         "uma linha por grupo", "agregar"),
        (COR_ENCADEAR, ".reset_index()", "tira o agrupamento do índice", "agregar"),
        (COR_ENCADEAR, ".head(n)", "corta as primeiras linhas", "ver"),
    ]
    tabela_de_funcoes(
        slide, linhas, Inches(2.05), Inches(2.72),
        (Inches(0.13), Inches(4.75), Inches(4.62), Inches(2.00)),
        ("", "função", "o que faz", "para quê"),
    )

    caixa(slide, MARGEM, Inches(4.98), FAIXA, Inches(1.66),
          preenchimento=QUASE_BRANCO)
    caixa(slide, MARGEM, Inches(4.98), Inches(0.09), Inches(1.66),
          preenchimento=COR_ENCADEAR)

    tb = texto_livre(slide, Inches(1.50), Inches(5.14), Inches(4.6), Inches(1.36))
    escrever(tb.text_frame, [
        {"texto": "O formato", "tamanho": 14, "bold": True, "depois": 5},
        {"texto": "resultado = (\n    tabela\n    .operacao_1(...)\n    .operacao_2(...)\n)",
         "fonte": "Consolas", "tamanho": 11, "entrelinhas": 1.0, "depois": 0},
    ])

    tb = texto_livre(slide, Inches(6.55), Inches(5.20), Inches(5.85), Inches(1.28))
    escrever(tb.text_frame, [
        {"texto": "Os parênteses existem para você poder quebrar a linha: sem "
                  "eles, a quebra encerra o comando.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "entrelinhas": 1.12, "depois": 8},
        {"texto": "Dentro do encadeamento, olhe para as colunas com lambda d:, e "
                  "nunca pelo nome da tabela original.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "entrelinhas": 1.12, "depois": 0},
    ])
    return slide


def s04_assinatura(prs, lays):
    slide = slide_com_titulo(
        prs, lays, "Toda função tem mais do que você usa", EYEBROW
    )

    imagem_ajustada(
        slide, os.path.join(ASSETS, "doc_sort_values.png"),
        MARGEM, Inches(2.02), Inches(6.55), Inches(4.45),
    )

    tb = texto_livre(slide, Inches(8.15), Inches(2.10), Inches(4.55), Inches(4.2))
    escrever(tb.text_frame, [
        {"texto": "Você escreve", "tamanho": 12.5, "cor": CINZA_ESCURO, "depois": 3},
        {"texto": ".sort_values(\"pena_anos\",\n              ascending=False)",
         "fonte": "Consolas", "tamanho": 12, "entrelinhas": 1.05, "depois": 12},
        {"texto": "A função aceita oito parâmetros.",
         "tamanho": 14, "bold": True, "depois": 10},
        {"texto": "A primeira linha da página é a assinatura, e ela responde três "
                  "perguntas: quais argumentos existem, em que ordem, e qual o "
                  "valor padrão de cada um.",
         "tamanho": 13, "entrelinhas": 1.15, "depois": 10},
        {"texto": "Repare no na_position='last': é ele que decide onde ficam os "
                  "valores faltantes quando você ordena. Ninguém adivinha isso, "
                  "está escrito ali.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "entrelinhas": 1.15, "depois": 0},
    ])
    return slide


def s05_como_ler(prs, lays):
    slide = slide_com_titulo(prs, lays, "As quatro partes de uma página", EYEBROW)

    imagem_ajustada(
        slide, os.path.join(ASSETS, "doc_assign.png"),
        MARGEM, Inches(2.02), Inches(8.05), Inches(4.35),
    )

    partes = [
        ("Parameters", "o que a função aceita, com o tipo e o padrão de cada um"),
        ("Returns", "o que ela devolve, e isso decide o que dá para encadear depois"),
        ("See also", "as funções vizinhas, que muitas vezes são a que você queria"),
        ("Examples", "código pronto para copiar e adaptar, no fim da página"),
    ]
    y = Inches(2.15)
    for i, (titulo, detalhe) in enumerate(partes):
        tb = texto_livre(slide, Inches(9.45), y + i * Inches(1.02),
                         Inches(3.25), Inches(0.95))
        escrever(tb.text_frame, [
            {"texto": titulo, "fonte": DISPLAY, "tamanho": 16, "bold": True,
             "cor": VERMELHO, "depois": 2},
            {"texto": detalhe, "tamanho": 12, "cor": CINZA_ESCURO,
             "entrelinhas": 1.12, "depois": 0},
        ])

    tb = texto_livre(slide, MARGEM, Inches(6.42), FAIXA, Inches(0.28))
    escrever(tb.text_frame, [
        {"texto": "Comece pelos Examples e volte para os Parameters quando "
                  "precisar mudar alguma coisa.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def s06_onde_procurar(prs, lays):
    slide = slide_com_titulo(prs, lays, "Onde procurar", EYEBROW)

    imagem_ajustada(
        slide, os.path.join(ASSETS, "doc_indice.png"),
        MARGEM, Inches(2.02), Inches(8.05), Inches(4.05),
    )

    caminhos = [
        ("No site", "pandas.pydata.org, em API reference. A lista à direita da "
                    "página separa por assunto: descriptive stats, groupby, "
                    "reshaping, missing data."),
        ("No notebook", "escreva o nome da função com ? no fim e rode a célula. "
                        "A documentação abre ali mesmo, sem trocar de janela."),
        ("Na busca", "procure por \"pandas\" mais o que você quer fazer, em "
                     "inglês. O primeiro resultado costuma ser a página oficial."),
    ]
    y = Inches(2.12)
    for i, (titulo, detalhe) in enumerate(caminhos):
        tb = texto_livre(slide, Inches(9.45), y + i * Inches(1.38),
                         Inches(3.25), Inches(1.30))
        escrever(tb.text_frame, [
            {"texto": titulo, "fonte": DISPLAY, "tamanho": 16, "bold": True,
             "cor": VERMELHO, "depois": 2},
            {"texto": detalhe, "tamanho": 11.5, "cor": CINZA_ESCURO,
             "entrelinhas": 1.12, "depois": 0},
        ])

    caixa(slide, MARGEM, Inches(6.22), Inches(8.05), Inches(0.42),
          preenchimento=QUASE_BRANCO)
    tb = texto_livre(slide, Inches(1.40), Inches(6.30), Inches(7.7), Inches(0.28))
    escrever(tb.text_frame, [
        {"texto": "criminal.groupby?          help(pd.cut)",
         "fonte": "Consolas", "tamanho": 13, "depois": 0},
    ])
    return slide


def main():
    prs = deck_limpo()
    lays = layouts(prs)

    for construir in (
        s01_capa, s02_funcoes_aula_2_3, s03_funcoes_aula_4, s04_assinatura,
        s05_como_ler, s06_onde_procurar,
    ):
        construir(prs, lays)

    gravar(prs, SAIDA, titulo=TITULO_DECK)


if __name__ == "__main__":
    import sys

    sys.stdout.reconfigure(encoding="utf-8")
    main()

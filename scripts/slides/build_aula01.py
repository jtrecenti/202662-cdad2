"""
Monta os slides da Aula 1 de Ciencia de Dados Aplicada ao Direito II (2026.2).

Base: template oficial do Insper (Apresentacao didatica - Final.potx, ja
convertido para .pptx em cdadeng/assets_insper). Cores e tipografia seguem o
Manual de Uso da Marca (insper-guia-de-marca.pdf, secao 2.1).

Uso:
    python build_aula01.py
"""

from __future__ import annotations

import os

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from insper import (
    AMARELO, AUTOR, BASE, BRANCO, CINZA, CINZA_CLARO, CINZA_ESCURO, DISPLAY,
    FAIXA, LARANJA, MARGEM, OA_BASES, OA_DESCRITIVAS, OA_MODELOS, OA_NEGOCIO,
    OA_PLANOS, PRETO, ROXO, TEXTO, TOPO, TURQUESA, VERDE, VERMELHO, caixa,
    celula, deck_limpo, escrever, gravar, imagem_ajustada, layouts,
    limpar_placeholders_vazios, slide_limpo, slide_secao, tabela_sem_estilo,
    texto_livre,
)
from insper import slide_com_titulo as _slide_com_titulo

AQUI = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(AQUI, "assets")
SAIDA = os.path.join(AQUI, "aula01.pptx")

EYEBROW = "Ciência de Dados Aplicada ao Direito II  ·  Aula 1  ·  Sobre a disciplina"
TITULO_DECK = "Aula 1: sobre a disciplina"
BUSCA_GOOGLE = "https://www.google.com/search?q=jurimetria"


def slide_com_titulo(prs, lays, titulo, eyebrow=EYEBROW):
    return _slide_com_titulo(prs, lays, titulo, eyebrow)


# ---------------------------------------------------------------- conteudo

OBJETIVOS = [
    (OA_PLANOS, "elaborar planos de estudo para solucionar problemas jurídicos",
     "operacionalizando conceitos, formulando perguntas de pesquisa, "
     "selecionando fontes de dados e planejando a coleta"),
    (OA_BASES, "construir e organizar bases de dados jurídicas",
     "com raspagem de dados, expressões regulares, transformação de dados, "
     "anotação manual e inteligência artificial"),
    (OA_DESCRITIVAS, "realizar análises descritivas relevantes para público leigo",
     "escolhendo o gráfico e a medida adequados a cada situação, "
     "para se comunicar com pessoas não técnicas"),
    (OA_MODELOS, "desenvolver modelos estatísticos para auxiliar na tomada de decisão",
     "interpretando parâmetros, avaliando significância e validando o desempenho "
     "com as métricas relevantes"),
    (OA_NEGOCIO, "aplicar soluções de ciência de dados em problemas de negócio",
     "provisionamento de carteiras, estratégias processuais e discussão de "
     "interpretabilidade, viés e LGPD"),
]

AVALIACAO = [
    ("PI", "Prova Intermediária", "25%", AMARELO,
     "Aulas 14 e 15\n24/09 e 29/09"),
    ("PF", "Prova Final", "35%", VERMELHO,
     "Aulas 30 e 31\n26/11 e 01/12"),
    ("PC", "Pesquisa de Campo", "20%", TURQUESA,
     "Estudo empírico em grupo:\nrelatório e debate"),
    ("PA", "Projetos Aplicados", "20%", ROXO,
     "8 projetos ao longo\ndo semestre, 12,5% cada"),
]

# (numero, data, tema, projeto, cor do objetivo de aprendizagem)
AULAS_PI = [
    ("01", "11/08", "Sobre a disciplina: ciência de dados no Direito e variável aleatória", "-", OA_NEGOCIO),
    ("02", "13/08", "Tipos de variáveis e sua representação no pandas", "-", OA_BASES),
    ("03", "18/08", "Medidas de tendência central e definição de variáveis a partir do texto", "01", OA_BASES),
    ("04", "20/08", "Transformação e agregação de dados com o pandas", "-", OA_DESCRITIVAS),
    ("05", "25/08", "Gramática de gráficos e introdução ao plotnine", "-", OA_DESCRITIVAS),
    ("06", "27/08", "Gráficos bivariados: dispersão, linhas e barras agrupadas", "02", OA_DESCRITIVAS),
    ("07", "01/09", "Probabilidade e incerteza: espaço amostral e independência", "-", OA_MODELOS),
    ("08", "03/09", "Distribuições de probabilidade e variáveis aleatórias", "-", OA_MODELOS),
    ("09", "08/09", "Inferência: teorema do limite central e intervalo de confiança", "-", OA_MODELOS),
    ("10", "10/09", "Intervalos de confiança e testes de hipóteses", "03", OA_MODELOS),
    ("11", "15/09", "Fontes de dados jurídicas e planejamento de pesquisa", "-", OA_PLANOS),
    ("12", "17/09", "Estruturação de dados e formulação da pergunta de pesquisa", "-", OA_PLANOS),
    ("13", "22/09", "Revisão integrada: variáveis, descritiva, probabilidade e inferência", "04", OA_MODELOS),
    ("14 e 15", "24 e 29/09", "PROVA INTERMEDIÁRIA", "-", None),
]

AULAS_PF = [
    ("16", "06/10", "Testes de hipóteses aplicados: teste t, proporções e qui-quadrado", "-", OA_MODELOS),
    ("17", "08/10", "Análise de variância (ANOVA) e comparações múltiplas", "-", OA_MODELOS),
    ("18", "13/10", "Correlação e regressão linear simples", "05", OA_MODELOS),
    ("19", "20/10", "Regressão linear múltipla, dummies e causalidade", "-", OA_MODELOS),
    ("20", "22/10", "Regressão logística e razão de chances", "-", OA_MODELOS),
    ("21", "27/10", "Matriz de confusão, medidas de erro e ponto de corte", "06", OA_MODELOS),
    ("22", "29/10", "Introdução a machine learning, validação cruzada e curva ROC", "-", OA_MODELOS),
    ("23", "03/11", "Classificação com scikit-learn e regularização", "-", OA_MODELOS),
    ("24", "05/11", "Modelos baseados em árvores: árvores e florestas aleatórias", "07 e 08", OA_MODELOS),
    ("25", "10/11", "Interpretabilidade, justiça algorítmica, viés e LGPD", "-", OA_NEGOCIO),
    ("26", "12/11", "Redes neurais, aprendizado profundo e embeddings", "-", OA_NEGOCIO),
    ("27 e 28", "17 e 19/11", "Apresentação e debate da Pesquisa de Campo",
     "Relatório", OA_NEGOCIO),
    ("29", "24/11", "Revisão geral e encerramento dos projetos aplicados", "-", OA_NEGOCIO),
    ("30 e 31", "26/11 e 01/12", "PROVA FINAL", "-", None),
]

DINAMICAS = [
    ("Python para desenvolver soluções",
     "o código é ferramenta de análise, não o objeto de estudo"),
    ("Projetos aplicados em sala",
     "em apps da disciplina, com correção automática ou feedback por rubrica"),
    ("Atividades e discussões em grupo",
     "dinâmicas de escolha e crítica de métodos, gráficos e pontos de corte"),
    ("Laboratórios de desenvolvimento intensivo",
     "aulas dedicadas à construção da pesquisa de campo e dos projetos"),
]


# ---------------------------------------------------------------- slides


def slide_01_noticia(prs, lays):
    slide = slide_limpo(prs, lays)
    imagem_ajustada(
        slide, os.path.join(ASSETS, "noticia_tse_g1.jpg"),
        MARGEM, Inches(0.85), FAIXA, Inches(5.25),
    )
    tb = texto_livre(slide, MARGEM, Inches(6.35), Inches(9.0), Inches(0.5))
    escrever(tb.text_frame, [
        {"texto": "g1 · Eleições 2026 · 14/07/2026", "tamanho": 12, "bold": True,
         "cor": VERMELHO, "depois": 2},
        {"texto": "Presidente do TSE propõe um selo de acerto para as pesquisas eleitorais.",
         "tamanho": 12, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def slide_02_artigo(prs, lays):
    slide = slide_com_titulo(prs, lays, "Portaria TSE: o Selo Acurácia Eleitoral")

    caixa(slide, MARGEM, Inches(2.35), Inches(0.10), Inches(2.05),
          preenchimento=VERMELHO)

    tb = texto_livre(slide, Inches(1.55), Inches(2.35), Inches(10.30), Inches(2.05))
    escrever(tb.text_frame, [
        {"texto": "Art. 1º", "fonte": DISPLAY, "tamanho": 20, "bold": True,
         "cor": VERMELHO, "depois": 8},
        {"texto": "Fica instituído o Selo Acurácia Eleitoral, destinado ao "
                  "reconhecimento e a valorização das empresas de pesquisa "
                  "eleitoral, cujas estimativas apresentem maior aderência aos "
                  "resultados oficiais proclamados pela Justiça Eleitoral.",
         "tamanho": 21, "entrelinhas": 1.25, "depois": 0},
    ])

    caixa(slide, MARGEM, Inches(5.30), FAIXA, Inches(1.05),
          preenchimento=RGBColor(0xF5, 0xF5, 0xF5))
    tb = texto_livre(slide, Inches(1.42), Inches(5.48), Inches(11.0), Inches(0.75))
    escrever(tb.text_frame, [
        {"texto": "\"Maior aderência aos resultados\" é uma medida. Qual?",
         "tamanho": 17, "bold": True, "depois": 3},
        {"texto": "Erro médio? Erro no primeiro colocado? Cobertura da margem "
                  "de erro? Cada escolha premia um instituto diferente.",
         "tamanho": 14, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def slide_03_jurimetria(prs, lays):
    slide = slide_limpo(prs, lays)
    imagem_ajustada(
        slide, os.path.join(ASSETS, "busca_jurimetria.jpg"),
        MARGEM, Inches(0.62), FAIXA, Inches(5.55),
    )
    tb = texto_livre(slide, MARGEM, Inches(6.35), Inches(10.0), Inches(0.5))
    escrever(tb.text_frame, [
        {"texto": BUSCA_GOOGLE, "tamanho": 12, "bold": True, "cor": VERMELHO,
         "link": BUSCA_GOOGLE, "sublinhado": True, "depois": 2},
        {"texto": "A pergunta já tem nome, método e uma comunidade que a estuda.",
         "tamanho": 12, "cor": CINZA_ESCURO, "depois": 0},
    ])
    return slide


def slide_04_pergunta(prs, lays):
    return slide_secao(
        prs, lays,
        "Quem é responsável por elaborar argumentos, decisões e análises "
        "estatisticamente corretas no âmbito do direito?",
        tamanho=40,
    )


def slide_05_objetivos(prs, lays):
    slide = slide_com_titulo(prs, lays, "Objetivos de aprendizagem")

    tb = texto_livre(slide, MARGEM, Inches(2.02), FAIXA, Inches(0.32))
    escrever(tb.text_frame, [
        {"texto": "Ao final da disciplina, você deverá ser capaz de:",
         "tamanho": 15, "cor": CINZA_ESCURO, "depois": 0},
    ])

    y = Inches(2.52)
    altura = Inches(0.86)
    passo = Inches(0.87)
    for cor, titulo, detalhe in OBJETIVOS:
        caixa(slide, MARGEM, y, Inches(0.13), altura, preenchimento=cor)
        tb = texto_livre(slide, Inches(1.44), y, Inches(11.0), altura)
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 17, "bold": True, "depois": 2},
            {"texto": detalhe, "tamanho": 13, "cor": CINZA_ESCURO,
             "entrelinhas": 1.05, "depois": 0},
        ])
        y = y + passo
    return slide


def slide_06_avaliacao(prs, lays):
    slide = slide_com_titulo(prs, lays, "Como você será avaliado")

    tb = texto_livre(slide, MARGEM, Inches(2.02), FAIXA, Inches(0.32))
    escrever(tb.text_frame, [
        {"texto": "Quatro avaliações, com pesos que somam 100%.",
         "tamanho": 15, "cor": CINZA_ESCURO, "depois": 0},
    ])

    largura = Inches(2.71)
    vao = Inches(0.22)
    y = Inches(2.60)
    altura = Inches(2.65)
    for i, (sigla, nome, peso, cor, detalhe) in enumerate(AVALIACAO):
        x = MARGEM + i * (largura + vao)
        caixa(slide, x, y, largura, altura,
              preenchimento=RGBColor(0xF7, 0xF7, 0xF7))
        caixa(slide, x, y, largura, Inches(0.09), preenchimento=cor)

        tb = texto_livre(slide, x + Inches(0.22), y + Inches(0.34),
                         largura - Inches(0.44), altura - Inches(0.5))
        escrever(tb.text_frame, [
            {"texto": sigla, "fonte": DISPLAY, "tamanho": 30, "bold": True,
             "cor": PRETO, "depois": 2},
            {"texto": nome, "tamanho": 14, "bold": True, "depois": 10},
            {"texto": peso, "fonte": DISPLAY, "tamanho": 26, "bold": True,
             "cor": VERMELHO, "depois": 8},
            {"texto": detalhe, "tamanho": 11.5, "cor": CINZA_ESCURO,
             "entrelinhas": 1.1, "depois": 0},
        ])

    tb = texto_livre(slide, MARGEM, Inches(5.48), FAIXA, Inches(1.1))
    escrever(tb.text_frame, [
        {"texto": "A nota final é a soma das quatro com esses pesos. "
                  "Não há nota mínima por bloco.",
         "tamanho": 13, "bold": True, "entrelinhas": 1.15, "depois": 6},
        {"texto": "A Pesquisa de Campo (PC) é um estudo empírico em grupo, com um "
                  "método que não vemos em aula: 7 temas, 14 grupos, relatório e debate.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "entrelinhas": 1.15, "depois": 4},
        {"texto": "Os Projetos Aplicados (PA) são 8 entregas, a maior parte em sala, "
                  "com correção automática ou rubrica visível na tela.",
         "tamanho": 12.5, "cor": CINZA_ESCURO, "entrelinhas": 1.15, "depois": 0},
    ])
    return slide


def legenda_oa(slide, y):
    """Faixa compacta lembrando a que objetivo cada cor corresponde."""
    itens = [
        (OA_PLANOS, "planos de estudo"),
        (OA_BASES, "bases de dados"),
        (OA_DESCRITIVAS, "análises descritivas"),
        (OA_MODELOS, "modelos estatísticos"),
        (OA_NEGOCIO, "problemas de negócio"),
    ]
    x = MARGEM
    for cor, rotulo in itens:
        caixa(slide, x, y + Inches(0.045), Inches(0.16), Inches(0.16),
              preenchimento=cor)
        tb = texto_livre(slide, x + Inches(0.22), y, Inches(1.95), Inches(0.26))
        escrever(tb.text_frame, [
            {"texto": rotulo, "tamanho": 10.5, "cor": CINZA_ESCURO, "depois": 0},
        ])
        x = x + Inches(2.30)


def tabela_de_aulas(slide, aulas, topo, altura):
    linhas = len(aulas) + 1
    tabela_shape = slide.shapes.add_table(
        linhas, 5, MARGEM, topo, FAIXA, altura
    )
    tabela = tabela_shape.table
    tabela_sem_estilo(tabela)

    for largura, i in zip(
        (Inches(0.13), Inches(0.85), Inches(1.28), Inches(8.14), Inches(1.10)),
        range(5),
    ):
        tabela.columns[i].width = largura

    alt_cabecalho = Inches(0.26)
    alt_linha = int((altura - alt_cabecalho) / len(aulas))
    tabela.rows[0].height = alt_cabecalho
    for r in range(1, linhas):
        tabela.rows[r].height = alt_linha

    cabecalho = ("", "Aula", "Data", "Tema", "Projeto")
    for c, titulo in enumerate(cabecalho):
        celula(tabela, 0, c, titulo, tamanho=10, bold=True, cor=BRANCO,
               fundo=PRETO,
               alinhamento=PP_ALIGN.LEFT if c == 3 else PP_ALIGN.CENTER)

    for r, (numero, data, tema, projeto, cor) in enumerate(aulas, start=1):
        prova = cor is None
        zebra = RGBColor(0xF4, 0xF4, 0xF4) if r % 2 == 0 else None
        fundo_linha = CINZA_CLARO if prova else zebra
        celula(tabela, r, 0, "", fundo=cor if cor else CINZA)
        celula(tabela, r, 1, numero, bold=True, fundo=fundo_linha,
               alinhamento=PP_ALIGN.CENTER)
        celula(tabela, r, 2, data, cor=CINZA_ESCURO,
               fundo=fundo_linha, alinhamento=PP_ALIGN.CENTER)
        celula(tabela, r, 3, tema, bold=prova, fundo=fundo_linha)
        celula(tabela, r, 4, projeto,
               bold=projeto not in ("-", ""),
               cor=VERMELHO if projeto not in ("-", "") else CINZA,
               fundo=fundo_linha, alinhamento=PP_ALIGN.CENTER)
    return tabela


def slide_08_aulas_pi(prs, lays):
    slide = slide_com_titulo(prs, lays, "Aulas até a Prova Intermediária")
    legenda_oa(slide, Inches(1.98))
    tabela_de_aulas(slide, AULAS_PI, Inches(2.36), Inches(4.36))
    return slide


def slide_09_aulas_pf(prs, lays):
    slide = slide_com_titulo(prs, lays, "Aulas até a Prova Final")
    legenda_oa(slide, Inches(1.98))
    tabela_de_aulas(slide, AULAS_PF, Inches(2.36), Inches(4.36))
    return slide


def slide_10_dinamicas(prs, lays):
    slide = slide_com_titulo(prs, lays, "Dinâmicas")

    y = Inches(2.20)
    altura = Inches(0.80)
    passo = Inches(0.86)
    for titulo, detalhe in DINAMICAS:
        caixa(slide, MARGEM, y, Inches(0.13), altura, preenchimento=VERMELHO)
        tb = texto_livre(slide, Inches(1.44), y, Inches(11.0), altura)
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        escrever(tb.text_frame, [
            {"texto": titulo, "tamanho": 17, "bold": True, "depois": 2},
            {"texto": detalhe, "tamanho": 13, "cor": CINZA_ESCURO, "depois": 0},
        ])
        y = y + passo

    caixa(slide, MARGEM, Inches(5.80), FAIXA, Inches(0.88),
          preenchimento=VERMELHO)
    tb = texto_livre(slide, Inches(1.50), Inches(5.98), Inches(10.85), Inches(0.6))
    escrever(tb.text_frame, [
        {"texto": "O foco da disciplina não é programação, e sim ciência de "
                  "dados aplicada ao direito.",
         "fonte": DISPLAY, "tamanho": 20, "bold": True, "cor": BRANCO,
         "alinhamento": PP_ALIGN.CENTER, "depois": 0},
    ])
    return slide


# ---------------------------------------------------------------- main


def main():
    prs = deck_limpo()
    lays = layouts(prs)

    slide_01_noticia(prs, lays)
    slide_02_artigo(prs, lays)
    slide_03_jurimetria(prs, lays)
    slide_04_pergunta(prs, lays)
    slide_05_objetivos(prs, lays)
    slide_06_avaliacao(prs, lays)
    slide_08_aulas_pi(prs, lays)
    slide_09_aulas_pf(prs, lays)
    slide_10_dinamicas(prs, lays)

    gravar(prs, SAIDA, titulo=TITULO_DECK)


if __name__ == "__main__":
    import sys

    sys.stdout.reconfigure(encoding="utf-8")
    main()

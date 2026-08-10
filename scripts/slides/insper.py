"""Ferramentas comuns para montar os slides da disciplina no template Insper.

Base: template oficial (Apresentacao didatica - Final.potx, ja convertido para
.pptx em cdadeng/assets_insper). Cores e tipografia seguem o Manual de Uso da
Marca (insper-guia-de-marca.pdf, secao 2.1).

Quem usa: build_aula01.py, build_aula02.py, ...
"""

from __future__ import annotations

import os
import re
import shutil
import zipfile
from datetime import datetime, timezone

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = r"C:\Users\jtrec\Documents\insper\cdadeng\assets_insper\template_insper.pptx"
RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"

# ---------------------------------------------------------------- marca

VERMELHO = RGBColor(0xE5, 0x05, 0x05)
PRETO = RGBColor(0x00, 0x00, 0x00)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
TURQUESA = RGBColor(0x3A, 0xCC, 0x9F)
VERDE = RGBColor(0x92, 0xD0, 0x53)
AMARELO = RGBColor(0xFF, 0xCC, 0x00)
LARANJA = RGBColor(0xF8, 0x9D, 0x49)
ROSA = RGBColor(0xF4, 0x7D, 0xCD)
ROXO = RGBColor(0x73, 0x0D, 0x9F)
CINZA_CLARO = RGBColor(0xDC, 0xDC, 0xDC)
CINZA = RGBColor(0xAB, 0xAB, 0xAB)
CINZA_ESCURO = RGBColor(0x5B, 0x5B, 0x5B)
QUASE_BRANCO = RGBColor(0xF7, 0xF7, 0xF7)

DISPLAY = "GT Ultra Fine Regular"
TEXTO = "Inter"

# cores dos objetivos de aprendizagem, herdadas do semestre anterior
# (verde = planos de estudo, ciano/turquesa = bases de dados,
#  amarelo = modelos, roxo/magenta = problemas de negocio)
OA_PLANOS, OA_BASES, OA_DESCRITIVAS, OA_MODELOS, OA_NEGOCIO = (
    VERDE, TURQUESA, LARANJA, AMARELO, ROXO
)

DISCIPLINA = "Ciência de Dados Aplicada ao Direito II"
AUTOR = "Julio Trecenti"

# geometria util lida do layout "Slide padrao 1"
MARGEM = Inches(1.17)
FAIXA = Inches(11.50)
TOPO = Inches(2.10)
BASE = Inches(6.72)


# ---------------------------------------------------------------- estrutura


def deck_limpo() -> Presentation:
    """Abre o template e descarta os slides de exemplo, preservando os masters."""
    prs = Presentation(TEMPLATE)
    lista = prs.slides._sldIdLst
    for sld in list(lista):
        prs.part.drop_rel(sld.get(RID))
        lista.remove(sld)
    return prs


def layouts(prs) -> dict:
    return {lay.name: lay for m in prs.slide_masters for lay in m.slide_layouts}


def limpar_placeholders_vazios(slide):
    for shape in list(slide.placeholders):
        if shape.has_text_frame and not shape.text_frame.text.strip():
            shape._element.getparent().remove(shape._element)


# ---------------------------------------------------------------- texto


def escrever(frame, linhas, *, fonte=TEXTO, tamanho=18, cor=PRETO, bold=False,
             espaco_antes=0, espaco_depois=6, entrelinhas=1.0, alinhamento=None):
    """Escreve paragrafos num text frame, um item de `linhas` por paragrafo.

    Cada item pode ser uma string ou um dict com sobrescritas locais.
    """
    frame.word_wrap = True
    for i, item in enumerate(linhas):
        cfg = {"texto": item} if isinstance(item, str) else dict(item)
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.space_before = Pt(cfg.get("antes", espaco_antes))
        p.space_after = Pt(cfg.get("depois", espaco_depois))
        p.line_spacing = cfg.get("entrelinhas", entrelinhas)
        al = cfg.get("alinhamento", alinhamento)
        if al is not None:
            p.alignment = al
        run = p.add_run()
        run.text = cfg["texto"]
        if cfg.get("link"):
            run.hyperlink.address = cfg["link"]
        f = run.font
        f.name = cfg.get("fonte", fonte)
        f.size = Pt(cfg.get("tamanho", tamanho))
        f.bold = cfg.get("bold", bold)
        # a cor precisa vir depois do hyperlink: o PowerPoint aplica a cor de
        # link do tema quando o run nasce com hlinkClick e nenhuma cor propria
        f.color.rgb = cfg.get("cor", cor)
        if cfg.get("sublinhado"):
            f.underline = True
    return frame


def caixa(slide, x, y, cx, cy, *, preenchimento=None, borda=None, largura_borda=1.0,
          forma=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(forma, x, y, cx, cy)
    if preenchimento is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = preenchimento
    if borda is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = borda
        sh.line.width = Pt(largura_borda)
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    sh.text_frame.margin_left = Inches(0.16)
    sh.text_frame.margin_right = Inches(0.16)
    sh.text_frame.margin_top = Inches(0.10)
    sh.text_frame.margin_bottom = Inches(0.10)
    return sh


def texto_livre(slide, x, y, cx, cy):
    tb = slide.shapes.add_textbox(x, y, cx, cy)
    tb.text_frame.word_wrap = True
    tb.text_frame.margin_left = 0
    tb.text_frame.margin_right = 0
    tb.text_frame.margin_top = 0
    tb.text_frame.margin_bottom = 0
    return tb


# ---------------------------------------------------------------- slides


def slide_capa(prs, lays, *, subtitulo, subtema=""):
    slide = prs.slides.add_slide(lays["Capa 1"])
    valores = {0: DISCIPLINA, 13: subtitulo, 15: subtema, 14: AUTOR}
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if idx in valores and valores[idx]:
            ph.text_frame.text = valores[idx]
    limpar_placeholders_vazios(slide)
    return slide


def slide_secao(prs, lays, titulo, *, layout="Tema da Seção 1", tamanho=44,
                cor=BRANCO):
    slide = prs.slides.add_slide(lays[layout])
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx != 0:
            continue
        ph.left, ph.width = Inches(1.20), Inches(10.93)
        ph.top, ph.height = Inches(2.30), Inches(2.90)
        frame = ph.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame.text = titulo
        # titulo com quebra de linha vira mais de um paragrafo, e o que nao for
        # formatado aqui herda os 60pt do layout
        for p in frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.15
            for r in p.runs:
                r.font.name = DISPLAY
                r.font.size = Pt(tamanho)
                r.font.bold = False
                r.font.color.rgb = cor
    limpar_placeholders_vazios(slide)
    return slide


def slide_com_titulo(prs, lays, titulo, eyebrow, *, tamanho=30):
    slide = prs.slides.add_slide(lays["Slide padrão 1"])
    for ph in list(slide.placeholders):
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text_frame.text = titulo
            r = ph.text_frame.paragraphs[0].runs[0]
            r.font.name = DISPLAY
            r.font.size = Pt(tamanho)
            r.font.bold = True
            r.font.color.rgb = PRETO
        elif idx == 15:
            ph.text_frame.text = eyebrow
            r = ph.text_frame.paragraphs[0].runs[0]
            r.font.name = TEXTO
            r.font.size = Pt(11)
            r.font.color.rgb = CINZA_ESCURO
    limpar_placeholders_vazios(slide)
    return slide


def slide_limpo(prs, lays):
    return prs.slides.add_slide(lays["Slide padrão 2"])


def imagem_ajustada(slide, caminho, x, y, cx_max, cy_max, *, centralizar=True,
                    moldura=True):
    """Insere a imagem cabendo na caixa, preservando a proporcao."""
    from PIL import Image

    with Image.open(caminho) as im:
        w, h = im.size
    escala = min(cx_max / w, cy_max / h)
    cx, cy = int(w * escala), int(h * escala)
    if centralizar:
        x = x + (cx_max - cx) // 2
        y = y + (cy_max - cy) // 2
    pic = slide.shapes.add_picture(caminho, x, y, cx, cy)
    if moldura:
        pic.line.color.rgb = CINZA_CLARO
        pic.line.width = Pt(0.75)
    return pic


def lista_com_barras(slide, itens, *, topo, altura, passo, cor_padrao=VERMELHO,
                     tamanho_titulo=17, tamanho_detalhe=13, largura_barra=0.13,
                     x_texto=1.44, largura_texto=11.0):
    """Lista em que cada item ganha uma barra colorida a esquerda.

    `itens` e uma lista de (cor, titulo, detalhe); cor None usa `cor_padrao`.
    """
    y = topo
    for item in itens:
        cor, titulo, detalhe = item if len(item) == 3 else (None, *item)
        caixa(slide, MARGEM, y, Inches(largura_barra), altura,
              preenchimento=cor or cor_padrao)
        tb = texto_livre(slide, Inches(x_texto), y, Inches(largura_texto), altura)
        tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        linhas = [{"texto": titulo, "tamanho": tamanho_titulo, "bold": True, "depois": 2}]
        if detalhe:
            linhas.append({"texto": detalhe, "tamanho": tamanho_detalhe,
                           "cor": CINZA_ESCURO, "entrelinhas": 1.05, "depois": 0})
        escrever(tb.text_frame, linhas)
        y = y + passo


def celula(tabela, linha, coluna, texto, *, tamanho=9.5, bold=False, cor=PRETO,
           fundo=None, alinhamento=PP_ALIGN.LEFT):
    c = tabela.cell(linha, coluna)
    c.text = ""
    c.margin_left = Inches(0.06)
    c.margin_right = Inches(0.06)
    c.margin_top = 0
    c.margin_bottom = 0
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fundo is None:
        c.fill.background()
    else:
        c.fill.solid()
        c.fill.fore_color.rgb = fundo
    p = c.text_frame.paragraphs[0]
    p.alignment = alinhamento
    run = p.add_run()
    run.text = texto
    run.font.name = TEXTO
    run.font.size = Pt(tamanho)
    run.font.bold = bold
    run.font.color.rgb = cor
    return c


def tabela_sem_estilo(tabela):
    """Desliga o estilo listrado herdado do tema; o visual vem das cores da marca."""
    from lxml import etree

    tbl = tabela._tbl
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    pr = tbl.find(f"{ns}tblPr")
    if pr is not None:
        pr.set("firstRow", "0")
        pr.set("bandRow", "0")
        for filho in list(pr):
            if etree.QName(filho).localname == "tableStyleId":
                pr.remove(filho)


# ------------------------------------------------- limpeza dos metadados

# o .potx do Insper chega com 74 slides de exemplo. Ao remove-los sobra um
# rastro que o PowerPoint continua exibindo (secoes) ou guardando (tags do
# Articulate, miniatura, autoria do template). Nada disso e nosso.
DESCARTAR = re.compile(
    r"^(docProps/thumbnail\.jpeg|docProps/custom\.xml|ppt/tags/tag\d+\.xml)$"
)
NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
NS_P14 = "{http://schemas.microsoft.com/office/powerpoint/2010/main}"

APP_XML = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/\
extended-properties" xmlns:vt="http://schemas.openxmlformats.org/\
officeDocument/2006/docPropsVTypes"><Application>Microsoft Office PowerPoint\
</Application><PresentationFormat>Widescreen</PresentationFormat><Slides>\
{slides}</Slides><Notes>0</Notes><HiddenSlides>0</HiddenSlides><MMClips>0\
</MMClips><ScaleCrop>false</ScaleCrop><Company>Insper</Company>\
<LinksUpToDate>false</LinksUpToDate><SharedDoc>false</SharedDoc>\
<HyperlinksChanged>false</HyperlinksChanged><AppVersion>16.0000</AppVersion>\
</Properties>"""


def _sem_secoes_e_tags(xml: bytes) -> bytes:
    """Tira do XML a lista de secoes do template e os ganchos de tags."""
    from lxml import etree

    raiz = etree.fromstring(xml)
    for lista in raiz.iter(f"{NS_P14}sectionLst"):
        ext = lista.getparent()          # <p:ext uri="{521415D9-...}">
        ext.getparent().remove(ext)
        break
    for cust in list(raiz.iter(f"{NS_P}custDataLst")):
        if cust.find(f"{NS_P}tags") is not None:
            cust.getparent().remove(cust)
    return etree.tostring(raiz, xml_declaration=True, encoding="UTF-8",
                          standalone=True)


def _sem_rels(xml: bytes) -> bytes:
    """Remove as relacoes que apontam para as partes descartadas."""
    return re.sub(
        r"<Relationship\b[^>]*Target=\"[^\"]*"
        r"(?:tags/tag\d+\.xml|thumbnail\.jpeg|custom\.xml)\"[^>]*/>",
        "", xml.decode("utf-8"),
    ).encode("utf-8")


def limpar_pacote(caminho: str, n_slides: int):
    """Reescreve o .pptx sem os resquicios do arquivo de origem do template."""
    temporario = caminho + ".tmp"
    with zipfile.ZipFile(caminho) as origem, \
            zipfile.ZipFile(temporario, "w", zipfile.ZIP_DEFLATED) as destino:
        for item in origem.infolist():
            nome = item.filename
            if DESCARTAR.match(nome):
                continue
            dados = origem.read(nome)
            if nome == "docProps/app.xml":
                dados = APP_XML.format(slides=n_slides).encode("utf-8")
            elif nome == "[Content_Types].xml":
                dados = re.sub(
                    r"<Override\b[^>]*PartName=\"/(?:ppt/tags/tag\d+|"
                    r"docProps/custom)\.xml\"[^>]*/>",
                    "", dados.decode("utf-8"),
                ).encode("utf-8")
            elif nome.endswith(".rels"):
                dados = _sem_rels(dados)
            elif nome in ("ppt/presentation.xml",) or (
                nome.startswith(("ppt/slideMasters/slideMaster",
                                 "ppt/slideLayouts/slideLayout"))
                and nome.endswith(".xml")
            ):
                dados = _sem_secoes_e_tags(dados)
            destino.writestr(item, dados)
    shutil.move(temporario, caminho)


def gravar(prs, caminho: str, *, titulo: str, assunto: str = DISCIPLINA) -> int:
    """Grava o deck com metadados nossos e sem os resquicios do template."""
    props = prs.core_properties
    props.title = titulo
    props.subject = assunto
    props.author = AUTOR
    props.last_modified_by = AUTOR
    props.category = ""
    props.comments = ""
    props.keywords = ""
    props.revision = 1
    # as datas que vinham no arquivo eram as da criacao do template
    agora = datetime.now(timezone.utc).replace(microsecond=0, tzinfo=None)
    props.created = agora
    props.modified = agora

    n_slides = len(prs.slides._sldIdLst)
    prs.save(caminho)
    limpar_pacote(caminho, n_slides)
    print(f"{n_slides} slides -> {caminho}")
    return n_slides
